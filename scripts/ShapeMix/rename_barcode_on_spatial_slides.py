#!/usr/bin/env python3
"""Rename "barcode" to "spot barcode" on the two spatial-ATAC slides only.

Only the two slide XML payloads (and their notes payloads) inserted by
insert_spatial_atac_slides.py are rewritten; every other package part is
copied byte-for-byte from the original file.

Caution: python-pptx renumbers slide partnames sequentially at load time, so
the in-memory partnames do not match the names stored in the ZIP (the deck's
slide order no longer matches its stored part numbering). Parts are therefore
matched by content markers, never by name.
"""

from __future__ import annotations

import re
from copy import deepcopy
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

KICKER_A = "SPATIAL ATAC"
KICKER_B = "WHY DECONVOLUTION"

# Exact single-run replacements per slide.
REPLACEMENTS_A = {
    "Same Tn5 chemistry as ATAC-seq — plus a barcode that records where each fragment came from.":
        "Same Tn5 chemistry as ATAC-seq — plus a spot barcode that records where each fragment came from.",
    "each circle = one spot with its own DNA barcode":
        "each circle = one spot with its own spot barcode",
    "Each fragment is sequenced with its spot's barcode":
        "Each fragment is sequenced with its spot barcode",
    "sorting by barcode gives peak counts at every tissue location":
        "sorting by spot barcode gives peak counts at every tissue location",
}
REPLACEMENTS_B = {
    "A barcoded spot is about 50 µm across; a cell is about 10 µm. Every cell under a spot shares the same barcode.":
        "A spot is about 50 µm across; a cell is about 10 µm. Every cell under a spot shares the same spot barcode.",
    "ONE BLENDED PROFILE PER BARCODE":
        "ONE BLENDED PROFILE PER SPOT BARCODE",
    "OFTEN ~10 CELLS SHARE ONE BARCODE":
        "OFTEN ~10 CELLS SHARE ONE SPOT BARCODE",
}

NOTES_A = (
    "Spatial ATAC applies the same Tn5 cut-and-tag chemistry the audience already saw, but "
    "on an intact tissue slice placed over a grid of barcoded spots. Every fragment released "
    "inside a spot is labeled with its spot barcode, so after sequencing the fragments can "
    "be sorted by location. The output is an ATAC peak-count table for every spot — an "
    "open-chromatin map of the tissue. This slide introduces the technology; the next slide "
    "explains the resolution catch that motivates deconvolution."
)
NOTES_B = (
    "The catch is resolution: a spot is roughly 50 micrometers across while a cell is "
    "roughly 10 micrometers, so one spot barcode usually collects fragments from several "
    "cells — often around ten, matching the pseudo-spot design used later in the benchmark. "
    "The example shows 3 T cells and 1 B cell under one spot producing a single summed "
    "profile of [8, 4, 7, 5]; these are exactly the numbers the following deconvolution "
    "slide solves, recovering the recipe 75% T + 25% B. Because per-cell identities are "
    "hidden in the sum, spot-level analysis needs deconvolution."
)


def slide_texts(slide) -> list[str]:
    return [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]


def apply_replacements(slide, mapping) -> int:
    applied = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text in mapping:
                    new_text = mapping[run.text]
                    run.text = new_text
                    applied += 1
                    # The blended-profile pill needs a wider box for the
                    # longer wording; nudge it left and shrink the font a bit.
                    if new_text == "ONE BLENDED PROFILE PER SPOT BARCODE":
                        shape.left = Inches(9.38)
                        shape.width = Inches(3.05)
                        run.font.size = Pt(8.2)
    return applied


def find_slide_part(zipfile: ZipFile, marker: str) -> str:
    matches = [
        name
        for name in zipfile.namelist()
        if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
        and marker.encode("utf-8") in zipfile.read(name)
    ]
    if len(matches) != 1:
        raise RuntimeError(f"Expected one slide part containing {marker!r}; found {matches}")
    return matches[0]


def notes_part_for(zipfile: ZipFile, slide_part: str) -> str:
    rels_name = f"ppt/slides/_rels/{slide_part.rsplit('/', 1)[-1]}.rels"
    rels = zipfile.read(rels_name).decode("utf-8")
    targets = re.findall(r'Target="\.\./(notesSlides/notesSlide\d+\.xml)"', rels)
    if len(targets) != 1:
        raise RuntimeError(f"Expected one notes target in {rels_name}; found {targets}")
    return f"ppt/{targets[0]}"


def ensure_notes_body(presentation, slide):
    """Give the slide's notes page a body placeholder if it lacks one.

    The deck's notes master has no body placeholder, so notes pages created
    by python-pptx come up empty. Clone the body shape from the title
    slide's notes page, which does have one.
    """
    from pptx.enum.shapes import PP_PLACEHOLDER

    notes = slide.notes_slide
    if notes.notes_text_frame is not None:
        return notes
    donor = presentation.slides[0].notes_slide
    donor_sp = None
    for placeholder in donor.placeholders:
        if placeholder.placeholder_format.type == PP_PLACEHOLDER.BODY:
            donor_sp = placeholder._element
            break
    if donor_sp is None:
        raise RuntimeError("No donor notes body placeholder found on slide 1")
    notes.shapes._spTree.append(deepcopy(donor_sp))
    if notes.notes_text_frame is None:
        raise RuntimeError("Cloned notes body placeholder was not recognized")
    return notes


def main() -> None:
    presentation = Presentation(DECK)
    slides = presentation.slides
    slide_a, slide_b = slides[4], slides[5]
    if KICKER_A not in " ".join(slide_texts(slide_a)):
        raise RuntimeError("Slide 5 is not the spatial-ATAC slide")
    if KICKER_B not in " ".join(slide_texts(slide_b)):
        raise RuntimeError("Slide 6 is not the why-deconvolution slide")

    if apply_replacements(slide_a, REPLACEMENTS_A) != len(REPLACEMENTS_A):
        raise RuntimeError("Slide 5 text did not match the expected wording")
    if apply_replacements(slide_b, REPLACEMENTS_B) != len(REPLACEMENTS_B):
        raise RuntimeError("Slide 6 text did not match the expected wording")

    for slide, notes_text in ((slide_a, NOTES_A), (slide_b, NOTES_B)):
        notes = ensure_notes_body(presentation, slide)
        notes.notes_text_frame.text = notes_text

    with TemporaryDirectory(prefix="shapemix_rename_barcode_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"
        presentation.save(generated)

        with ZipFile(DECK, "r") as original, ZipFile(
            generated, "r"
        ) as regenerated, ZipFile(assembled, "w", compression=ZIP_DEFLATED) as output:
            # Match the two edited slides in BOTH packages by content, since
            # python-pptx renumbers stored partnames on load/save.
            payload_map = {}
            for marker in (KICKER_A, KICKER_B):
                orig_slide = find_slide_part(original, marker)
                gen_slide = find_slide_part(regenerated, marker)
                payload_map[orig_slide] = gen_slide
                payload_map[notes_part_for(original, orig_slide)] = notes_part_for(
                    regenerated, gen_slide
                )

            for info in original.infolist():
                if info.filename in payload_map:
                    data = regenerated.read(payload_map[info.filename])
                else:
                    data = original.read(info.filename)
                output.writestr(info, data)

        check = Presentation(assembled)
        joined_a = " ".join(slide_texts(check.slides[4]))
        joined_b = " ".join(slide_texts(check.slides[5]))
        for new_text in REPLACEMENTS_A.values():
            if new_text not in joined_a:
                raise RuntimeError(f"Missing on slide 5 after save: {new_text[:50]}...")
        for new_text in REPLACEMENTS_B.values():
            if new_text not in joined_b:
                raise RuntimeError(f"Missing on slide 6 after save: {new_text[:50]}...")
        joined_7 = " ".join(slide_texts(check.slides[6]))
        if "smoothie" not in joined_7:
            raise RuntimeError("Slide 7 is no longer the deconvolution smoothie slide")
        if len(check.slides) != len(presentation.slides):
            raise RuntimeError("Slide count changed unexpectedly")

        assembled.replace(DECK)

    print(f"Renamed barcode → spot barcode on slides 5 and 6: {DECK}")


if __name__ == "__main__":
    main()
