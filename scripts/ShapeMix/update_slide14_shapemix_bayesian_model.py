#!/usr/bin/env python3
"""Rebuild slide 14 in the compact slide-9 Bayesian style.

Only slide 14's XML payload is replaced. Every other package part remains
byte-for-byte unchanged.
"""

from __future__ import annotations

from hashlib import sha256
from pathlib import Path
from tempfile import TemporaryDirectory
from xml.etree import ElementTree
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation

from insert_shapemix_bayesian_model_slide import DECK, TITLE, build, slide_texts

SLIDE_INDEX = 13
OLD_TITLE = "ShapeMix adds fragment-length bins to the Bayesian model"


def package_hashes(path: Path) -> dict[str, str]:
    with ZipFile(path) as archive:
        return {
            name: sha256(archive.read(name)).hexdigest()
            for name in archive.namelist()
        }


def assemble(
    original_path: Path,
    generated_path: Path,
    output_path: Path,
    original_slide_part: str,
    generated_slide_part: str,
) -> None:
    with ZipFile(original_path, "r") as original, ZipFile(
        generated_path, "r"
    ) as generated, ZipFile(output_path, "w", compression=ZIP_DEFLATED) as output:
        if generated_slide_part not in generated.namelist():
            raise RuntimeError(f"Generated package is missing {generated_slide_part}")
        for info in original.infolist():
            data = (
                generated.read(generated_slide_part)
                if info.filename == original_slide_part
                else original.read(info.filename)
            )
            output.writestr(info, data)


def slide_part_at(path: Path, slide_index: int) -> str:
    """Resolve a slide position to its actual package filename."""
    presentation_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    document_rel_ns = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    )
    package_rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"

    with ZipFile(path) as archive:
        presentation = ElementTree.fromstring(
            archive.read("ppt/presentation.xml")
        )
        slide_ids = presentation.find(f"{{{presentation_ns}}}sldIdLst")
        if slide_ids is None or slide_index >= len(slide_ids):
            raise RuntimeError(f"Slide position {slide_index + 1} does not exist")
        relationship_id = slide_ids[slide_index].attrib[
            f"{{{document_rel_ns}}}id"
        ]

        relationships = ElementTree.fromstring(
            archive.read("ppt/_rels/presentation.xml.rels")
        )
        target = None
        for relationship in relationships.findall(
            f"{{{package_rel_ns}}}Relationship"
        ):
            if relationship.attrib.get("Id") == relationship_id:
                target = relationship.attrib.get("Target")
                break
        if not target:
            raise RuntimeError(
                f"No slide relationship found for position {slide_index + 1}"
            )
        return target.lstrip("/") if target.startswith("/") else f"ppt/{target}"


def main() -> None:
    before = package_hashes(DECK)
    current = Presentation(DECK)
    if len(current.slides) != 33:
        raise RuntimeError(f"Expected 33 slides; found {len(current.slides)}")
    if not any(OLD_TITLE in text for text in slide_texts(current.slides[SLIDE_INDEX])):
        raise RuntimeError("Slide 14 is not the expected ShapeMix model slide")
    original_slide_part = slide_part_at(DECK, SLIDE_INDEX)

    with TemporaryDirectory(prefix="shapemix_slide14_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"

        presentation = Presentation(DECK)
        slide = presentation.slides[SLIDE_INDEX]
        for shape in list(slide.shapes):
            shape._element.getparent().remove(shape._element)
        build(slide)
        presentation.save(generated)
        generated_slide_part = slide_part_at(generated, SLIDE_INDEX)
        assemble(
            DECK,
            generated,
            assembled,
            original_slide_part,
            generated_slide_part,
        )

        check = Presentation(assembled)
        if len(check.slides) != 33:
            raise RuntimeError("Slide count changed while updating slide 14")
        joined = " ".join(slide_texts(check.slides[SLIDE_INDEX]))
        for needle in (
            TITLE,
            "R_LB",
            "n_lb",
            "N_LB",
            "p(z | N_LB, N)",
            "POSTERIOR",
            "LIKELIHOOD",
            "PRIOR",
        ):
            if needle not in joined:
                raise RuntimeError(f"Revised slide 14 is missing {needle!r}")
        for forbidden in ("Multinomial", "C × P", "S × P"):
            if forbidden in joined:
                raise RuntimeError(f"Revised slide 14 contains {forbidden!r}")

        after = package_hashes(assembled)
        for name, digest in before.items():
            if name == original_slide_part:
                continue
            if after.get(name) != digest:
                raise RuntimeError(f"Unexpected change outside slide 14: {name}")

        assembled.replace(DECK)

    print(f"Rebuilt slide 14 in the compact Bayes style: {DECK}")
    print("Verified: every package part except slide 14 is byte-identical")


if __name__ == "__main__":
    main()
