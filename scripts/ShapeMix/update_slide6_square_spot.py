#!/usr/bin/env python3
"""Redraw slide 6's left panel: square spot, two cell colors only.

The zoomed spot becomes a square (echoing the navy-highlighted square on the
grid slide it zooms into) and the cells inside use only the two example
colors — teal-green T cells and blue B cells, six and two, i.e. the same
75/25 mixture the rest of the example uses.

Only the why-deconvolution slide part and its notes part are rewritten;
every other package part is copied byte-for-byte from the original file.
Parts are matched by content because python-pptx renumbers stored partnames
on load.
"""

from __future__ import annotations

import re
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

KICKER = "WHY DECONVOLUTION"

NAVY = "13233B"
SLATE = "52627A"
TEAL = "18A999"
BLUE = "3B82F6"
WHITE = "FFFFFF"

FONT = "Aptos"

EMU_PER_INCH = 914400

# Square spot: same bounds the dashed circle used to occupy.
SPOT_X, SPOT_Y, SPOT_SIDE = 1.945, 2.675, 2.85

# Eight cells, six teal T + two blue B = the example's 75/25 mixture.
CELLS = [
    (2.62, 3.42, TEAL),
    (3.32, 3.30, BLUE),
    (4.02, 3.52, TEAL),
    (2.42, 4.12, TEAL),
    (3.12, 4.02, TEAL),
    (3.87, 4.15, TEAL),
    (2.72, 4.82, BLUE),
    (4.07, 4.72, TEAL),
]
CELL_D = 0.52

NOTES = (
    "The catch is resolution: a spot is roughly 50 micrometers across while a cell is "
    "roughly 10 micrometers, so one spot barcode usually collects fragments from several "
    "cells — often around ten, matching the pseudo-spot design used later in the benchmark. "
    "The zoomed spot is drawn as a square to match the barcoded grid on the previous slide, "
    "and it contains only the two cell types used throughout this example: six teal-green "
    "T cells and two blue B cells, a 75/25 mixture. The right panel simplifies that to "
    "3 T + 1 B — the same 75/25 ratio — producing the single summed profile [8, 4, 7, 5]; "
    "these are exactly the numbers the following deconvolution slide solves, recovering "
    "75% T + 25% B. Because per-cell identities are hidden in the sum, spot-level analysis "
    "needs deconvolution."
)


def C(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#").upper())


def inches(emu: int) -> float:
    return emu / EMU_PER_INCH


def add_shape(slide, kind, x, y, w, h, fill, line=None, line_width=1.0):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(fill)
    shape.line.color.rgb = C(line if line else fill)
    shape.line.width = Pt(line_width)
    return shape


def add_circle(slide, x, y, d, fill, line=WHITE, line_width=1.5):
    return add_shape(slide, MSO_SHAPE.OVAL, x, y, d, d, fill, line=line, line_width=line_width)


def add_cell(slide, x, y, d, color):
    add_circle(slide, x, y, d, color, line=WHITE, line_width=2)
    add_circle(slide, x + d * 0.34, y + d * 0.33, d * 0.27, WHITE, line=WHITE, line_width=0.5)


def add_text(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = C(color)
    return box


def add_line(slide, x1, y1, x2, y2, color, width):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = C(color)
    line.line.width = Pt(width)
    return line


def slide_texts(slide) -> list[str]:
    return [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]


def remove_old_left_diagram(slide) -> None:
    """Remove the dashed circle spot, its cells, and the cell-size callout."""
    removed_ovals = 0
    removed_other = 0
    for shape in list(slide.shapes):
        x, y = inches(shape.left), inches(shape.top)
        in_left_panel = x < 6.40 and 1.90 <= y <= 6.40
        if not in_left_panel:
            continue
        is_oval = "Oval" in shape.name
        is_cell_label = shape.has_text_frame and "one cell" in shape.text_frame.text
        is_callout_line = (
            shape.shape_type == MSO_SHAPE_TYPE.LINE and y < 4.0 and inches(shape.height) < 0.4
        )
        if is_oval:
            removed_ovals += 1
        elif is_cell_label or is_callout_line:
            removed_other += 1
        else:
            continue
        shape._element.getparent().remove(shape._element)
    # 1 spot circle + 9 cells + 9 highlight dots
    if removed_ovals != 19 or removed_other != 2:
        raise RuntimeError(
            f"Unexpected removal counts: {removed_ovals} ovals, {removed_other} callout shapes"
        )


def draw_square_spot(slide) -> None:
    # Square spot, styled like the navy-highlighted grid square it zooms into.
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        SPOT_X,
        SPOT_Y,
        SPOT_SIDE,
        SPOT_SIDE,
        WHITE,
        line=NAVY,
        line_width=2.5,
    )
    for cx, cy, color in CELLS:
        add_cell(slide, cx - CELL_D / 2, cy - CELL_D / 2, CELL_D, color)
    # Cell-size callout, redrawn after the square so it stays on top.
    add_text(slide, "one cell\n≈ 10 µm", 4.72, 3.02, 1.28, 0.44, size=10, color=SLATE, bold=True)
    add_line(slide, 4.78, 3.34, 4.24, 3.58, color=SLATE, width=1.0)


def set_notes(slide) -> None:
    tf = slide.notes_slide.notes_text_frame
    if tf is None:
        raise RuntimeError("Notes body placeholder missing on the why-deconvolution slide")
    tf.text = NOTES


def find_slide_part(zipfile: ZipFile, marker: str) -> str:
    matches = [
        name
        for name in zipfile.namelist()
        if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
        and marker.encode("utf-8") in zipfile.read(name)
    ]
    if len(matches) != 1:
        raise RuntimeError(f"Expected one slide part containing {marker!r}; found {matches}")
    return matches[0]


def notes_part_for(zipfile: ZipFile, slide_part: str) -> str:
    rels_name = f"ppt/slides/_rels/{slide_part.rsplit('/', 1)[-1]}.rels"
    rels = zipfile.read(rels_name).decode("utf-8")
    targets = re.findall(r'Target="\.\./(notesSlides/notesSlide\d+\.xml)"', rels)
    if len(targets) != 1:
        raise RuntimeError(f"Expected one notes target in {rels_name}; found {targets}")
    return f"ppt/{targets[0]}"


def main() -> None:
    presentation = Presentation(DECK)
    slide = presentation.slides[5]
    if KICKER not in " ".join(slide_texts(slide)):
        raise RuntimeError("Slide 6 is not the why-deconvolution slide")

    remove_old_left_diagram(slide)
    draw_square_spot(slide)
    set_notes(slide)

    with TemporaryDirectory(prefix="shapemix_square_spot_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"
        presentation.save(generated)

        with ZipFile(DECK, "r") as original, ZipFile(
            generated, "r"
        ) as regenerated, ZipFile(assembled, "w", compression=ZIP_DEFLATED) as output:
            orig_slide = find_slide_part(original, KICKER)
            gen_slide = find_slide_part(regenerated, KICKER)
            payload_map = {
                orig_slide: gen_slide,
                notes_part_for(original, orig_slide): notes_part_for(regenerated, gen_slide),
            }
            for info in original.infolist():
                if info.filename in payload_map:
                    data = regenerated.read(payload_map[info.filename])
                else:
                    data = original.read(info.filename)
                output.writestr(info, data)

        check = Presentation(assembled)
        joined = " ".join(slide_texts(check.slides[5]))
        for expected in (KICKER, "One spot under the microscope", "one cell"):
            if expected not in joined:
                raise RuntimeError(f"Missing on saved slide 6: {expected!r}")
        if "SPATIAL ATAC" not in " ".join(slide_texts(check.slides[4])):
            raise RuntimeError("Slide 5 changed unexpectedly")
        if len(check.slides) != len(presentation.slides):
            raise RuntimeError("Slide count changed unexpectedly")

        assembled.replace(DECK)

    print(f"Redrew slide 6's spot as a square with two cell colors: {DECK}")


if __name__ == "__main__":
    main()
