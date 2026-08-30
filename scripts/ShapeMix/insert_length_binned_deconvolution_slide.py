#!/usr/bin/env python3
"""Insert a length-binned clone immediately after slide 5.

The final package is assembled from the original ZIP so every pre-existing
slide and notes XML payload remains byte-for-byte unchanged.
"""

from __future__ import annotations

from copy import deepcopy
from hashlib import sha256
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

TITLE = "Deconvolution works like identifying ingredients in a smoothie"
MARKER = "LENGTH BINS:"

NAVY = "13233B"
SLATE = "52627A"
CORAL = "F26B5B"
GOLD = "F4B942"
PURPLE = "8B5CF6"
WHITE = "FFFFFF"
FONT = "Aptos"


# Counts are ordered bottom-to-top: short, medium, long. Each tuple preserves
# the solid bar's original total. The mixed profile is exactly 75% T + 25% B
# within every bin, not only after summing the bins.
BIN_COUNTS = {
    # T reference: totals 10, 2, 8, 4.
    11: (5, 4, 1),
    14: (1, 1, 0),
    17: (4, 2, 2),
    20: (0, 2, 2),
    # B reference: totals 2, 10, 4, 8.
    29: (1, 0, 1),
    32: (5, 1, 4),
    35: (0, 2, 2),
    38: (4, 2, 2),
    # Unknown mixed spot: totals 8, 4, 7, 5.
    48: (4, 3, 1),
    51: (2, 1, 1),
    54: (3, 2, 2),
    57: (1, 2, 2),
}


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def slide_text(slide) -> list[str]:
    return [shape.text for shape in slide.shapes if hasattr(shape, "text_frame")]


def shape_by_id(slide, shape_id: int):
    matches = [shape for shape in slide.shapes if shape.shape_id == shape_id]
    if len(matches) != 1:
        raise RuntimeError(f"Expected one shape with id {shape_id}; found {len(matches)}")
    return matches[0]


def set_textbox(
    shape,
    text: str,
    *,
    size: float,
    text_color: str,
    bold: bool = True,
    align=PP_ALIGN.CENTER,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(text_color)


def add_pill(slide, text, x, y, width, height, fill, text_color, size):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(fill)
    shape.line.width = Pt(0.8)
    tf = shape.text_frame
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = 0
    tf.margin_bottom = 0
    set_textbox(shape, text, size=size, text_color=text_color)
    return shape


def add_segment(slide, name, left, top, width, height, fill):
    segment = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    segment.name = name
    segment.fill.solid()
    segment.fill.fore_color.rgb = rgb(fill)
    segment.line.color.rgb = rgb(fill)
    segment.line.width = Pt(0.35)
    return segment


def replace_solid_bars_with_stacks(slide) -> None:
    colors = (CORAL, GOLD, PURPLE)
    for bar_id, bins in BIN_COUNTS.items():
        bar = shape_by_id(slide, bar_id)
        total = sum(bins)
        if total <= 0:
            raise RuntimeError(f"Bar {bar_id} has no counts")

        left, width = bar.left, bar.width
        full_height = bar.height
        bottom = bar.top + bar.height
        parent = bar._element.getparent()
        parent.remove(bar._element)

        nonzero = [(index, count) for index, count in enumerate(bins) if count]
        cursor = bottom
        used_height = 0
        for position, (bin_index, count) in enumerate(nonzero):
            if position == len(nonzero) - 1:
                segment_height = full_height - used_height
            else:
                segment_height = round(full_height * count / total)
                used_height += segment_height
            cursor -= segment_height
            add_segment(
                slide,
                f"Length-bin stack bar {bar_id} bin {bin_index + 1}",
                left,
                cursor,
                width,
                segment_height,
                colors[bin_index],
            )
        if cursor != bottom - full_height:
            raise RuntimeError(f"Stack geometry failed for bar {bar_id}")


def build_variant_in_memory(presentation: Presentation):
    source_matches = []
    variant_matches = []
    for index, slide in enumerate(presentation.slides):
        texts = slide_text(slide)
        if TITLE in texts:
            if MARKER in texts:
                variant_matches.append((index, slide))
            else:
                source_matches.append((index, slide))
    if variant_matches:
        raise RuntimeError("The length-binned slide already exists; refusing to add a duplicate")
    if len(source_matches) != 1:
        raise RuntimeError(f"Expected one source slide; found {len(source_matches)}")

    source_index, source = source_matches[0]
    if source_index != 4:
        raise RuntimeError(f"Expected the source to be slide 5; found position {source_index + 1}")

    new_slide = presentation.slides.add_slide(source.slide_layout)
    new_slide.background.fill.solid()
    new_slide.background.fill.fore_color.rgb = source.background.fill.fore_color.rgb
    for shape in source.shapes:
        new_slide.shapes._spTree.insert_element_before(
            deepcopy(shape.element), "p:extLst"
        )

    # Place the newly created slide immediately after its solid-bar source.
    slide_id_list = presentation.slides._sldIdLst
    new_slide_id = slide_id_list[-1]
    slide_id_list.remove(new_slide_id)
    slide_id_list.insert(source_index + 1, new_slide_id)

    replace_solid_bars_with_stacks(new_slide)

    # Shared legend in the space previously used by the reference caption.
    caption = shape_by_id(new_slide, 42)
    caption.left = Inches(0.94)
    caption.top = Inches(5.68)
    caption.width = Inches(0.95)
    caption.height = Inches(0.24)
    set_textbox(caption, MARKER, size=9.2, text_color=NAVY)
    add_pill(new_slide, "SHORT", 1.93, 5.68, 0.67, 0.24, CORAL, WHITE, 7.4)
    add_pill(new_slide, "MEDIUM", 2.68, 5.68, 0.82, 0.24, GOLD, NAVY, 7.4)
    add_pill(new_slide, "LONG", 3.58, 5.68, 0.66, 0.24, PURPLE, WHITE, 7.4)

    mixed_note = shape_by_id(new_slide, 61)
    set_textbox(
        mixed_note,
        "P1: 8 = 4 short + 3 medium + 1 long",
        size=8.1,
        text_color=SLATE,
        bold=True,
    )

    takeaway = shape_by_id(new_slide, 79)
    set_textbox(
        takeaway,
        "SAME PEAK TOTALS  •  EXTRA CLUE: SEGMENT-LENGTH COMPOSITION",
        size=9.7,
        text_color=WHITE,
    )

    try:
        new_slide.notes_slide.notes_text_frame.text = (
            "This slide repeats the same two-cell Bayesian deconvolution example as the previous "
            "slide, but splits every peak count into segment-length bins. Coral is short, gold is "
            "medium, and purple is long. The total bar heights remain T = [10, 2, 8, 4], "
            "B = [2, 10, 4, 8], and mixed = [8, 4, 7, 5]. The mixed profile is still exactly "
            "75% T plus 25% B within every individual length bin. For example, mixed Peak 1 is "
            "8 counts = 4 short + 3 medium + 1 long. The colored pieces do not add counts; they "
            "retain extra fragment-length information hidden by an ordinary peak-total table."
        )
    except Exception:
        pass

    return new_slide


def payload_hashes(path: Path, prefix: str) -> dict[str, str]:
    with ZipFile(path) as archive:
        return {
            name: sha256(archive.read(name)).hexdigest()
            for name in archive.namelist()
            if name.startswith(prefix) and name.endswith(".xml")
        }


def assemble_without_rewriting_existing_slides(
    original_path: Path, generated_path: Path, output_path: Path
) -> None:
    replaced = {
        "[Content_Types].xml",
        "ppt/presentation.xml",
        "ppt/_rels/presentation.xml.rels",
        "docProps/app.xml",
    }
    new_parts = {
        "ppt/slides/slide25.xml",
        "ppt/slides/_rels/slide25.xml.rels",
        "ppt/notesSlides/notesSlide25.xml",
        "ppt/notesSlides/_rels/notesSlide25.xml.rels",
    }

    with ZipFile(original_path, "r") as original, ZipFile(
        generated_path, "r"
    ) as generated, ZipFile(output_path, "w", compression=ZIP_DEFLATED) as output:
        original_names = set(original.namelist())
        generated_names = set(generated.namelist())
        missing = (replaced | new_parts) - generated_names
        if missing:
            raise RuntimeError(f"Generated package is missing required parts: {sorted(missing)}")
        unexpected_existing = new_parts & original_names
        if unexpected_existing:
            raise RuntimeError(f"Target names already exist: {sorted(unexpected_existing)}")

        for info in original.infolist():
            data = generated.read(info.filename) if info.filename in replaced else original.read(info.filename)
            output.writestr(info, data)
        for name in sorted(new_parts):
            output.writestr(generated.getinfo(name), generated.read(name))


def main() -> None:
    before_slides = payload_hashes(DECK, "ppt/slides/slide")
    before_notes = payload_hashes(DECK, "ppt/notesSlides/notesSlide")

    with TemporaryDirectory(prefix="shapemix_length_bins_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"

        presentation = Presentation(DECK)
        build_variant_in_memory(presentation)
        presentation.save(generated)
        assemble_without_rewriting_existing_slides(DECK, generated, assembled)

        check = Presentation(assembled)
        if len(check.slides) != 25:
            raise RuntimeError(f"Expected 25 slides; found {len(check.slides)}")
        sixth_text = slide_text(check.slides[5])
        if TITLE not in sixth_text or MARKER not in sixth_text:
            raise RuntimeError("The inserted slide is not in position 6")

        after_slides = payload_hashes(assembled, "ppt/slides/slide")
        after_notes = payload_hashes(assembled, "ppt/notesSlides/notesSlide")
        for name, digest in before_slides.items():
            if after_slides.get(name) != digest:
                raise RuntimeError(f"Existing slide payload changed: {name}")
        for name, digest in before_notes.items():
            if after_notes.get(name) != digest:
                raise RuntimeError(f"Existing notes payload changed: {name}")

        assembled.replace(DECK)

    print(f"Inserted length-binned clone after slide 5: {DECK}")


if __name__ == "__main__":
    main()
