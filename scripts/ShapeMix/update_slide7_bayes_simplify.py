#!/usr/bin/env python3
"""Rebuild the Bayes slide (position 7) to its current approved design.

This reproduces slide 7 exactly as reviewed on Aug 30, 2026, including the
manual font pass done in PowerPoint:
- title: "Bayesian model to find the number of cells for each cell type"
  (24 pt Aptos Display);
- plain-language definitions of z and N at 14 pt (no dimensions);
- compact formula p(z | N) proportional-to p(N | z) x p(z) with small
  colored term labels;
- one-line callouts for posterior / likelihood / prior at 14 pt;
- no MAP mentions, no takeaway bar, no footer source text (hairline only).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from insert_count_only_bayesian_slides import (
    CREAM,
    FONT,
    FONT_DISPLAY,
    INK,
    KICKER,
    MID,
    NAVY,
    PURPLE,
    PURPLE_PALE,
    SLATE,
    TEAL,
    TEAL_PALE,
    GOLD_PALE,
    C,
    add_card,
    add_shape,
    add_text,
    set_bg,
)
from update_slide9_bayes_compact import DARK_GOLD, add_line, slide_texts

REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

SLIDE_TITLE = "Bayesian model to find the number of cells for each cell type"


def add_definition(slide, symbol, text, x, y):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(9.3), Inches(0.337))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = f"{symbol}   "
    r1.font.name = FONT
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = C(NAVY)
    r2 = p.add_run()
    r2.text = text
    r2.font.name = FONT
    r2.font.size = Pt(14)
    r2.font.color.rgb = C(SLATE)
    return box


def rebuild(slide) -> None:
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)

    set_bg(slide, CREAM)
    add_text(slide, KICKER, 0.68, 0.28, 8.8, 0.28, size=10.5, color=TEAL, bold=True)
    add_text(
        slide, SLIDE_TITLE, 0.65, 0.55, 11.8, 0.62,
        size=24, color=NAVY, bold=True, font=FONT_DISPLAY,
    )

    add_definition(
        slide, "z",
        "number of cells for each cell type in the spot — what we want to find",
        2.0, 1.202,
    )
    add_definition(
        slide, "N",
        "the cut-site counts actually observed at the peaks — what we see",
        2.015, 1.546,
    )

    terms = [
        ("p(z | N)", TEAL, 2.80, 2.30, "posterior", 3.95),
        ("∝", NAVY, 5.20, 0.50, None, None),
        ("p(N | z)", DARK_GOLD, 5.80, 2.30, "likelihood", 6.95),
        ("×", NAVY, 8.20, 0.50, None, None),
        ("p(z)", PURPLE, 8.80, 1.70, "prior", 9.65),
    ]
    for text, color, x, w, label, center in terms:
        add_text(
            slide, text, x, 2.30, w, 0.68,
            size=30, color=color, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
        )
        if label:
            add_text(
                slide, label, center - 0.9, 3.04, 1.8, 0.28,
                size=11, color=color, bold=True, align=PP_ALIGN.CENTER,
            )

    callouts = [
        (0.70, TEAL, TEAL_PALE, "POSTERIOR",
         "best supported z after seeing data N", 0.316,
         3.95, 2.575),
        (4.79, DARK_GOLD, GOLD_PALE, "LIKELIHOOD",
         "given z, what is the probability of seeing data N", 0.552,
         6.95, 6.665),
        (8.88, PURPLE, PURPLE_PALE, "PRIOR",
         "how plausible z is before seeing any data", 0.552,
         9.65, 10.755),
    ]
    y, w, h = 4.00, 3.75, 1.15
    for x, color, pale, heading, body, body_h, term_center, card_center in callouts:
        add_line(slide, term_center, 3.36, card_center, y, color)
        add_card(slide, x, y, w, h, fill=pale, line=color)
        add_text(slide, heading, x + 0.16, y + 0.10, w - 0.32, 0.26,
                 size=11, color=color, bold=True)
        add_text(slide, body, x + 0.16, y + 0.42, w - 0.32, body_h,
                 size=14, color=INK)

    # Bottom hairline only; the footer source text was removed by review.
    hairline = add_shape(slide, MSO_SHAPE.RECTANGLE, 0.68, 7.13, 11.97, 0.011, MID)
    hairline.line.width = Pt(0.25)


def main() -> None:
    presentation = Presentation(DECK)
    matches = [
        index
        for index, slide in enumerate(presentation.slides)
        if any(SLIDE_TITLE in text for text in slide_texts(slide))
    ]
    if matches != [6]:
        raise RuntimeError(f"Expected the Bayes slide only at position 7; found {matches}")

    rebuild(presentation.slides[6])
    presentation.save(DECK)

    check = Presentation(DECK)
    joined = " ".join(t for t in slide_texts(check.slides[6]) if t.strip())
    if "MAP" in joined or "Model reference" in joined or "ONE SCORE" in joined:
        raise RuntimeError("Slide 7 contains removed content")
    for needle in (
        SLIDE_TITLE,
        "posterior",
        "p(z | N)",
        "best supported z after seeing data N",
        "given z, what is the probability of seeing data N",
        "how plausible z is before seeing any data",
        "number of cells for each cell type in the spot",
        "what we see",
    ):
        if needle not in joined:
            raise RuntimeError(f"Slide 7 is missing {needle!r}")
    print(f"Rebuilt Bayes slide at position 7 (current approved design): {DECK}")


if __name__ == "__main__":
    main()
