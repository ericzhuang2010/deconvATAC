#!/usr/bin/env python3
"""Insert a likelihood slide immediately after slide 8 (the prior slide).

Content follows Section 1.4 of
docs/ShapeMix/tutorials/ShapeMix_ATAC_Bayesian_simple.pdf:
the negative binomial N ~ NegBinomial(mean n = z . R,
inverse-dispersion = phi_ref * sum z), introducing the symbols R and n.

Existing slides are not modified; the new slide is added and moved into
position 9.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from insert_count_only_bayesian_slides import (
    CREAM,
    FONT,
    FONT_DISPLAY,
    KICKER,
    MID,
    NAVY,
    SLATE,
    TEAL,
    GOLD_PALE,
    C,
    add_bullets,
    add_card,
    add_notes,
    add_pill,
    add_text,
    set_bg,
)
from insert_prior_slide import PRIOR_TITLE
from update_slide9_bayes_compact import DARK_GOLD, slide_texts

REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

LIKELIHOOD_TITLE = "The likelihood: the probability of seeing N, given z"


def add_definition(slide, symbol, text, x, y):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(9.9), Inches(0.337))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = f"{symbol}   "
    r1.font.name = FONT
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = C(NAVY)
    r2 = p.add_run()
    r2.text = text
    r2.font.name = FONT
    r2.font.size = Pt(14)
    r2.font.color.rgb = C(SLATE)
    return box


def build(slide) -> None:
    set_bg(slide, CREAM)
    add_text(slide, KICKER, 0.68, 0.28, 8.8, 0.28, size=10.5, color=TEAL, bold=True)
    add_text(
        slide, LIKELIHOOD_TITLE, 0.65, 0.55, 11.8, 0.62,
        size=24, color=NAVY, bold=True, font=FONT_DISPLAY,
    )

    add_definition(
        slide, "R",
        "reference peak rate — expected cut sites per cell of each type, learned "
        "from labeled reference cells",
        2.0, 1.202,
    )
    add_definition(
        slide, "n",
        "predicted (mean) counts — what a candidate recipe z predicts: n = z · R",
        2.015, 1.546,
    )

    # Formula card.
    add_card(slide, 0.7, 2.1, 11.93, 1.75, fill=GOLD_PALE, line=DARK_GOLD)
    add_pill(slide, "LIKELIHOOD  p(N | z)", 0.95, 2.3, 2.35, 0.32, DARK_GOLD)
    add_text(
        slide,
        "N  ~  NegBinomial( mean = n ,   inverse-dispersion = φref · Σ z )",
        0.95, 2.78, 11.45, 0.5,
        size=21, color=NAVY, bold=True, align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "n = z · R is a dot product: each cell type's amount × its reference rate, "
        "added over the cell types",
        0.95, 3.36, 11.45, 0.32,
        size=12.5, color=SLATE, align=PP_ALIGN.CENTER,
    )

    add_text(slide, "How to read the formula", 0.71, 4.1, 5.9, 0.3,
             size=14, color=NAVY, bold=True)
    add_bullets(
        slide,
        [
            "n is the prediction: the counts we would expect if the spot truly had "
            "recipe z.",
            "N rarely equals n exactly — sequencing counts fluctuate around the "
            "prediction.",
            "φref is a fixed constant estimated from reference cells; multiplied by "
            "Σ z, it sets how much spread is allowed.",
        ],
        0.71, 4.46, 5.85, 2.2, size=12, bullet_color=DARK_GOLD,
    )

    add_text(slide, "Why the negative binomial?", 6.85, 4.1, 5.8, 0.3,
             size=14, color=NAVY, bold=True)
    add_bullets(
        slide,
        [
            "N is a count — a non-negative whole number — so it needs a count "
            "distribution.",
            "A Poisson model forces the spread to equal the mean; real sequencing "
            "counts vary more than that (overdispersion).",
            "The negative binomial adds a controllable extra-spread term — its "
            "variance is n + n² / (φref · Σ z) — so noisy real data fit honestly.",
        ],
        6.85, 4.46, 5.8, 2.2, size=12, bullet_color=DARK_GOLD,
    )

    hairline = add_card(slide, 0.68, 7.13, 11.97, 0.011, fill=MID, line=MID)
    hairline.line.width = Pt(0.25)

    add_notes(
        slide,
        "The likelihood is the data-fit term of Bayes' rule from slide 7. R is the "
        "fixed reference peak rate learned from labeled reference cells; n = z · R "
        "is the predicted mean count for a candidate recipe. The observed N is "
        "modeled as a negative binomial draw centered on n, with spread controlled "
        "by the fixed phi_ref times the total abundance. The negative binomial is "
        "chosen because N is a count and real sequencing counts are overdispersed: "
        "a Poisson would force variance = mean, while the negative binomial allows "
        "variance n + n^2/(phi_ref * sum z). Source: "
        "ShapeMix_ATAC_Bayesian_simple.pdf, Section 1.4.",
    )


def main() -> None:
    presentation = Presentation(DECK)
    for slide in presentation.slides:
        if any(LIKELIHOOD_TITLE in text for text in slide_texts(slide)):
            raise RuntimeError("The likelihood slide already exists; refusing to duplicate")

    prior_positions = [
        index
        for index, slide in enumerate(presentation.slides)
        if any(PRIOR_TITLE in text for text in slide_texts(slide))
    ]
    if prior_positions != [7]:
        raise RuntimeError(f"Expected the prior slide only at position 8; found {prior_positions}")

    layout = presentation.slides[7].slide_layout
    new_slide = presentation.slides.add_slide(layout)
    build(new_slide)

    slide_id_list = presentation.slides._sldIdLst
    new_id = slide_id_list[-1]
    slide_id_list.remove(new_id)
    slide_id_list.insert(8, new_id)

    presentation.save(DECK)

    check = Presentation(DECK)
    if len(check.slides) != 30:
        raise RuntimeError(f"Expected 30 slides; found {len(check.slides)}")
    texts = " ".join(t for t in slide_texts(check.slides[8]) if t.strip())
    for needle in (LIKELIHOOD_TITLE, "NegBinomial", "reference peak rate",
                   "predicted (mean) counts", "How to read the formula",
                   "Why the negative binomial?"):
        if needle not in texts:
            raise RuntimeError(f"New slide 9 is missing {needle!r}")
    if not any(PRIOR_TITLE in t for t in slide_texts(check.slides[7])):
        raise RuntimeError("Slide 8 is no longer the prior slide")
    if not any("Four symbols" in t for t in slide_texts(check.slides[9])):
        raise RuntimeError("Slide 10 is no longer the symbols slide")
    print(f"Inserted the likelihood slide at position 9: {DECK}")


if __name__ == "__main__":
    main()
