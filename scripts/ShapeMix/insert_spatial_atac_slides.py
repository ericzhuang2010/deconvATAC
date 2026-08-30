#!/usr/bin/env python3
"""Insert two spatial-ATAC slides between slide 4 and slide 5.

New slide A ("Spatial ATAC") introduces the assay: the familiar Tn5
chemistry applied to an intact tissue slice sitting on a grid of barcoded
spots. New slide B ("Why deconvolution") shows the resolution catch: a spot
is wider than a cell, so every barcode reports a blended profile — which is
exactly the mixture the following deconvolution slides learn to unmix.

The final package is assembled from the original ZIP so every pre-existing
slide and notes payload remains byte-for-byte unchanged; only the package
metadata parts (presentation.xml and friends) are rewritten, plus the new
slide parts.
"""

from __future__ import annotations

import re
from hashlib import sha256
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

ANCHOR_BEFORE = "ATAC-seq turns accessible DNA into peak counts"
ANCHOR_AFTER = "Deconvolution works like identifying ingredients in a smoothie"
KICKER_A = "SPATIAL ATAC"
KICKER_B = "WHY DECONVOLUTION"

# Shared deck palette (identical to build_shapemix_high_school_deck.py).
NAVY = "13233B"
INK = "213047"
SLATE = "52627A"
TEAL = "18A999"
TEAL_DARK = "0C7F75"
TEAL_PALE = "DDF4F1"
BLUE = "3B82F6"
BLUE_PALE = "E7F0FE"
PURPLE = "8B5CF6"
CORAL = "F26B5B"
GOLD = "F4B942"
GOLD_PALE = "FFF3D6"
CREAM = "F7F4EC"
WHITE = "FFFFFF"
PALE = "F4F7FA"
MID = "D7E0E8"

DNA_COLOR = "654640"
LOBE_COLORS = ("6D43D6", PURPLE, "7650DE", PURPLE)
LOBE_LINE = "5634B5"
ENZYME_BLUE = "68B5E8"

FONT = "Aptos"
FONT_DISPLAY = "Aptos Display"


def C(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#").upper())


def set_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C(color)


def add_shape(slide, kind, x, y, w, h, fill, line=None, line_width=1.0):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(fill)
    shape.line.color.rgb = C(line if line else fill)
    shape.line.width = Pt(line_width)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=20,
    color=INK,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = C(color)
    return box


def add_rich_text(slide, runs, x, y, w, h, size=20, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    for content, kwargs in runs:
        r = p.add_run()
        r.text = content
        r.font.name = kwargs.get("font", FONT)
        r.font.size = Pt(kwargs.get("size", size))
        r.font.bold = kwargs.get("bold", False)
        r.font.color.rgb = C(kwargs.get("color", color))
    return box


def add_line(slide, x1, y1, x2, y2, color=MID, width=1.5):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = C(color)
    line.line.width = Pt(width)
    return line


def add_circle(slide, x, y, d, fill, line=WHITE, line_width=1.5):
    return add_shape(slide, MSO_SHAPE.OVAL, x, y, d, d, fill, line=line, line_width=line_width)


def add_pill(slide, text, x, y, w, h, fill, color=WHITE, size=13, line=None):
    shp = add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill,
        line=line if line else fill,
        line_width=1,
    )
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = C(color)
    return shp


def add_card(slide, x, y, w, h, fill=WHITE, line=MID, line_width=1.0):
    return add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line=line, line_width=line_width
    )


def add_title(slide, title, kicker):
    add_text(slide, kicker.upper(), 0.68, 0.28, 8.8, 0.28, size=10.5, color=TEAL, bold=True)
    add_text(
        slide,
        title,
        0.65,
        0.55,
        11.8,
        0.62,
        size=27,
        color=NAVY,
        bold=True,
        font=FONT_DISPLAY,
    )


def add_footer(slide, source):
    add_line(slide, 0.68, 7.13, 12.65, 7.13, color=MID, width=0.8)
    add_text(slide, source, 0.7, 7.17, 11.8, 0.18, size=7.5, color=SLATE)


def add_notes(slide, text):
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        pass


def add_nucleosome(slide, center_x, center_y, scale=1.0):
    """Histone core with a visible DNA wrap — same drawing as slides 2–3."""
    lobe_w = 0.18 * scale
    lobe_h = 0.48 * scale
    for offset, color in zip((-0.22, -0.075, 0.075, 0.22), LOBE_COLORS):
        lobe = add_shape(
            slide,
            MSO_SHAPE.OVAL,
            center_x + offset * scale - lobe_w / 2,
            center_y - lobe_h / 2,
            lobe_w,
            lobe_h,
            color,
            line=LOBE_LINE,
            line_width=0.6,
        )
        lobe.rotation = 18
    wrap = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(center_x - 0.38 * scale),
        Inches(center_y - 0.31 * scale),
        Inches(0.76 * scale),
        Inches(0.62 * scale),
    )
    wrap.fill.background()
    wrap.line.color.rgb = C(DNA_COLOR)
    wrap.line.width = Pt(1.45)
    wrap.rotation = 9
    add_line(
        slide,
        center_x - 0.34 * scale,
        center_y + 0.12 * scale,
        center_x + 0.34 * scale,
        center_y - 0.10 * scale,
        color=DNA_COLOR,
        width=1.25,
    )
    return wrap


def add_tn5_marker(slide, x, y):
    """Two-lobed Tn5 enzyme with a coral cut tick — same drawing as slides 2–3."""
    add_circle(slide, x - 0.11, y - 0.14, 0.16, ENZYME_BLUE, line=WHITE, line_width=0.6)
    add_circle(slide, x - 0.01, y - 0.10, 0.16, ENZYME_BLUE, line=WHITE, line_width=0.6)
    add_line(slide, x + 0.02, y + 0.01, x + 0.02, y + 0.15, color=CORAL, width=1.2)


def add_atac_fragment(slide, x, y, w, color):
    """Fragment with adapter caps and endpoint dots — same drawing as slide 3."""
    add_line(slide, x, y, x + w, y, color=color, width=3.8)
    cap_w = min(0.12, w * 0.24)
    add_line(slide, x, y, x + cap_w, y, color=CORAL, width=3.8)
    add_line(slide, x + w - cap_w, y, x + w, y, color=CORAL, width=3.8)
    for end_x in (x, x + w):
        add_circle(slide, end_x - 0.035, y - 0.035, 0.07, NAVY, line=NAVY, line_width=0.3)


def add_cell(slide, x, y, d, color):
    """Cell circle with the small white highlight used on the tissue slides."""
    add_circle(slide, x, y, d, color, line=WHITE, line_width=2)
    add_circle(slide, x + d * 0.34, y + d * 0.33, d * 0.27, WHITE, line=WHITE, line_width=0.5)


# ---------------------------------------------------------------------------
# Slide A — what spatial ATAC is
# ---------------------------------------------------------------------------

def build_slide_a(slide):
    set_bg(slide, CREAM)
    add_title(slide, "Spatial ATAC maps open chromatin across a tissue", KICKER_A)
    add_text(
        slide,
        "Same Tn5 chemistry as ATAC-seq — plus a barcode that records where each fragment came from.",
        0.72,
        1.14,
        11.84,
        0.44,
        size=19,
        color=SLATE,
        align=PP_ALIGN.CENTER,
    )

    cards = [
        (0.72, BLUE_PALE, BLUE, "1 · BARCODED TISSUE", BLUE, WHITE),
        (4.89, TEAL_PALE, TEAL, "2 · TN5 CUTS IN PLACE", TEAL, WHITE),
        (9.06, GOLD_PALE, GOLD, "3 · FRAGMENT + ADDRESS", GOLD, NAVY),
    ]
    for x, fill, line, pill_text, pill_fill, pill_text_color in cards:
        add_card(slide, x, 1.86, 3.77, 4.30, fill=fill, line=line)
        add_pill(slide, pill_text, x + 0.26, 2.08, 2.10, 0.34, pill_fill, color=pill_text_color, size=9.5)

    # Chevrons between the steps, as on the algorithm-overview slide.
    for chev_x in (4.53, 8.70):
        add_shape(slide, MSO_SHAPE.CHEVRON, chev_x, 3.80, 0.32, 0.55, TEAL, line=TEAL)

    # Card 1: tissue slice on a grid of barcoded spots (slide-7 grid vocabulary).
    x0 = 0.72
    add_text(
        slide,
        "A thin tissue slice sits on a grid of barcoded spots",
        x0 + 0.25,
        2.56,
        3.27,
        0.62,
        size=15,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    grid_colors = [
        BLUE, BLUE, TEAL, TEAL,
        PURPLE, BLUE, CORAL, TEAL,
        PURPLE, CORAL, CORAL, GOLD,
        PURPLE, PURPLE, GOLD, GOLD,
    ]
    k = 0
    for row in range(4):
        for col in range(4):
            gx = x0 + 0.785 + col * 0.60
            gy = 3.32 + row * 0.60
            add_circle(slide, gx, gy, 0.40, grid_colors[k], line=WHITE, line_width=2)
            k += 1
    add_text(
        slide,
        "each circle = one spot with its own DNA barcode",
        x0 + 0.25,
        5.62,
        3.27,
        0.44,
        size=10.5,
        color=SLATE,
        align=PP_ALIGN.CENTER,
    )

    # Card 2: the familiar chromatin + Tn5 drawing, working inside the tissue.
    x1 = 4.89
    add_text(
        slide,
        "Inside every spot, Tn5 tags accessible DNA",
        x1 + 0.25,
        2.56,
        3.27,
        0.62,
        size=15,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_line(slide, x1 + 0.30, 4.30, x1 + 3.47, 4.30, color=DNA_COLOR, width=1.7)
    for center_x in (x1 + 0.85, x1 + 2.95):
        add_nucleosome(slide, center_x, 4.30, scale=0.8)
    for marker_x in (x1 + 1.70, x1 + 2.25):
        add_tn5_marker(slide, marker_x, 4.30)
    add_pill(slide, "Tn5", x1 + 1.70, 3.62, 0.56, 0.24, BLUE, size=8.2)
    add_text(
        slide,
        "the exact same cut-and-tag chemistry as bulk ATAC-seq",
        x1 + 0.25,
        5.62,
        3.27,
        0.44,
        size=10.5,
        color=SLATE,
        align=PP_ALIGN.CENTER,
    )

    # Card 3: sequenced fragments carry their spot's address.
    x2 = 9.06
    add_text(
        slide,
        "Each fragment is sequenced with its spot's barcode",
        x2 + 0.25,
        2.56,
        3.27,
        0.62,
        size=15,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    fragments = [
        (0.60, CORAL, "A2", 3.70),
        (1.05, GOLD, "B7", 4.30),
        (1.50, PURPLE, "C5", 4.90),
    ]
    for width, color, barcode, y in fragments:
        add_atac_fragment(slide, x2 + 0.45, y, width, color)
        add_pill(slide, barcode, x2 + 0.45 + width + 0.18, y - 0.11, 0.62, 0.22, NAVY, size=7.5)
    add_text(
        slide,
        "sorting by barcode gives peak counts at every tissue location",
        x2 + 0.25,
        5.58,
        3.27,
        0.50,
        size=10.5,
        color=SLATE,
        align=PP_ALIGN.CENTER,
    )

    add_pill(
        slide,
        "RESULT: AN OPEN-CHROMATIN MAP — ATAC PEAK COUNTS AT EVERY SPOT IN THE TISSUE",
        2.42,
        6.42,
        8.50,
        0.40,
        TEAL,
        size=10.5,
    )
    add_footer(slide, "Source: Deng et al., Nature (2022), spatial-ATAC-seq. Diagram is conceptual.")
    add_notes(
        slide,
        "Spatial ATAC applies the same Tn5 cut-and-tag chemistry the audience already saw, but "
        "on an intact tissue slice placed over a grid of barcoded spots. Every fragment released "
        "inside a spot is labeled with that spot's DNA barcode, so after sequencing the fragments "
        "can be sorted by location. The output is an ATAC peak-count table for every spot — an "
        "open-chromatin map of the tissue. This slide introduces the technology; the next slide "
        "explains the resolution catch that motivates deconvolution.",
    )


# ---------------------------------------------------------------------------
# Slide B — why deconvolution is needed
# ---------------------------------------------------------------------------

def build_slide_b(slide):
    set_bg(slide, WHITE)
    add_title(slide, "Spots are bigger than cells, so each spot reports a blend", KICKER_B)
    add_text(
        slide,
        "A barcoded spot is about 50 µm across; a cell is about 10 µm. Every cell under a spot shares the same barcode.",
        0.72,
        1.14,
        11.84,
        0.46,
        size=18.5,
        color=SLATE,
        align=PP_ALIGN.CENTER,
    )

    # Left: scale picture — one spot covering many cells.
    add_card(slide, 0.72, 1.86, 5.30, 4.60, fill=PALE, line=MID)
    add_text(slide, "One spot under the microscope", 0.98, 2.06, 4.80, 0.34, size=18, color=NAVY, bold=True)
    spot = add_shape(
        slide,
        MSO_SHAPE.OVAL,
        1.945,
        2.675,
        2.85,
        2.85,
        TEAL_PALE,
        line=TEAL,
        line_width=2,
    )
    spot.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    nuclei = [
        (2.62, 3.42, TEAL),
        (3.32, 3.30, BLUE),
        (4.02, 3.52, TEAL),
        (2.42, 4.12, GOLD),
        (3.12, 4.02, TEAL),
        (3.87, 4.15, CORAL),
        (2.72, 4.82, BLUE),
        (3.42, 4.72, TEAL),
        (4.07, 4.72, PURPLE),
    ]
    for cx, cy, color in nuclei:
        add_cell(slide, cx - 0.26, cy - 0.26, 0.52, color)
    add_text(slide, "one cell\n≈ 10 µm", 4.72, 3.02, 1.28, 0.44, size=10, color=SLATE, bold=True)
    add_line(slide, 4.78, 3.34, 4.24, 3.58, color=SLATE, width=1.0)
    # Scale bar under the spot.
    add_line(slide, 1.945, 5.72, 4.795, 5.72, color=SLATE, width=1.2)
    for tick_x in (1.945, 4.795):
        add_line(slide, tick_x, 5.66, tick_x, 5.78, color=SLATE, width=1.2)
    add_text(
        slide,
        "spot ≈ 50 µm",
        2.60,
        5.78,
        1.55,
        0.26,
        size=11,
        color=SLATE,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0.01,
    )
    add_pill(slide, "OFTEN ~10 CELLS SHARE ONE BARCODE", 1.42, 6.06, 3.90, 0.32, NAVY, size=9.5)

    # Right: consequence — one blended ATAC profile per barcode. The purple
    # outline echoes the "unknown mixed spot" card on the next slide, and the
    # blended counts [8, 4, 7, 5] are exactly the profile that slide solves.
    add_card(slide, 6.42, 1.86, 6.19, 4.60, fill=WHITE, line=PURPLE, line_width=1.2)
    add_text(slide, "What one spot reports", 6.70, 2.06, 3.40, 0.34, size=18, color=NAVY, bold=True)
    add_pill(slide, "ONE BLENDED PROFILE PER BARCODE", 9.55, 2.10, 2.85, 0.32, PURPLE, size=8.5)

    cell_row = [(7.05, TEAL, "T"), (7.55, TEAL, "T"), (8.05, TEAL, "T"), (8.72, BLUE, "B")]
    for cx, color, label in cell_row:
        add_cell(slide, cx - 0.20, 2.62, 0.40, color)
        add_text(
            slide,
            label,
            cx - 0.20,
            2.62,
            0.40,
            0.38,
            size=11,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0,
        )
    add_text(slide, "+", 8.31, 2.66, 0.26, 0.32, size=16, color=SLATE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_text(
        slide,
        "3 T cells + 1 B cell caught under the same spot",
        9.10,
        2.60,
        3.35,
        0.55,
        size=12.5,
        color=NAVY,
        bold=True,
    )
    add_shape(slide, MSO_SHAPE.DOWN_ARROW, 7.13, 3.35, 0.34, 0.40, TEAL, line=TEAL)

    baseline, unit = 5.42, 0.22
    values = (8, 4, 7, 5)
    for i, value in enumerate(values):
        bar_x = 7.65 + i * 1.04
        height = value * unit
        add_shape(slide, MSO_SHAPE.RECTANGLE, bar_x, baseline - height, 0.62, height, GOLD, line=GOLD)
        add_text(
            slide,
            str(value),
            bar_x - 0.08,
            baseline - height - 0.24,
            0.78,
            0.20,
            size=11,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
        add_text(
            slide,
            f"P{i + 1}",
            bar_x - 0.08,
            baseline + 0.05,
            0.78,
            0.18,
            size=9.5,
            color=SLATE,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
    add_line(slide, 7.55, baseline, 11.50, baseline, color=NAVY, width=0.8)
    add_text(
        slide,
        "Which cells contributed what is hidden in this blend",
        6.70,
        5.92,
        5.65,
        0.34,
        size=14,
        color=CORAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_card(slide, 0.72, 6.58, 11.90, 0.44, fill=NAVY, line=NAVY)
    add_rich_text(
        slide,
        [
            ("Un-mixing each spot back into a cell-type recipe is ", {"bold": True, "color": WHITE}),
            ("DECONVOLUTION", {"bold": True, "color": GOLD}),
            (" — the next slides show how.", {"bold": True, "color": WHITE}),
        ],
        1.00,
        6.63,
        11.34,
        0.32,
        size=13.5,
        align=PP_ALIGN.CENTER,
    )
    add_footer(
        slide,
        "Sources: Deng et al., Nature (2022); Ouologuem et al., Bioinformatics (2025). Sizes and counts are illustrative.",
    )
    add_notes(
        slide,
        "The catch is resolution: a barcoded spot is roughly 50 micrometers across while a cell "
        "is roughly 10 micrometers, so one barcode usually collects fragments from several cells "
        "— often around ten, matching the pseudo-spot design used later in the benchmark. The "
        "example shows 3 T cells and 1 B cell under one spot producing a single summed profile "
        "of [8, 4, 7, 5]; these are exactly the numbers the following deconvolution slide "
        "solves, recovering the recipe 75% T + 25% B. Because per-cell identities are hidden in "
        "the sum, spot-level analysis needs deconvolution.",
    )


# ---------------------------------------------------------------------------
# Insertion machinery (same byte-preserving approach as the earlier insert)
# ---------------------------------------------------------------------------

def slide_texts(slide) -> list[str]:
    return [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]


def build_variant_in_memory(presentation: Presentation):
    slides = presentation.slides
    for slide in slides:
        joined = " ".join(slide_texts(slide))
        if KICKER_A in joined or KICKER_B in joined:
            raise RuntimeError("Spatial ATAC slides already exist; refusing to add duplicates")

    if ANCHOR_BEFORE not in " ".join(slide_texts(slides[3])):
        raise RuntimeError("Slide 4 is not the expected ATAC-seq basics slide")
    if ANCHOR_AFTER not in " ".join(slide_texts(slides[4])):
        raise RuntimeError("Slide 5 is not the expected deconvolution smoothie slide")

    blank = presentation.slide_layouts[6]
    slide_id_list = presentation.slides._sldIdLst

    slide_a = presentation.slides.add_slide(blank)
    build_slide_a(slide_a)
    sld_id = slide_id_list[-1]
    slide_id_list.remove(sld_id)
    slide_id_list.insert(4, sld_id)

    slide_b = presentation.slides.add_slide(blank)
    build_slide_b(slide_b)
    sld_id = slide_id_list[-1]
    slide_id_list.remove(sld_id)
    slide_id_list.insert(5, sld_id)


def payload_hashes(path: Path, prefix: str) -> dict[str, str]:
    with ZipFile(path) as archive:
        return {
            name: sha256(archive.read(name)).hexdigest()
            for name in archive.namelist()
            if name.startswith(prefix) and name.endswith(".xml")
        }


def assemble_without_rewriting_existing_slides(
    original_path: Path, generated_path: Path, output_path: Path
) -> None:
    replaced = {
        "[Content_Types].xml",
        "ppt/presentation.xml",
        "ppt/_rels/presentation.xml.rels",
        "docProps/app.xml",
    }
    new_part_pattern = re.compile(
        r"^ppt/(slides/slide\d+\.xml"
        r"|slides/_rels/slide\d+\.xml\.rels"
        r"|notesSlides/notesSlide\d+\.xml"
        r"|notesSlides/_rels/notesSlide\d+\.xml\.rels)$"
    )

    with ZipFile(original_path, "r") as original, ZipFile(
        generated_path, "r"
    ) as generated, ZipFile(output_path, "w", compression=ZIP_DEFLATED) as output:
        original_names = set(original.namelist())
        generated_names = set(generated.namelist())

        missing = replaced - generated_names
        if missing:
            raise RuntimeError(f"Generated package is missing required parts: {sorted(missing)}")

        new_parts = generated_names - original_names
        bad = [name for name in new_parts if not new_part_pattern.match(name)]
        if bad:
            raise RuntimeError(f"Unexpected new parts outside slide/notes namespaces: {sorted(bad)}")
        if len([n for n in new_parts if n.startswith("ppt/slides/slide")]) != 2:
            raise RuntimeError(f"Expected exactly two new slide XML parts; got {sorted(new_parts)}")

        for info in original.infolist():
            data = (
                generated.read(info.filename)
                if info.filename in replaced
                else original.read(info.filename)
            )
            output.writestr(info, data)
        for name in sorted(new_parts):
            output.writestr(generated.getinfo(name), generated.read(name))


def main() -> None:
    before_slides = payload_hashes(DECK, "ppt/slides/slide")
    before_notes = payload_hashes(DECK, "ppt/notesSlides/notesSlide")
    before_count = len(Presentation(DECK).slides)

    with TemporaryDirectory(prefix="shapemix_spatial_atac_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"

        presentation = Presentation(DECK)
        build_variant_in_memory(presentation)
        presentation.save(generated)
        assemble_without_rewriting_existing_slides(DECK, generated, assembled)

        check = Presentation(assembled)
        if len(check.slides) != before_count + 2:
            raise RuntimeError(f"Expected {before_count + 2} slides; found {len(check.slides)}")
        if KICKER_A not in " ".join(slide_texts(check.slides[4])):
            raise RuntimeError("Inserted slide A is not in position 5")
        if KICKER_B not in " ".join(slide_texts(check.slides[5])):
            raise RuntimeError("Inserted slide B is not in position 6")
        if ANCHOR_BEFORE not in " ".join(slide_texts(check.slides[3])):
            raise RuntimeError("Slide 4 moved unexpectedly")
        if ANCHOR_AFTER not in " ".join(slide_texts(check.slides[6])):
            raise RuntimeError("The smoothie slide did not shift to position 7")

        after_slides = payload_hashes(assembled, "ppt/slides/slide")
        after_notes = payload_hashes(assembled, "ppt/notesSlides/notesSlide")
        for name, digest in before_slides.items():
            if after_slides.get(name) != digest:
                raise RuntimeError(f"Existing slide payload changed: {name}")
        for name, digest in before_notes.items():
            if after_notes.get(name) != digest:
                raise RuntimeError(f"Existing notes payload changed: {name}")

        assembled.replace(DECK)

    print(f"Inserted two spatial-ATAC slides after slide 4: {DECK}")


if __name__ == "__main__":
    main()
