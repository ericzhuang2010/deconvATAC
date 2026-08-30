#!/usr/bin/env python3
"""Patch slide 3's manually adjusted ATAC-seq chart without rebuilding the deck."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

EMU_PER_INCH = 914400
BASELINE = 6.23
UNIT_HEIGHT = 0.17

NAVY = "13233B"
WHITE = "FFFFFF"
SHORT = "F26B5B"
MEDIUM = "F4B942"
LONG = "8B5CF6"
FONT = "Aptos"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def inches(value) -> float:
    return value / EMU_PER_INCH


def shape_by_id(slide, shape_id: int):
    matches = [shape for shape in slide.shapes if shape.shape_id == shape_id]
    if len(matches) != 1:
        raise RuntimeError(f"Expected one slide-3 shape with id {shape_id}; found {len(matches)}")
    return matches[0]


def remove_generated_stack_shapes(slide) -> None:
    for shape in list(slide.shapes):
        if shape.name.startswith("Peak count stack "):
            shape._element.getparent().remove(shape._element)


def set_fill_and_line(shape, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(0.8)


def set_geometry(shape, x: float, y: float, width: float, height: float) -> None:
    shape.left = Inches(x)
    shape.top = Inches(y)
    shape.width = Inches(width)
    shape.height = Inches(height)


def replace_first_run(shape, text: str, size: float | None = None) -> None:
    paragraphs = shape.text_frame.paragraphs
    if not paragraphs or not paragraphs[0].runs:
        raise RuntimeError(f"Shape id {shape.shape_id} has no editable text run")
    run = paragraphs[0].runs[0]
    run.text = text
    if size is not None:
        run.font.size = Pt(size)


def replace_all_text(shape, text: str, size: float, color: str = WHITE) -> None:
    """Replace every prior run while retaining the shape's fill and margins."""
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = tf.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = rgb(color)


def add_segment(slide, name: str, x: float, y: float, width: float, height: float, color: str):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    shape.name = f"Peak count stack {name}"
    set_fill_and_line(shape, color)
    return shape


def add_segment_label(
    slide,
    name: str,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    color: str,
    size: float,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    box.name = f"Peak count stack label {name}"
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = tf.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = rgb(color)


def main() -> None:
    presentation = Presentation(DECK)
    if len(presentation.slides) < 3:
        raise RuntimeError("The presentation has fewer than three slides")

    slide = presentation.slides[2]
    slide_text = " ".join(
        shape.text for shape in slide.shapes if hasattr(shape, "text_frame")
    )
    if "ATAC-seq" not in slide_text or "PEAK COUNTS" not in slide_text:
        raise RuntimeError("Slide 3 is not the expected ATAC-seq peak-count slide")

    # Stable IDs from the user's manually adjusted slide. Geometry checks make
    # the patch fail safely if PowerPoint has replaced these chart objects.
    bars = {
        "p1": (shape_by_id(slide, 119), 3.05, 2.30),
        "p2": (shape_by_id(slide, 122), 6.05, 2.60),
        "p3": (shape_by_id(slide, 125), 9.45, 2.35),
    }
    for key, (shape, expected_x, expected_width) in bars.items():
        if abs(inches(shape.left) - expected_x) > 0.04:
            raise RuntimeError(f"{key} bar moved unexpectedly in x")
        if abs(inches(shape.width) - expected_width) > 0.04:
            raise RuntimeError(f"{key} bar width changed unexpectedly")
        if abs(inches(shape.top + shape.height) - BASELINE) > 0.04:
            raise RuntimeError(f"{key} bar no longer ends at the expected baseline")

    remove_generated_stack_shapes(slide)

    # Reuse the three original bars as the bottom stack segment so all manual
    # widths and x-alignment remain untouched.
    p1, p1_x, p1_w = bars["p1"]
    p2, p2_x, p2_w = bars["p2"]
    p3, p3_x, p3_w = bars["p3"]

    # Peak 1: 4 short endpoints at the bottom + 1 medium endpoint on top.
    p1_short_h = 4 * UNIT_HEIGHT
    p1_medium_h = 1 * UNIT_HEIGHT
    set_geometry(p1, p1_x, BASELINE - p1_short_h, p1_w, p1_short_h)
    set_fill_and_line(p1, SHORT)
    p1_medium_y = BASELINE - p1_short_h - p1_medium_h
    add_segment(slide, "P1 medium", p1_x, p1_medium_y, p1_w, p1_medium_h, MEDIUM)

    # Peak 2: 2 medium endpoints at the bottom + 2 long endpoints on top.
    p2_medium_h = 2 * UNIT_HEIGHT
    p2_long_h = 2 * UNIT_HEIGHT
    set_geometry(p2, p2_x, BASELINE - p2_medium_h, p2_w, p2_medium_h)
    set_fill_and_line(p2, MEDIUM)
    p2_long_y = BASELINE - p2_medium_h - p2_long_h
    add_segment(slide, "P2 long", p2_x, p2_long_y, p2_w, p2_long_h, LONG)

    # Peak 3: all 3 endpoints come from medium-length fragments.
    p3_medium_h = 3 * UNIT_HEIGHT
    p3_medium_y = BASELINE - p3_medium_h
    set_geometry(p3, p3_x, p3_medium_y, p3_w, p3_medium_h)
    set_fill_and_line(p3, MEDIUM)

    # Total labels above the stacks.
    for shape_id, total, x, width in (
        (120, 5, p1_x, p1_w),
        (123, 4, p2_x, p2_w),
        (126, 3, p3_x, p3_w),
    ):
        label = shape_by_id(slide, shape_id)
        replace_first_run(label, str(total))
        label.left = Inches(x)
        label.top = Inches(BASELINE - total * UNIT_HEIGHT - 0.24)
        label.width = Inches(width)

    # Segment numbers use the fragment-length colors defined in row 2.
    add_segment_label(slide, "P1 short", "4", p1_x, BASELINE - p1_short_h, p1_w, p1_short_h, WHITE, 12)
    add_segment_label(slide, "P1 medium", "1", p1_x, p1_medium_y, p1_w, p1_medium_h, NAVY, 7.5)
    add_segment_label(slide, "P2 medium", "2", p2_x, BASELINE - p2_medium_h, p2_w, p2_medium_h, NAVY, 10)
    add_segment_label(slide, "P2 long", "2", p2_x, p2_long_y, p2_w, p2_long_h, WHITE, 10)
    add_segment_label(slide, "P3 medium", "3", p3_x, p3_medium_y, p3_w, p3_medium_h, NAVY, 11)

    cue = shape_by_id(slide, 109)
    replace_all_text(cue, "HEIGHT = COUNT  •  COLOR = FRAGMENT LENGTH", size=7.5)

    rule = shape_by_id(slide, 128)
    replace_all_text(
        rule,
        "P1: 4 short + 1 medium = 5   |   P2: 2 medium + 2 long = 4   |   P3: 3 medium = 3",
        size=8.5,
    )

    try:
        notes = slide.notes_slide.notes_text_frame
        note_text = notes.text
        correction = (
            "Stacked endpoint counts after the aligned-fragment adjustment: "
            "Peak 1 has 4 short-fragment endpoints and 1 medium-fragment endpoint (5 total); "
            "Peak 2 has 2 medium and 2 long endpoints (4 total); "
            "Peak 3 has 3 medium endpoints (3 total). Two medium fragments place one end in each "
            "of neighboring peaks, so odd per-peak totals are possible even though every fragment "
            "has two endpoints."
        )
        if "Stacked endpoint counts after the aligned-fragment adjustment" not in note_text:
            notes.text = f"{note_text.rstrip()}\n\n{correction}".strip()
    except Exception:
        pass

    temporary = DECK.with_name(f".{DECK.stem}.stacked-counts.tmp.pptx")
    presentation.save(temporary)
    # Verify the round-tripped package before replacing the requested deck.
    check = Presentation(temporary)
    if len(check.slides) != len(presentation.slides):
        raise RuntimeError("Slide count changed while saving the patched deck")
    temporary.replace(DECK)
    print(f"Updated stacked peak counts on slide 3: {DECK}")


if __name__ == "__main__":
    main()
