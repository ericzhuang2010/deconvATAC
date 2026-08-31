#!/usr/bin/env python3
"""Repair the deck after an external PowerPoint re-save shuffled slides.

Observed damage:
- Position 9 shows the old big-box Bayes slide (should be the compact design).
- Position 12 shows a duplicate Bayes slide; the original THE PROBLEM slide
  was lost. It is restored from the bak5 snapshot.

The external re-save renumbered package parts, so part names no longer match
positions and byte-preserving zip surgery is unreliable. This script works at
the object level and performs a plain full save.
"""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation

from update_slide9_bayes_compact import TITLE as BAYES_TITLE
from update_slide9_bayes_compact import rebuild, slide_texts

REPO_ROOT = Path(__file__).resolve().parents[2]
PRESENTATIONS = REPO_ROOT / "docs" / "ShapeMix" / "presentations"
DECK = PRESENTATIONS / "ShapeMix_High_School_Research_Deck.pptx"
BACKUP = PRESENTATIONS / "ShapeMix_High_School_Research_Deck bak5.pptx"

PROBLEM_KICKER = "THE PROBLEM"
PROBLEM_TITLE = "One spatial spot can contain several cell types"


def clone_from(source, target) -> None:
    for shape in list(target.shapes):
        shape._element.getparent().remove(shape._element)
    target.background.fill.solid()
    target.background.fill.fore_color.rgb = source.background.fill.fore_color.rgb
    for shape in source.shapes:
        target.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")
    notes_text = ""
    if source.has_notes_slide:
        notes_text = source.notes_slide.notes_text_frame.text
    try:
        target.notes_slide.notes_text_frame.text = notes_text
    except Exception:
        pass


def main() -> None:
    backup = Presentation(BACKUP)
    problem_source = backup.slides[8]
    source_texts = slide_texts(problem_source)
    if PROBLEM_KICKER not in source_texts or not any(
        PROBLEM_TITLE in text for text in source_texts
    ):
        raise RuntimeError("bak5 position 9 is not THE PROBLEM slide; aborting")

    presentation = Presentation(DECK)
    if len(presentation.slides) != 28:
        raise RuntimeError(f"Expected 28 slides; found {len(presentation.slides)}")

    bayes_positions = [
        index
        for index, slide in enumerate(presentation.slides)
        if any(BAYES_TITLE in text for text in slide_texts(slide))
    ]
    if bayes_positions != [8, 11]:
        raise RuntimeError(f"Unexpected Bayes slide positions {bayes_positions}; aborting")

    # Position 9: rebuild as the compact-formula design.
    rebuild(presentation.slides[8])

    # Position 12: restore THE PROBLEM from the backup snapshot.
    clone_from(problem_source, presentation.slides[11])

    presentation.save(DECK)

    # Verify from a fresh load.
    check = Presentation(DECK)
    if len(check.slides) != 28:
        raise RuntimeError("Slide count changed during repair")
    nine = slide_texts(check.slides[8])
    if "posterior" not in nine:
        raise RuntimeError("Position 9 is not the compact Bayes design")
    twelve = slide_texts(check.slides[11])
    if PROBLEM_KICKER not in twelve or not any(PROBLEM_TITLE in t for t in twelve):
        raise RuntimeError("Position 12 is not THE PROBLEM slide")
    bayes_after = [
        index
        for index, slide in enumerate(check.slides)
        if any(BAYES_TITLE in text for text in slide_texts(slide))
    ]
    if bayes_after != [8]:
        raise RuntimeError(f"Bayes slide should appear only at position 9; found {bayes_after}")

    print(f"Repaired positions 9 (compact Bayes) and 12 (THE PROBLEM restored): {DECK}")


if __name__ == "__main__":
    main()
