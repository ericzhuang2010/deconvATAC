#!/usr/bin/env python3
"""Insert the expanded ShapeMix likelihood and multinomial slide after slide 14.

This performs the insertion directly in the PPTX package relationship tables,
so all pre-existing slide and notes payloads remain byte-for-byte unchanged,
including decks whose internal slide filenames do not match their positions.
"""

from __future__ import annotations

import re
from hashlib import sha256
from pathlib import Path
from tempfile import TemporaryDirectory
from xml.etree import ElementTree
from zipfile import ZIP_DEFLATED, ZipFile, ZipInfo

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from insert_count_only_bayesian_slides import (
    BLUE,
    BLUE_PALE,
    C,
    CREAM,
    CORAL,
    CORAL_PALE,
    FONT_DISPLAY,
    GOLD,
    GOLD_PALE,
    INK,
    MID,
    NAVY,
    PURPLE,
    PURPLE_PALE,
    SLATE,
    TEAL,
    TEAL_PALE,
    WHITE,
    add_card,
    add_pill,
    add_shape,
    add_text,
    set_bg,
)

REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

INSERT_AFTER_INDEX = 13  # zero-based index of slide 14
PREVIOUS_TITLE = "Bayesian model for ShapeMix"
NEXT_TITLE = "Four symbols describe one spot's counts"
TITLE = "ShapeMix adds a likelihood for the length-bin split"
DARK_GOLD = "B8860B"

PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DOCUMENT_REL_NS = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
)
PACKAGE_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
EXTENDED_PROPS_NS = (
    "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
)
VT_NS = "http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"
SLIDE_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
)
SLIDE_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
)


def add_centered_text(slide, text, x, y, w, h, *, size, color, bold=False):
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_callout(slide, x, heading, body, fill, line):
    add_card(slide, x, 2.63, 3.75, 1.18, fill=fill, line=line)
    add_text(
        slide,
        heading,
        x + 0.16,
        2.73,
        3.43,
        0.26,
        size=11,
        color=line,
        bold=True,
    )
    add_text(
        slide,
        body,
        x + 0.16,
        3.05,
        3.43,
        0.56,
        size=12.5,
        color=INK,
    )


def add_definition(slide, x, symbol, text, fill, line):
    add_card(slide, x, 5.68, 3.75, 0.68, fill=fill, line=line)
    add_text(
        slide,
        symbol,
        x + 0.14,
        5.81,
        0.98,
        0.30,
        size=12,
        color=line,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        text,
        x + 1.14,
        5.77,
        2.43,
        0.39,
        size=10.5,
        color=INK,
        valign=MSO_ANCHOR.MIDDLE,
    )


def build(slide) -> None:
    set_bg(slide, CREAM)
    add_text(
        slide,
        "SHAPEMIX • BAYESIAN MODEL",
        0.68,
        0.28,
        8.8,
        0.28,
        size=10.5,
        color=TEAL,
        bold=True,
    )
    add_text(
        slide,
        TITLE,
        0.65,
        0.55,
        11.8,
        0.62,
        size=26,
        color=NAVY,
        bold=True,
        font=FONT_DISPLAY,
    )

    # Expanded Bayes formula: posterior = length-bin likelihood × total-count
    # likelihood × prior.
    terms = [
        ("p(z | N_LB, N)", TEAL, 0.70, 2.70, "posterior"),
        ("∝", NAVY, 3.40, 0.46, None),
        ("p(N_LB | N,z)", CORAL, 3.86, 2.72, "length-bin likelihood"),
        ("×", NAVY, 6.58, 0.42, None),
        ("p(N | z)", DARK_GOLD, 7.00, 1.90, "total-count likelihood"),
        ("×", NAVY, 8.90, 0.42, None),
        ("p(z)", PURPLE, 9.32, 1.55, "prior"),
    ]
    for text, color, x, width, label in terms:
        add_centered_text(
            slide,
            text,
            x,
            1.34,
            width,
            0.66,
            size=24,
            color=color,
            bold=True,
        )
        if label:
            add_centered_text(
                slide,
                label,
                x,
                2.02,
                width,
                0.24,
                size=10,
                color=color,
                bold=True,
            )

    add_callout(
        slide,
        0.70,
        "POSTERIOR",
        "best supported z after seeing both the total and length-bin data",
        TEAL_PALE,
        TEAL,
    )
    add_callout(
        slide,
        4.79,
        "LIKELIHOOD",
        "two data-fit terms: length-bin split × total peak count",
        GOLD_PALE,
        DARK_GOLD,
    )
    add_callout(
        slide,
        8.88,
        "PRIOR",
        "how plausible z is before seeing any data",
        PURPLE_PALE,
        PURPLE,
    )

    add_text(
        slide,
        "The new length-bin term uses a multinomial distribution",
        0.71,
        4.08,
        8.7,
        0.32,
        size=14,
        color=NAVY,
        bold=True,
    )
    add_card(slide, 0.70, 4.48, 11.92, 0.90, fill=CORAL_PALE, line=CORAL)
    add_centered_text(
        slide,
        "N_LB | N,z  ~  Multinomial( N, n_lb / n )  =  Multinomial( N, z · R_LB / n )",
        0.92,
        4.61,
        11.48,
        0.52,
        size=18,
        color=NAVY,
        bold=True,
    )

    add_definition(
        slide,
        0.70,
        "N",
        "fixed total number of cut sites at the peak",
        BLUE_PALE,
        BLUE,
    )
    add_definition(
        slide,
        4.79,
        "n_lb / n",
        "predicted fractions across the length bins; they sum to 1",
        TEAL_PALE,
        TEAL,
    )
    add_definition(
        slide,
        8.88,
        "N_LB",
        "observed bin counts; they sum to N",
        GOLD_PALE,
        DARK_GOLD,
    )

    add_centered_text(
        slide,
        "Same cut sites, one extra question: how were the N counts split by fragment length?",
        0.70,
        6.55,
        11.92,
        0.32,
        size=12,
        color=SLATE,
        bold=True,
    )

    hairline = add_shape(slide, MSO_SHAPE.RECTANGLE, 0.68, 7.13, 11.97, 0.011, MID)
    hairline.line.width = Pt(0.25)


def slide_texts(slide) -> list[str]:
    return [
        shape.text_frame.text.strip()
        for shape in slide.shapes
        if shape.has_text_frame and shape.text_frame.text.strip()
    ]


def package_hashes(path: Path) -> dict[str, str]:
    with ZipFile(path) as archive:
        return {
            name: sha256(archive.read(name)).hexdigest()
            for name in archive.namelist()
        }


def serialize_xml(root) -> bytes:
    return ElementTree.tostring(
        root,
        encoding="utf-8",
        xml_declaration=True,
    )


def slide_part_at(path: Path, slide_index: int) -> str:
    with ZipFile(path) as archive:
        presentation = ElementTree.fromstring(
            archive.read("ppt/presentation.xml")
        )
        slide_ids = presentation.find(f"{{{PRESENTATION_NS}}}sldIdLst")
        if slide_ids is None or slide_index >= len(slide_ids):
            raise RuntimeError(f"Slide position {slide_index + 1} does not exist")
        relationship_id = slide_ids[slide_index].attrib[
            f"{{{DOCUMENT_REL_NS}}}id"
        ]
        relationships = ElementTree.fromstring(
            archive.read("ppt/_rels/presentation.xml.rels")
        )
        for relationship in relationships.findall(
            f"{{{PACKAGE_REL_NS}}}Relationship"
        ):
            if relationship.attrib.get("Id") == relationship_id:
                target = relationship.attrib["Target"]
                return (
                    target.lstrip("/")
                    if target.startswith("/")
                    else f"ppt/{target}"
                )
    raise RuntimeError(f"Could not resolve slide position {slide_index + 1}")


def next_slide_part(original_names: set[str]) -> str:
    numbers = []
    for name in original_names:
        match = re.fullmatch(r"ppt/slides/slide(\d+)\.xml", name)
        if match:
            numbers.append(int(match.group(1)))
    return f"ppt/slides/slide{max(numbers) + 1}.xml"


def modify_presentation(data: bytes, insert_index: int, relationship_id: str) -> bytes:
    root = ElementTree.fromstring(data)
    slide_ids = root.find(f"{{{PRESENTATION_NS}}}sldIdLst")
    if slide_ids is None:
        raise RuntimeError("Presentation has no slide ID list")
    new_slide_id = max(int(item.attrib["id"]) for item in slide_ids) + 1
    item = ElementTree.Element(f"{{{PRESENTATION_NS}}}sldId")
    item.set("id", str(new_slide_id))
    item.set(f"{{{DOCUMENT_REL_NS}}}id", relationship_id)
    slide_ids.insert(insert_index, item)
    return serialize_xml(root)


def modify_relationships(data: bytes, relationship_id: str, slide_part: str) -> bytes:
    root = ElementTree.fromstring(data)
    relationship = ElementTree.SubElement(
        root,
        f"{{{PACKAGE_REL_NS}}}Relationship",
    )
    relationship.set("Id", relationship_id)
    relationship.set("Type", SLIDE_REL_TYPE)
    relationship.set("Target", f"slides/{Path(slide_part).name}")
    return serialize_xml(root)


def modify_content_types(data: bytes, slide_part: str) -> bytes:
    root = ElementTree.fromstring(data)
    override = ElementTree.SubElement(
        root,
        f"{{{CONTENT_TYPES_NS}}}Override",
    )
    override.set("PartName", f"/{slide_part}")
    override.set("ContentType", SLIDE_CONTENT_TYPE)
    return serialize_xml(root)


def modify_app_metadata(data: bytes, slide_count: int) -> bytes:
    root = ElementTree.fromstring(data)
    slides = root.find(f"{{{EXTENDED_PROPS_NS}}}Slides")
    if slides is not None:
        slides.text = str(slide_count)

    heading_pairs = root.find(f"{{{EXTENDED_PROPS_NS}}}HeadingPairs")
    if heading_pairs is not None:
        vector = heading_pairs.find(f"{{{VT_NS}}}vector")
        if vector is not None:
            children = list(vector)
            for index, child in enumerate(children[:-1]):
                text = child.find(f"{{{VT_NS}}}lpstr")
                if text is not None and text.text == "Slide Titles":
                    count = children[index + 1].find(f"{{{VT_NS}}}i4")
                    if count is not None:
                        count.text = str(slide_count)

    titles = root.find(f"{{{EXTENDED_PROPS_NS}}}TitlesOfParts")
    if titles is not None:
        vector = titles.find(f"{{{VT_NS}}}vector")
        if vector is not None:
            target_size = slide_count + 5  # four fonts plus one theme
            while len(vector) < target_size:
                item = ElementTree.SubElement(vector, f"{{{VT_NS}}}lpstr")
                item.text = "PowerPoint Presentation"
            vector.set("size", str(len(vector)))
    return serialize_xml(root)


def next_relationship_id(relationships_data: bytes) -> str:
    root = ElementTree.fromstring(relationships_data)
    numbers = []
    for relationship in root.findall(f"{{{PACKAGE_REL_NS}}}Relationship"):
        match = re.fullmatch(r"rId(\d+)", relationship.attrib.get("Id", ""))
        if match:
            numbers.append(int(match.group(1)))
    return f"rId{max(numbers) + 1}"


def make_zip_info(source: ZipInfo, filename: str) -> ZipInfo:
    info = ZipInfo(filename, date_time=source.date_time)
    info.compress_type = ZIP_DEFLATED
    info.external_attr = source.external_attr
    info.create_system = source.create_system
    return info


def assemble_with_insert(
    original_path: Path,
    generated_path: Path,
    output_path: Path,
    generated_slide_part: str,
    new_slide_part: str,
    insert_index: int,
    slide_count: int,
) -> None:
    with ZipFile(original_path, "r") as original, ZipFile(
        generated_path, "r"
    ) as generated, ZipFile(output_path, "w", compression=ZIP_DEFLATED) as output:
        original_names = set(original.namelist())
        relationships_data = original.read(
            "ppt/_rels/presentation.xml.rels"
        )
        relationship_id = next_relationship_id(relationships_data)

        replacements = {
            "[Content_Types].xml": modify_content_types(
                original.read("[Content_Types].xml"),
                new_slide_part,
            ),
            "ppt/presentation.xml": modify_presentation(
                original.read("ppt/presentation.xml"),
                insert_index,
                relationship_id,
            ),
            "ppt/_rels/presentation.xml.rels": modify_relationships(
                relationships_data,
                relationship_id,
                new_slide_part,
            ),
            "docProps/app.xml": modify_app_metadata(
                original.read("docProps/app.xml"),
                slide_count,
            ),
        }
        for info in original.infolist():
            output.writestr(
                info,
                replacements.get(info.filename, original.read(info.filename)),
            )

        generated_rels = (
            f"ppt/slides/_rels/{Path(generated_slide_part).name}.rels"
        )
        new_rels = f"ppt/slides/_rels/{Path(new_slide_part).name}.rels"
        slide_info = generated.getinfo(generated_slide_part)
        output.writestr(
            make_zip_info(slide_info, new_slide_part),
            generated.read(generated_slide_part),
        )
        if generated_rels in generated.namelist():
            rels_info = generated.getinfo(generated_rels)
            output.writestr(
                make_zip_info(rels_info, new_rels),
                generated.read(generated_rels),
            )

        if new_slide_part in original_names:
            raise RuntimeError(f"New slide part already exists: {new_slide_part}")


def main() -> None:
    original = Presentation(DECK)
    original_count = len(original.slides)
    if original_count != 33:
        raise RuntimeError(f"Expected 33 slides; found {original_count}")
    for slide in original.slides:
        if any(TITLE in text for text in slide_texts(slide)):
            raise RuntimeError("The multinomial slide already exists")
    if not any(PREVIOUS_TITLE in text for text in slide_texts(original.slides[13])):
        raise RuntimeError("Slide 14 is not the expected ShapeMix model slide")
    if not any(NEXT_TITLE in text for text in slide_texts(original.slides[14])):
        raise RuntimeError("Slide 15 is not the expected symbols slide")

    before = package_hashes(DECK)
    with ZipFile(DECK) as archive:
        new_slide_part = next_slide_part(set(archive.namelist()))

    with TemporaryDirectory(prefix="shapemix_multinomial_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"

        presentation = Presentation(DECK)
        new_slide = presentation.slides.add_slide(
            presentation.slides[13].slide_layout
        )
        build(new_slide)
        presentation.save(generated)
        generated_slide_part = slide_part_at(generated, original_count)

        assemble_with_insert(
            DECK,
            generated,
            assembled,
            generated_slide_part,
            new_slide_part,
            INSERT_AFTER_INDEX + 1,
            original_count + 1,
        )

        check = Presentation(assembled)
        if len(check.slides) != original_count + 1:
            raise RuntimeError(
                f"Expected {original_count + 1} slides; found {len(check.slides)}"
            )
        new_text = " ".join(slide_texts(check.slides[14]))
        for needle in (
            TITLE,
            "p(N_LB | N,z)",
            "p(N | z)",
            "Multinomial",
            "n_lb / n",
            "Same cut sites",
        ):
            if needle not in new_text:
                raise RuntimeError(f"New slide 15 is missing {needle!r}")
        if not any(PREVIOUS_TITLE in text for text in slide_texts(check.slides[13])):
            raise RuntimeError("The original slide 14 was changed")
        if not any(NEXT_TITLE in text for text in slide_texts(check.slides[15])):
            raise RuntimeError("The original slide 15 did not move intact to position 16")

        after = package_hashes(assembled)
        allowed_changes = {
            "[Content_Types].xml",
            "ppt/presentation.xml",
            "ppt/_rels/presentation.xml.rels",
            "docProps/app.xml",
        }
        for name, digest in before.items():
            if name in allowed_changes:
                continue
            if after.get(name) != digest:
                raise RuntimeError(f"Unexpected change to existing package part: {name}")

        assembled.replace(DECK)

    print(f"Inserted expanded-likelihood multinomial slide after slide 14: {DECK}")
    print(f"Slide count: {original_count} -> {original_count + 1}")
    print("Verified: every pre-existing slide and notes payload is byte-identical")


if __name__ == "__main__":
    main()
