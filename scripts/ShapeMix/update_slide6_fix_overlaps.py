#!/usr/bin/env python3
"""Fix overlapping text/boxes on slide 6 (why-deconvolution).

1. The "one cell ≈ 10 µm" callout label overlapped the square spot's right
   border: the label moves right of the square, and its pointer line is
   redrawn from the new position.
2. The "ONE BLENDED PROFILE PER SPOT BARCODE" pill overlapped the card
   heading "What one spot reports": the pill moves onto its own line below
   the heading, and the cell-icon row shifts down to keep clear of it.

Only the why-deconvolution slide part (and its notes, which travel with it)
is rewritten; every other package part is copied byte-for-byte. Parts are
matched by content because python-pptx renumbers stored partnames on load.
"""

from __future__ import annotations

import re
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

KICKER = "WHY DECONVOLUTION"
PILL_TEXT = "ONE BLENDED PROFILE PER SPOT BARCODE"
SLATE = "52627A"
EMU_PER_INCH = 914400
ICON_ROW_SHIFT = 0.28


def inches(emu: int) -> float:
    return emu / EMU_PER_INCH


def slide_texts(slide) -> list[str]:
    return [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]


def fix_cell_size_callout(slide) -> None:
    label = None
    old_line = None
    for shape in slide.shapes:
        if shape.has_text_frame and "one cell" in shape.text_frame.text:
            label = shape
        elif (
            shape.shape_type == MSO_SHAPE_TYPE.LINE
            and inches(shape.left) < 6.4
            and inches(shape.top) < 4.0
            and inches(shape.height) < 0.4
        ):
            old_line = shape
    if label is None or old_line is None:
        raise RuntimeError("Cell-size callout label or pointer line not found")

    # Clear of the square's right edge at x = 4.795, inside the card at 6.02.
    label.left = Inches(4.88)
    label.width = Inches(1.10)

    line_color = old_line.line.color.rgb
    line_width = old_line.line.width
    old_line._element.getparent().remove(old_line._element)
    new_line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(4.86), Inches(3.44), Inches(4.32), Inches(3.55)
    )
    new_line.line.color.rgb = line_color
    new_line.line.width = line_width or Pt(1.0)


def fix_blended_profile_pill(slide) -> None:
    pill = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text == PILL_TEXT:
            pill = shape
            break
    if pill is None:
        raise RuntimeError("Blended-profile pill not found")
    # Own line below the heading; right-aligned inside the card.
    pill.left = Inches(9.55)
    pill.top = Inches(2.44)
    pill.width = Inches(2.95)

    # Shift the icon row (cells, T/B letters, "+", and the 3T+1B text) down
    # so the pill's new line stays clear of it.
    shifted = 0
    for shape in slide.shapes:
        if shape is pill:
            continue
        x, y = inches(shape.left), inches(shape.top)
        if x > 6.5 and 2.55 < y < 2.85:
            shape.top = Inches(y + ICON_ROW_SHIFT)
            shifted += 1
    # 4 cell circles + 4 highlight dots + 4 letter boxes + "+" + example text
    if shifted != 14:
        raise RuntimeError(f"Expected to shift 14 icon-row shapes; shifted {shifted}")


def find_slide_part(zipfile: ZipFile, marker: str) -> str:
    matches = [
        name
        for name in zipfile.namelist()
        if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
        and marker.encode("utf-8") in zipfile.read(name)
    ]
    if len(matches) != 1:
        raise RuntimeError(f"Expected one slide part containing {marker!r}; found {matches}")
    return matches[0]


def main() -> None:
    presentation = Presentation(DECK)
    slide = presentation.slides[5]
    if KICKER not in " ".join(slide_texts(slide)):
        raise RuntimeError("Slide 6 is not the why-deconvolution slide")

    fix_cell_size_callout(slide)
    fix_blended_profile_pill(slide)

    with TemporaryDirectory(prefix="shapemix_fix_overlaps_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"
        presentation.save(generated)

        with ZipFile(DECK, "r") as original, ZipFile(
            generated, "r"
        ) as regenerated, ZipFile(assembled, "w", compression=ZIP_DEFLATED) as output:
            orig_slide = find_slide_part(original, KICKER)
            gen_slide = find_slide_part(regenerated, KICKER)
            payload_map = {orig_slide: gen_slide}
            for info in original.infolist():
                if info.filename in payload_map:
                    data = regenerated.read(payload_map[info.filename])
                else:
                    data = original.read(info.filename)
                output.writestr(info, data)

        check = Presentation(assembled)
        joined = " ".join(slide_texts(check.slides[5]))
        for expected in (KICKER, PILL_TEXT, "one cell"):
            if expected not in joined:
                raise RuntimeError(f"Missing on saved slide 6: {expected!r}")
        if "SPATIAL ATAC" not in " ".join(slide_texts(check.slides[4])):
            raise RuntimeError("Slide 5 changed unexpectedly")
        if len(check.slides) != len(presentation.slides):
            raise RuntimeError("Slide count changed unexpectedly")

        assembled.replace(DECK)

    print(f"Fixed overlapping shapes on slide 6: {DECK}")


if __name__ == "__main__":
    main()
