#!/usr/bin/env python3
"""Redraw slide 5 as a grid of square spots containing blue/green cells.

The revised figure drops the Tn5/fragment cards and shows only what the user
asked for: a grid of square barcoded spots, each covering a few cells. Cells
use exactly two colors — teal (green) for T cells and blue for B cells — to
match the deconvolution example on slides 6-7. One spot is highlighted with
3 T + 1 B cells: the exact mixture the next slide zooms into.

Only the spatial-ATAC slide part and its notes part are rewritten; every
other package part is copied byte-for-byte from the original file. Parts are
matched by content because python-pptx renumbers stored partnames on load.
"""

from __future__ import annotations

import re
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

KICKER = "SPATIAL ATAC"
TITLE = "Spatial ATAC maps open chromatin across a tissue"
FOOTER_SOURCE = "Source: Deng et al., Nature (2022), spatial-ATAC-seq. Diagram is conceptual."

NAVY = "13233B"
INK = "213047"
SLATE = "52627A"
TEAL = "18A999"
BLUE = "3B82F6"
GOLD = "F4B942"
WHITE = "FFFFFF"
PALE = "F4F7FA"
MID = "D7E0E8"

FONT = "Aptos"

NOTES = (
    "Spatial ATAC divides the tissue into a grid of barcoded spots. Each square of the grid "
    "is one spot with its own spot barcode, and the colored dots are the cells that spot "
    "covers — only two cell types are drawn, teal-green T cells and blue B cells, to match "
    "the deconvolution example on the following slides. Every spot reports one ATAC "
    "peak-count profile summed over all of its cells. The spot outlined in dark navy holds "
    "three T cells and one B cell: the next slide zooms into exactly that situation and "
    "shows why the blended profile must be unmixed."
)

# Per-square cell layouts: (col, row) -> list of (dx, dy, color) inside the square.
P2 = ((0.16, 0.22), (0.44, 0.40))
P3 = ((0.14, 0.14), (0.46, 0.20), (0.26, 0.46))
P4 = ((0.12, 0.12), (0.46, 0.16), (0.16, 0.46), (0.46, 0.44))

GRID_CELLS = [
    # row 0
    [(P3, (BLUE, BLUE, TEAL))],
    [(P2, (BLUE, BLUE))],
    [(P4, (TEAL, TEAL, TEAL, BLUE))],  # highlighted spot: 3 T + 1 B
    [(P3, (TEAL, TEAL, BLUE))],
    [(P2, (TEAL, TEAL))],
    # row 1
    [(P3, (BLUE, BLUE, BLUE))],
    [(P2, (BLUE, TEAL))],
    [(P3, (TEAL, BLUE, TEAL))],
    [(P3, (TEAL, TEAL, TEAL))],
    [(P2, (TEAL, BLUE))],
    # row 2
    [(P2, (BLUE, TEAL))],
    [(P3, (BLUE, BLUE, TEAL))],
    [(P2, (TEAL, BLUE))],
    [(P3, (TEAL, TEAL, BLUE))],
    [(P4, (TEAL, TEAL, TEAL, TEAL))],
    # row 3
    [(P2, (BLUE, BLUE))],
    [(P4, (BLUE, BLUE, TEAL, TEAL))],
    [(P2, (BLUE, TEAL))],
    [(P3, (TEAL, BLUE, BLUE))],
    [(P3, (TEAL, TEAL, BLUE))],
]
HIGHLIGHT_INDEX = 2  # row 0, col 2

GRID_COLS, GRID_ROWS = 5, 4
SQUARE = 0.82
GAP = 0.08
CELL_D = 0.22


def C(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#").upper())


def add_shape(slide, kind, x, y, w, h, fill, line=None, line_width=1.0):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(fill)
    shape.line.color.rgb = C(line if line else fill)
    shape.line.width = Pt(line_width)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=20,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = C(color)
    return box


def add_circle(slide, x, y, d, fill, line=WHITE, line_width=1.5):
    return add_shape(slide, MSO_SHAPE.OVAL, x, y, d, d, fill, line=line, line_width=line_width)


def add_pill(slide, text, x, y, w, h, fill, color=WHITE, size=13):
    shp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line=fill, line_width=1)
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = C(color)
    return shp


def add_card(slide, x, y, w, h, fill=WHITE, line=MID, line_width=1.0):
    return add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line=line, line_width=line_width
    )


def slide_texts(slide) -> list[str]:
    return [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]


def clear_slide_body(slide) -> None:
    """Remove everything except the kicker, title, and footer."""
    keep_texts = {KICKER, TITLE, FOOTER_SOURCE}
    kept = set()
    for shape in list(slide.shapes):
        if shape.has_text_frame and shape.text_frame.text in keep_texts:
            kept.add(shape.text_frame.text)
            continue
        if shape.shape_type == 9 and abs(shape.top / 914400 - 7.13) < 0.03:
            continue  # footer rule line
        shape._element.getparent().remove(shape._element)
    if kept != keep_texts:
        raise RuntimeError(f"Missing expected header/footer shapes: {keep_texts - kept}")


def draw_figure(slide) -> None:
    add_text(
        slide,
        "The tissue is divided into a grid of barcoded spots — each spot collects ATAC signal from the cells it covers.",
        0.72,
        1.14,
        11.84,
        0.44,
        size=18.5,
        color=SLATE,
        align=PP_ALIGN.CENTER,
    )

    # Left: the spot grid.
    add_card(slide, 0.86, 1.80, 6.50, 4.62, fill=PALE, line=MID)
    add_text(slide, "Tissue section on a barcoded grid", 1.12, 2.02, 5.9, 0.34, size=18, color=NAVY, bold=True)

    grid_w = GRID_COLS * SQUARE + (GRID_COLS - 1) * GAP
    gx0 = 0.86 + (6.50 - grid_w) / 2
    gy0 = 2.56
    for index, cell_spec in enumerate(GRID_CELLS):
        col = index % GRID_COLS
        row = index // GRID_COLS
        sq_x = gx0 + col * (SQUARE + GAP)
        sq_y = gy0 + row * (SQUARE + GAP)
        highlighted = index == HIGHLIGHT_INDEX
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            sq_x,
            sq_y,
            SQUARE,
            SQUARE,
            WHITE,
            line=NAVY if highlighted else SLATE,
            line_width=2.5 if highlighted else 1.0,
        )
        (offsets, colors) = cell_spec[0]
        for (dx, dy), color in zip(offsets, colors):
            add_circle(slide, sq_x + dx, sq_y + dy, CELL_D, color, line=WHITE, line_width=1)

    caption_y = gy0 + GRID_ROWS * SQUARE + (GRID_ROWS - 1) * GAP + 0.08
    add_text(
        slide,
        "each square = one spot with its own spot barcode",
        1.08,
        caption_y,
        3.85,
        0.26,
        size=10.5,
        color=SLATE,
        margin=0.01,
    )
    add_circle(slide, 5.06, caption_y + 0.03, 0.18, TEAL, line=WHITE, line_width=1)
    add_text(slide, "T cell", 5.28, caption_y, 0.62, 0.26, size=10.5, color=INK, margin=0.01)
    add_circle(slide, 5.98, caption_y + 0.03, 0.18, BLUE, line=WHITE, line_width=1)
    add_text(slide, "B cell", 6.20, caption_y, 0.62, 0.26, size=10.5, color=INK, margin=0.01)

    # Right: how to read the grid.
    add_card(slide, 7.72, 1.80, 4.90, 4.62, fill=WHITE, line=MID)
    add_text(slide, "How to read the grid", 8.00, 2.02, 4.3, 0.34, size=18, color=NAVY, bold=True)
    rows = [
        ("1", BLUE, "One spot barcode per square", "every square of the grid carries its own DNA barcode"),
        ("2", TEAL, "Cells sit inside spots", "a spot usually covers several cells — here two types, blue and green"),
        ("3", GOLD, "One profile per spot", "each spot reports a single ATAC peak-count profile for all its cells together"),
    ]
    for i, (num, color, heading, body) in enumerate(rows):
        y = 2.58 + i * 1.22
        add_circle(slide, 8.02, y, 0.50, color, line=color, line_width=1)
        add_text(
            slide,
            num,
            8.02,
            y + 0.01,
            0.50,
            0.40,
            size=17,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0,
        )
        add_text(slide, heading, 8.70, y - 0.04, 3.75, 0.32, size=15.5, color=NAVY, bold=True)
        add_text(slide, body, 8.70, y + 0.30, 3.75, 0.62, size=12.5, color=SLATE)

    add_pill(
        slide,
        "RESULT: ATAC PEAK COUNTS AT EVERY SPOT IN THE TISSUE",
        3.57,
        6.56,
        6.20,
        0.40,
        TEAL,
        size=10.5,
    )


def set_notes(slide) -> None:
    tf = slide.notes_slide.notes_text_frame
    if tf is None:
        raise RuntimeError("Notes body placeholder missing on the spatial slide")
    tf.text = NOTES


def find_slide_part(zipfile: ZipFile, marker: str) -> str:
    matches = [
        name
        for name in zipfile.namelist()
        if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
        and marker.encode("utf-8") in zipfile.read(name)
    ]
    if len(matches) != 1:
        raise RuntimeError(f"Expected one slide part containing {marker!r}; found {matches}")
    return matches[0]


def notes_part_for(zipfile: ZipFile, slide_part: str) -> str:
    rels_name = f"ppt/slides/_rels/{slide_part.rsplit('/', 1)[-1]}.rels"
    rels = zipfile.read(rels_name).decode("utf-8")
    targets = re.findall(r'Target="\.\./(notesSlides/notesSlide\d+\.xml)"', rels)
    if len(targets) != 1:
        raise RuntimeError(f"Expected one notes target in {rels_name}; found {targets}")
    return f"ppt/{targets[0]}"


def main() -> None:
    presentation = Presentation(DECK)
    slide = presentation.slides[4]
    if KICKER not in " ".join(slide_texts(slide)):
        raise RuntimeError("Slide 5 is not the spatial-ATAC slide")

    clear_slide_body(slide)
    draw_figure(slide)
    set_notes(slide)

    forbidden = re.compile(r"tn5|fragment", re.IGNORECASE)
    for text in slide_texts(slide):
        if text != FOOTER_SOURCE and forbidden.search(text):
            raise RuntimeError(f"Forbidden wording remains on slide 5: {text!r}")

    with TemporaryDirectory(prefix="shapemix_spot_grid_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"
        presentation.save(generated)

        with ZipFile(DECK, "r") as original, ZipFile(
            generated, "r"
        ) as regenerated, ZipFile(assembled, "w", compression=ZIP_DEFLATED) as output:
            orig_slide = find_slide_part(original, KICKER)
            gen_slide = find_slide_part(regenerated, KICKER)
            payload_map = {
                orig_slide: gen_slide,
                notes_part_for(original, orig_slide): notes_part_for(regenerated, gen_slide),
            }
            for info in original.infolist():
                if info.filename in payload_map:
                    data = regenerated.read(payload_map[info.filename])
                else:
                    data = original.read(info.filename)
                output.writestr(info, data)

        check = Presentation(assembled)
        joined = " ".join(slide_texts(check.slides[4]))
        for expected in (KICKER, TITLE, "Tissue section on a barcoded grid", "How to read the grid"):
            if expected not in joined:
                raise RuntimeError(f"Missing on saved slide 5: {expected!r}")
        if "WHY DECONVOLUTION" not in " ".join(slide_texts(check.slides[5])):
            raise RuntimeError("Slide 6 changed unexpectedly")
        if len(check.slides) != len(presentation.slides):
            raise RuntimeError("Slide count changed unexpectedly")

        assembled.replace(DECK)

    print(f"Redrew slide 5 as a square spot grid: {DECK}")


if __name__ == "__main__":
    main()
