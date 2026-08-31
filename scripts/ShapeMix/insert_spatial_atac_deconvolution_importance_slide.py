#!/usr/bin/env python3
"""Insert a slide on why deconvolution matters immediately after slide 6.

The final package is assembled from the original ZIP so every pre-existing
slide and notes payload remains byte-for-byte unchanged.
"""

from __future__ import annotations

import re
from hashlib import sha256
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from insert_count_only_bayesian_slides import (
    BLUE,
    BLUE_PALE,
    C,
    CREAM,
    FONT,
    FONT_DISPLAY,
    GOLD,
    GOLD_PALE,
    INK,
    MID,
    NAVY,
    PURPLE,
    PURPLE_PALE,
    SLATE,
    TEAL,
    TEAL_PALE,
    WHITE,
    add_card,
    add_pill,
    add_shape,
    add_text,
    set_bg,
)

REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

INSERT_AFTER_INDEX = 5  # zero-based index of slide 6
KICKER = "WHY DECONVOLUTION"
TITLE = "Why deconvolution matters for spatial ATAC-seq"
NEXT_TITLE = "Bayesian model to find the number of cells for each cell type"
GREEN = "43AA6A"
GREEN_PALE = "E5F5E9"


def add_centered_text(slide, text, x, y, w, h, *, size, color, bold=False):
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_spot_panel(slide):
    x, y, w, h = 0.70, 1.72, 3.50, 2.75
    add_card(slide, x, y, w, h, fill=WHITE, line=MID)
    add_pill(slide, "ONE SPATIAL SPOT", x + 0.82, y + 0.18, 1.86, 0.31, BLUE)

    # Ten stylized cells: six T-like, three B-like, one monocyte-like.
    colors = [TEAL] * 6 + [PURPLE] * 3 + [GOLD]
    positions = [
        (1.11, 2.45), (1.63, 2.38), (2.15, 2.48), (2.68, 2.39), (3.18, 2.50),
        (1.35, 2.98), (1.90, 3.02), (2.45, 2.96), (2.99, 3.04), (2.18, 3.48),
    ]
    for (cx, cy), color in zip(positions, colors):
        cell = add_shape(slide, MSO_SHAPE.OVAL, cx, cy, 0.38, 0.38, color, line=WHITE)
        cell.line.width = Pt(1.1)

    add_centered_text(
        slide,
        "many cells and cell types\nshare the same spot",
        x + 0.25,
        y + 2.03,
        w - 0.50,
        0.52,
        size=12,
        color=SLATE,
        bold=True,
    )


def add_mixed_signal_panel(slide):
    x, y, w, h = 4.92, 1.72, 3.50, 2.75
    add_card(slide, x, y, w, h, fill=WHITE, line=MID)
    add_pill(slide, "MIXED ATAC SIGNAL", x + 0.76, y + 0.18, 1.98, 0.31, NAVY)

    # One observed peak-count profile. The bars are deliberately one color:
    # the assay records the mixture, not each cell type's contribution.
    baseline_y = 3.57
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.46, baseline_y, 2.58, 0.015, MID)
    heights = [0.55, 1.18, 0.75, 1.42, 0.88]
    for i, height in enumerate(heights):
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            x + 0.58 + i * 0.49,
            baseline_y - height,
            0.25,
            height,
            NAVY,
        )

    add_centered_text(
        slide,
        "all cut sites are counted together\n→ one blended profile",
        x + 0.25,
        y + 2.03,
        w - 0.50,
        0.52,
        size=12,
        color=SLATE,
        bold=True,
    )


def add_deconvolution_panel(slide):
    x, y, w, h = 9.12, 1.72, 3.50, 2.75
    add_card(slide, x, y, w, h, fill=TEAL_PALE, line=TEAL)
    add_pill(slide, "DECONVOLUTION", x + 0.88, y + 0.18, 1.74, 0.31, TEAL)

    add_centered_text(
        slide,
        "estimated cell-type composition",
        x + 0.25,
        y + 0.70,
        w - 0.50,
        0.30,
        size=12,
        color=NAVY,
        bold=True,
    )

    # A 60/30/10 composition bar.
    bar_x, bar_y, bar_w, bar_h = x + 0.42, y + 1.18, 2.66, 0.52
    segments = [
        (0.60, TEAL, "60% T"),
        (0.30, PURPLE, "30% B"),
        (0.10, GOLD, ""),
    ]
    cursor = bar_x
    for fraction, color, label in segments:
        seg_w = bar_w * fraction
        add_shape(slide, MSO_SHAPE.RECTANGLE, cursor, bar_y, seg_w, bar_h, color)
        if label:
            add_centered_text(
                slide,
                label,
                cursor,
                bar_y,
                seg_w,
                bar_h,
                size=10,
                color=WHITE if color != GOLD else NAVY,
                bold=True,
            )
        cursor += seg_w
    add_centered_text(
        slide,
        "10% mono",
        x + 0.42,
        y + 1.78,
        2.66,
        0.24,
        size=10.5,
        color=SLATE,
    )
    add_centered_text(
        slide,
        "estimate who contributed\nto the mixed signal",
        x + 0.25,
        y + 2.08,
        w - 0.50,
        0.48,
        size=12,
        color=NAVY,
        bold=True,
    )


def add_impact_card(slide, x, title, body, fill, line):
    add_card(slide, x, 4.83, 3.78, 1.42, fill=fill, line=line)
    add_text(slide, title, x + 0.18, 4.98, 3.42, 0.28, size=11.5, color=line, bold=True)
    add_text(slide, body, x + 0.18, 5.32, 3.42, 0.74, size=11, color=INK)


def build(slide) -> None:
    set_bg(slide, CREAM)
    add_text(slide, KICKER, 0.68, 0.28, 8.8, 0.28, size=10.5, color=TEAL, bold=True)
    add_text(
        slide,
        TITLE,
        0.65,
        0.55,
        11.8,
        0.62,
        size=27,
        color=NAVY,
        bold=True,
        font=FONT_DISPLAY,
    )
    add_text(
        slide,
        "A spatial spot usually contains several cells, so its accessibility profile is a mixture.",
        0.71,
        1.17,
        11.9,
        0.38,
        size=14,
        color=SLATE,
    )

    add_spot_panel(slide)
    add_centered_text(slide, "→", 4.25, 2.72, 0.58, 0.55, size=28, color=TEAL, bold=True)
    add_mixed_signal_panel(slide)
    add_centered_text(slide, "→", 8.47, 2.72, 0.58, 0.55, size=28, color=TEAL, bold=True)
    add_deconvolution_panel(slide)

    add_impact_card(
        slide,
        0.70,
        "MAP CELL-TYPE COMPOSITION",
        "Reveal where different cell populations are enriched across the tissue.",
        BLUE_PALE,
        BLUE,
    )
    add_impact_card(
        slide,
        4.78,
        "COMPARE REGIONS FAIRLY",
        "A peak can change because the cell mixture changed—not only because chromatin changed.",
        GOLD_PALE,
        "B8860B",
    )
    add_impact_card(
        slide,
        8.86,
        "CONNECT SIGNAL TO BIOLOGY",
        "Identify which cell populations may drive spatial patterns in tissue or disease.",
        GREEN_PALE,
        GREEN,
    )

    add_centered_text(
        slide,
        "Goal: estimate each spot's cell-type composition—not assign every cut site to one cell.",
        0.70,
        6.55,
        11.93,
        0.30,
        size=11,
        color=SLATE,
    )
    hairline = add_shape(slide, MSO_SHAPE.RECTANGLE, 0.68, 7.13, 11.97, 0.011, MID)
    hairline.line.width = Pt(0.25)


def slide_texts(slide) -> list[str]:
    return [
        shape.text_frame.text.strip()
        for shape in slide.shapes
        if shape.has_text_frame and shape.text_frame.text.strip()
    ]


def part_hashes(path: Path, prefixes: tuple[str, ...]) -> dict[str, str]:
    with ZipFile(path) as archive:
        return {
            name: sha256(archive.read(name)).hexdigest()
            for name in archive.namelist()
            if name.startswith(prefixes)
        }


NEW_PART_PATTERN = re.compile(
    r"^ppt/slides/(?:_rels/)?slide\d+\.xml(?:\.rels)?$"
)
REPLACED_PARTS = {
    "[Content_Types].xml",
    "ppt/presentation.xml",
    "ppt/_rels/presentation.xml.rels",
    "docProps/app.xml",
}


def assemble(original_path: Path, generated_path: Path, output_path: Path) -> None:
    with ZipFile(original_path, "r") as original, ZipFile(
        generated_path, "r"
    ) as generated, ZipFile(output_path, "w", compression=ZIP_DEFLATED) as output:
        original_names = set(original.namelist())
        generated_names = set(generated.namelist())
        new_parts = generated_names - original_names
        bad = [name for name in new_parts if not NEW_PART_PATTERN.match(name)]
        if bad:
            raise RuntimeError(f"Unexpected new package parts: {sorted(bad)}")
        missing = REPLACED_PARTS - generated_names
        if missing:
            raise RuntimeError(f"Generated package is missing parts: {sorted(missing)}")

        for info in original.infolist():
            data = (
                generated.read(info.filename)
                if info.filename in REPLACED_PARTS
                else original.read(info.filename)
            )
            output.writestr(info, data)
        for name in sorted(new_parts):
            output.writestr(generated.getinfo(name), generated.read(name))


def main() -> None:
    original = Presentation(DECK)
    original_count = len(original.slides)
    for slide in original.slides:
        if any(TITLE in text for text in slide_texts(slide)):
            raise RuntimeError("The deconvolution-importance slide already exists")
    if not any(NEXT_TITLE in text for text in slide_texts(original.slides[6])):
        raise RuntimeError("Slide 7 is not the expected Bayesian-model slide")

    preserved_before = part_hashes(
        DECK,
        ("ppt/slides/", "ppt/notesSlides/"),
    )

    with TemporaryDirectory(prefix="shapemix_deconv_importance_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"

        presentation = Presentation(DECK)
        new_slide = presentation.slides.add_slide(presentation.slides[5].slide_layout)
        build(new_slide)
        slide_ids = presentation.slides._sldIdLst
        new_id = slide_ids[-1]
        slide_ids.remove(new_id)
        slide_ids.insert(INSERT_AFTER_INDEX + 1, new_id)
        presentation.save(generated)

        assemble(DECK, generated, assembled)

        check = Presentation(assembled)
        if len(check.slides) != original_count + 1:
            raise RuntimeError(
                f"Expected {original_count + 1} slides; found {len(check.slides)}"
            )
        new_text = " ".join(slide_texts(check.slides[6]))
        for needle in (
            TITLE,
            "MIXED ATAC SIGNAL",
            "DECONVOLUTION",
            "COMPARE REGIONS FAIRLY",
        ):
            if needle not in new_text:
                raise RuntimeError(f"New slide 7 is missing {needle!r}")
        if not any(NEXT_TITLE in text for text in slide_texts(check.slides[7])):
            raise RuntimeError("The original slide 7 did not move intact to position 8")

        preserved_after = part_hashes(
            assembled,
            ("ppt/slides/", "ppt/notesSlides/"),
        )
        for name, digest in preserved_before.items():
            if preserved_after.get(name) != digest:
                raise RuntimeError(f"Existing slide or notes payload changed: {name}")

        assembled.replace(DECK)

    print(f"Inserted deconvolution-importance slide after slide 6: {DECK}")
    print(f"Slide count: {original_count} -> {original_count + 1}")
    print("Verified: every pre-existing slide and notes payload is byte-identical")


if __name__ == "__main__":
    main()
