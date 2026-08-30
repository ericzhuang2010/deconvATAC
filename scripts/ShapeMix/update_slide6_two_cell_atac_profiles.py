#!/usr/bin/env python3
"""Redesign only slide 6 with two reference cell types and ATAC profiles."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

NAVY = "13233B"
INK = "213047"
SLATE = "52627A"
TEAL = "18A999"
TEAL_PALE = "DDF4F1"
BLUE = "3B82F6"
BLUE_PALE = "E7F0FE"
PURPLE = "8B5CF6"
PURPLE_PALE = "EEE8FF"
GOLD = "F4B942"
WHITE = "FFFFFF"
MID = "D7E0E8"
PALE = "F4F7FA"

FONT = "Aptos"


def color(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_shape(slide, kind, x, y, width, height, fill, line=None, line_width=1.0):
    shape = slide.shapes.add_shape(
        kind, Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(line or fill)
    shape.line.width = Pt(line_width)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    width,
    height,
    *,
    size=16,
    text_color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.03,
):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(text_color)
    return box


def add_pill(slide, text, x, y, width, height, fill, *, text_color=WHITE, size=10):
    shape = add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        width,
        height,
        fill,
        line=fill,
    )
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = tf.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color(text_color)
    return shape


def add_line(slide, x1, y1, x2, y2, *, line_color=MID, width=1.0):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color(line_color)
    line.line.width = Pt(width)
    return line


def add_profile_chart(
    slide,
    values,
    x,
    baseline,
    *,
    bar_width,
    gap,
    unit_height,
    bar_color,
    value_size=8.5,
    peak_size=7.5,
):
    for index, value in enumerate(values):
        bar_x = x + index * (bar_width + gap)
        bar_height = value * unit_height
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            bar_x,
            baseline - bar_height,
            bar_width,
            bar_height,
            bar_color,
            line=bar_color,
            line_width=0.5,
        )
        add_text(
            slide,
            str(value),
            bar_x - 0.08,
            baseline - bar_height - 0.18,
            bar_width + 0.16,
            0.17,
            size=value_size,
            text_color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
        add_text(
            slide,
            f"P{index + 1}",
            bar_x - 0.08,
            baseline + 0.04,
            bar_width + 0.16,
            0.17,
            size=peak_size,
            text_color=SLATE,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
    chart_width = 4 * bar_width + 3 * gap
    add_line(slide, x - 0.06, baseline, x + chart_width + 0.06, baseline, line_color=NAVY, width=0.8)


def add_reference_card(slide, *, y, label, values, accent, pale):
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        0.76,
        y,
        3.18,
        1.52,
        WHITE,
        line=accent,
        line_width=1.2,
    )
    add_shape(slide, MSO_SHAPE.OVAL, 1.00, y + 0.47, 0.52, 0.52, accent, line=WHITE)
    add_shape(slide, MSO_SHAPE.OVAL, 1.18, y + 0.64, 0.17, 0.17, WHITE, line=WHITE, line_width=0.4)
    add_text(
        slide,
        label,
        1.65,
        y + 0.20,
        1.18,
        0.28,
        size=15,
        text_color=accent,
        bold=True,
    )
    add_pill(slide, "REFERENCE", 1.65, y + 0.55, 0.94, 0.23, pale, text_color=accent, size=7.6)
    add_profile_chart(
        slide,
        values,
        2.64,
        y + 1.16,
        bar_width=0.20,
        gap=0.14,
        unit_height=0.055,
        bar_color=GOLD,
    )


def add_stacked_recipe(slide):
    x, y, width, height = 9.34, 3.36, 2.88, 0.58
    t_width = width * 0.75
    b_width = width - t_width
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, t_width, height, TEAL, line=TEAL)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + t_width, y, b_width, height, BLUE, line=BLUE)
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height)
    )
    border.fill.background()
    border.line.color.rgb = color(NAVY)
    border.line.width = Pt(1.0)
    add_text(
        slide,
        "75%",
        x,
        y,
        t_width,
        height,
        size=13,
        text_color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )
    add_text(
        slide,
        "25%",
        x + t_width,
        y,
        b_width,
        height,
        size=11,
        text_color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )


def find_target_slide(presentation):
    target_title = "Deconvolution works like identifying ingredients in a smoothie"
    matches = []
    for slide in presentation.slides:
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text_frame")]
        if target_title in texts:
            matches.append(slide)
    if len(matches) != 1:
        raise RuntimeError(f"Expected exactly one target slide; found {len(matches)}")
    return matches[0]


def main() -> None:
    presentation = Presentation(DECK)
    slide = find_target_slide(presentation)

    # Preserve the four header objects exactly and replace only the diagram.
    header_ids = {2, 3, 4, 5}
    if {shape.shape_id for shape in slide.shapes if shape.shape_id in header_ids} != header_ids:
        raise RuntimeError("Slide 6 header signatures changed; refusing to patch")
    for shape in list(slide.shapes):
        if shape.shape_id not in header_ids:
            shape._element.getparent().remove(shape._element)

    # Left: two reference cell types, explicitly ordered T above B.
    add_reference_card(
        slide,
        y=2.02,
        label="T cell",
        values=(10, 2, 8, 4),
        accent=TEAL,
        pale=TEAL_PALE,
    )
    add_reference_card(
        slide,
        y=4.02,
        label="B cell",
        values=(2, 10, 4, 8),
        accent=BLUE,
        pale=BLUE_PALE,
    )
    add_text(
        slide,
        "reference ATAC count profiles",
        0.94,
        5.70,
        2.82,
        0.24,
        size=10.5,
        text_color=SLATE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # The mixed spot is an observed ATAC profile—not colored cell dots.
    add_text(slide, "+", 3.95, 3.39, 0.36, 0.52, size=27, text_color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        4.38,
        2.02,
        3.22,
        3.52,
        WHITE,
        line=PURPLE,
        line_width=1.2,
    )
    add_text(
        slide,
        "UNKNOWN MIXED SPOT",
        4.70,
        2.28,
        2.58,
        0.28,
        size=14,
        text_color=PURPLE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "observed ATAC count profile",
        4.71,
        2.67,
        2.56,
        0.24,
        size=10.5,
        text_color=SLATE,
        align=PP_ALIGN.CENTER,
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 4.70, 3.16, 2.58, 2.04, PURPLE_PALE, line=MID)
    add_profile_chart(
        slide,
        (8, 4, 7, 5),
        4.96,
        4.82,
        bar_width=0.34,
        gap=0.22,
        unit_height=0.12,
        bar_color=GOLD,
        value_size=10,
        peak_size=8.5,
    )
    add_text(
        slide,
        "same four peaks • same count scale",
        4.79,
        5.25,
        2.40,
        0.20,
        size=8.7,
        text_color=SLATE,
        align=PP_ALIGN.CENTER,
    )

    # Model arrow.
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, 7.78, 3.26, 0.94, 0.66, TEAL, line=TEAL)
    add_pill(slide, "BAYESIAN\nMODEL", 7.78, 4.02, 0.94, 0.48, NAVY, size=8.2)

    # Right: two-cell recipe, T first and B second.
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        8.93,
        2.02,
        3.70,
        3.52,
        WHITE,
        line=NAVY,
        line_width=1.2,
    )
    add_text(
        slide,
        "Estimated recipe",
        9.23,
        2.31,
        3.10,
        0.34,
        size=17,
        text_color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "75% T  +  25% B",
        9.23,
        2.73,
        3.10,
        0.34,
        size=19,
        text_color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_stacked_recipe(slide)

    for row_y, label, percentage, accent in (
        (4.22, "T cell", "75%", TEAL),
        (4.72, "B cell", "25%", BLUE),
    ):
        add_shape(slide, MSO_SHAPE.OVAL, 9.36, row_y + 0.03, 0.20, 0.20, accent, line=accent)
        add_text(slide, label, 9.67, row_y - 0.03, 1.42, 0.28, size=12.5, text_color=INK)
        add_text(
            slide,
            percentage,
            11.42,
            row_y - 0.03,
            0.60,
            0.28,
            size=12.5,
            text_color=NAVY,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
    add_pill(
        slide,
        "mixed = 0.75 × T + 0.25 × B",
        9.26,
        5.10,
        3.04,
        0.28,
        NAVY,
        size=9.3,
    )

    add_pill(
        slide,
        "THE BAYESIAN MODEL FINDS THE T/B WEIGHTS THAT BEST RECONSTRUCT THE OBSERVED ATAC COUNTS",
        1.25,
        6.18,
        10.84,
        0.38,
        TEAL,
        size=9.5,
    )

    try:
        slide.notes_slide.notes_text_frame.text = (
            "This slide uses only two cell types. Read the T-cell reference above the B-cell "
            "reference, then compare both with the unknown spot's observed ATAC peak-count "
            "profile. The teaching counts are T = [10, 2, 8, 4], B = [2, 10, 4, 8], and "
            "mixed = [8, 4, 7, 5]. At every peak, the mixed profile equals 0.75 times the T "
            "reference plus 0.25 times the B reference. The three profiles all total 24 cut "
            "sites, so the example isolates composition rather than sequencing depth. Real "
            "data fluctuate, so a fitted reconstruction is usually approximate."
        )
    except Exception:
        pass

    temporary = DECK.with_name(f".{DECK.stem}.slide6-two-cell.tmp.pptx")
    presentation.save(temporary)
    check = Presentation(temporary)
    if len(check.slides) != len(presentation.slides):
        raise RuntimeError("Slide count changed while saving slide 6")
    temporary.replace(DECK)
    print(f"Updated only slide 6: {DECK}")


if __name__ == "__main__":
    main()
