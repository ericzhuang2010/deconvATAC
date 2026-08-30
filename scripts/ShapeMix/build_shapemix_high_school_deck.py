#!/usr/bin/env python3
"""Build an approachable, results-free ShapeMix-ATAC PowerPoint deck.

The deck uses only editable PowerPoint shapes and text.  Run with python-pptx
available on PYTHONPATH, for example:

    PYTHONPATH=/tmp/shapemix_pptx_deps .venv/bin/python \
        scripts/ShapeMix/build_shapemix_high_school_deck.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
OUT = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

W = 13.333
H = 7.5


def C(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor.from_string(value.upper())


NAVY = "13233B"
NAVY_2 = "1C3152"
INK = "213047"
SLATE = "52627A"
TEAL = "18A999"
TEAL_DARK = "0C7F75"
TEAL_PALE = "DDF4F1"
BLUE = "3B82F6"
BLUE_PALE = "E7F0FE"
PURPLE = "8B5CF6"
PURPLE_PALE = "EEE8FF"
CORAL = "F26B5B"
CORAL_PALE = "FDE9E6"
GOLD = "F4B942"
GOLD_PALE = "FFF3D6"
GREEN = "43AA6A"
GREEN_PALE = "E5F5E9"
CREAM = "F7F4EC"
WHITE = "FFFFFF"
PALE = "F4F7FA"
MID = "D7E0E8"
LIGHT_TEXT = "D7E2EE"

FONT = "Aptos"
FONT_DISPLAY = "Aptos Display"


def set_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C(color)


def add_shape(
    slide,
    kind,
    x,
    y,
    w,
    h,
    fill,
    line=None,
    line_width=1.0,
    radius=True,
):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(fill)
    shape.line.color.rgb = C(line if line else fill)
    shape.line.width = Pt(line_width)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=20,
    color=INK,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.04,
    italic=False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = C(color)
    return box


def add_rich_text(slide, runs, x, y, w, h, size=20, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    for content, kwargs in runs:
        r = p.add_run()
        r.text = content
        r.font.name = kwargs.get("font", FONT)
        r.font.size = Pt(kwargs.get("size", size))
        r.font.bold = kwargs.get("bold", False)
        r.font.italic = kwargs.get("italic", False)
        r.font.color.rgb = C(kwargs.get("color", color))
    return box


def add_bullets(slide, items, x, y, w, h, size=18, color=INK, bullet_color=TEAL):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        p.line_spacing = 1.05
        p.alignment = PP_ALIGN.LEFT
        rb = p.add_run()
        rb.text = "•  "
        rb.font.name = FONT
        rb.font.size = Pt(size)
        rb.font.bold = True
        rb.font.color.rgb = C(bullet_color)
        rt = p.add_run()
        rt.text = item
        rt.font.name = FONT
        rt.font.size = Pt(size)
        rt.font.color.rgb = C(color)
    return box


def add_line(slide, x1, y1, x2, y2, color=MID, width=1.5, dash=None):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = C(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = dash
    return line


def add_circle(slide, x, y, d, fill, line=WHITE, line_width=1.5):
    return add_shape(
        slide, MSO_SHAPE.OVAL, x, y, d, d, fill, line=line, line_width=line_width
    )


def add_pill(slide, text, x, y, w, h, fill, color=WHITE, size=13, line=None):
    shp = add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill,
        line=line if line else fill,
        line_width=1,
    )
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = C(color)
    return shp


def add_card(slide, x, y, w, h, fill=WHITE, line=MID, shadow=False):
    return add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill,
        line=line,
        line_width=1.0,
    )


def add_title(slide, title, kicker=None, dark=False, slide_num=None):
    if kicker:
        add_text(
            slide,
            kicker.upper(),
            0.68,
            0.28,
            8.8,
            0.28,
            size=10.5,
            color=TEAL if not dark else GOLD,
            bold=True,
        )
    add_text(
        slide,
        title,
        0.65,
        0.55 if kicker else 0.36,
        11.8,
        0.62,
        size=27,
        color=WHITE if dark else NAVY,
        bold=True,
        font=FONT_DISPLAY,
    )
    if slide_num is not None:
        add_text(
            slide,
            f"{slide_num:02d}",
            12.18,
            0.34,
            0.45,
            0.26,
            size=10,
            color=LIGHT_TEXT if dark else SLATE,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )


def add_footer(slide, source=None, dark=False):
    add_line(slide, 0.68, 7.13, 12.65, 7.13, color=NAVY_2 if dark else MID, width=0.8)
    if source:
        add_text(
            slide,
            source,
            0.7,
            7.17,
            11.8,
            0.18,
            size=7.5,
            color=LIGHT_TEXT if dark else SLATE,
        )


def add_notes(slide, text):
    try:
        tf = slide.notes_slide.notes_text_frame
        tf.text = text
    except Exception:
        # Notes support differs slightly across python-pptx releases.  The deck
        # remains fully usable if notes cannot be written.
        pass


def add_stacked_bar(slide, x, y, w, h, values, colors, labels=None, outline=NAVY):
    total = sum(values)
    cursor = x
    for i, (value, color) in enumerate(zip(values, colors)):
        seg_w = w * value / total
        rect = add_shape(slide, MSO_SHAPE.RECTANGLE, cursor, y, seg_w, h, color, line=color)
        if labels and seg_w > 0.45:
            add_text(
                slide,
                labels[i],
                cursor,
                y + 0.01,
                seg_w,
                h - 0.02,
                size=10,
                color=WHITE,
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
            )
        cursor += seg_w
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    border.fill.background()
    border.line.color.rgb = C(outline)
    border.line.width = Pt(1.0)
    return border


def add_fragment(slide, x, y, w, color, line_width=5, cut_color=INK):
    add_line(slide, x, y, x + w, y, color=color, width=line_width)
    add_line(slide, x, y - 0.06, x, y + 0.06, color=cut_color, width=1.2)
    add_line(slide, x + w, y - 0.06, x + w, y + 0.06, color=cut_color, width=1.2)


def add_nucleosome(slide, center_x, center_y, scale=1.0):
    """Draw a compact histone core with a visible DNA wrap."""
    dna_color = "654640"
    lobe_colors = ("6D43D6", PURPLE, "7650DE", PURPLE)
    lobe_w = 0.18 * scale
    lobe_h = 0.48 * scale
    for offset, color in zip((-0.22, -0.075, 0.075, 0.22), lobe_colors):
        lobe = add_shape(
            slide,
            MSO_SHAPE.OVAL,
            center_x + offset * scale - lobe_w / 2,
            center_y - lobe_h / 2,
            lobe_w,
            lobe_h,
            color,
            line="5634B5",
            line_width=0.6,
        )
        lobe.rotation = 18

    wrap = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(center_x - 0.38 * scale),
        Inches(center_y - 0.31 * scale),
        Inches(0.76 * scale),
        Inches(0.62 * scale),
    )
    wrap.fill.background()
    wrap.line.color.rgb = C(dna_color)
    wrap.line.width = Pt(1.45)
    wrap.rotation = 9
    add_line(
        slide,
        center_x - 0.34 * scale,
        center_y + 0.12 * scale,
        center_x + 0.34 * scale,
        center_y - 0.10 * scale,
        color=dna_color,
        width=1.25,
    )
    return wrap


def add_tn5_marker(slide, x, y):
    """Small two-lobed Tn5 marker placed on exposed linker DNA."""
    enzyme_blue = "68B5E8"
    add_circle(slide, x - 0.11, y - 0.14, 0.16, enzyme_blue, line=WHITE, line_width=0.6)
    add_circle(slide, x - 0.01, y - 0.10, 0.16, enzyme_blue, line=WHITE, line_width=0.6)
    add_line(slide, x + 0.02, y + 0.01, x + 0.02, y + 0.15, color=CORAL, width=1.2)


def add_atac_fragment(slide, x, y, w, color):
    """Fragment with visible endpoint dots and short adapter-colored caps."""
    add_line(slide, x, y, x + w, y, color=color, width=3.8)
    cap_w = min(0.12, w * 0.24)
    add_line(slide, x, y, x + cap_w, y, color=CORAL, width=3.8)
    add_line(slide, x + w - cap_w, y, x + w, y, color=CORAL, width=3.8)
    for end_x in (x, x + w):
        add_circle(slide, end_x - 0.035, y - 0.035, 0.07, NAVY, line=NAVY, line_width=0.3)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # 1 — Title
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_pill(slide, "PROJECT OVERVIEW • RESULTS-FREE DRAFT", 0.72, 0.55, 3.28, 0.38, TEAL, size=10.5)
    add_text(slide, "ShapeMix-ATAC", 0.72, 1.28, 7.2, 0.8, size=40, color=WHITE, bold=True, font=FONT_DISPLAY)
    add_text(
        slide,
        "Can fragment shape help us unmix spatial chromatin data?",
        0.74,
        2.16,
        6.5,
        1.08,
        size=26,
        color=LIGHT_TEXT,
        bold=False,
        font=FONT_DISPLAY,
    )
    add_text(slide, "Andy Zhuang  •  High School Science Research", 0.76, 6.54, 5.9, 0.35, size=14, color=LIGHT_TEXT)
    # Decorative mixed spot
    add_circle(slide, 8.55, 1.03, 3.45, NAVY_2, line=TEAL, line_width=3)
    cells = [
        (9.02, 1.55, 0.72, BLUE),
        (10.08, 1.31, 0.64, CORAL),
        (10.78, 2.03, 0.82, GOLD),
        (9.58, 2.45, 0.92, TEAL),
        (8.95, 3.05, 0.66, PURPLE),
        (10.60, 3.17, 0.70, BLUE),
        (9.66, 3.72, 0.62, CORAL),
    ]
    for x, y, d, color in cells:
        add_circle(slide, x, y, d, color, line=WHITE, line_width=2)
        add_circle(slide, x + d * 0.34, y + d * 0.32, d * 0.30, WHITE, line=WHITE, line_width=0.5)
    add_text(slide, "mixed spatial spot", 8.78, 4.64, 3.0, 0.35, size=14, color=LIGHT_TEXT, bold=True, align=PP_ALIGN.CENTER)
    # fragment stream
    for i, (length, color) in enumerate([(0.55, TEAL), (0.95, GOLD), (0.42, TEAL), (1.26, CORAL), (0.70, GOLD)]):
        add_fragment(slide, 8.68 + i * 0.70, 5.43 + (i % 2) * 0.28, length, color, line_width=4, cut_color=LIGHT_TEXT)
    add_notes(slide, "Open with the research question. Explain that this is a project overview: the biology, model idea, data, and test plan are ready to discuss, while results are intentionally not shown in this draft.")

    # 2 — Background
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    add_title(slide, "Chromatin controls which DNA instructions can be used", "Biology background", slide_num=2)
    add_card(slide, 0.68, 1.36, 3.54, 4.98, fill=WHITE, line=MID)
    add_text(slide, "Same DNA, different cells", 0.98, 1.72, 2.95, 0.52, size=23, color=NAVY, bold=True)
    add_bullets(
        slide,
        [
            "Most cells contain nearly the same genome.",
            "Cell types differ in which DNA regions are open and active.",
            "Those open regions help reveal cell identity and state.",
        ],
        0.98,
        2.48,
        2.92,
        2.64,
        size=17.5,
    )
    add_pill(slide, "CHROMATIN = DNA + PACKAGING", 1.06, 5.56, 2.75, 0.42, NAVY, size=10.5)
    # Closed/open comparison
    add_text(slide, "CLOSED", 4.83, 1.46, 2.7, 0.3, size=13, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "OPEN", 9.16, 1.46, 2.7, 0.3, size=13, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, 4.54, 1.83, 3.33, 3.45, fill=PALE, line=MID)
    add_card(slide, 8.83, 1.83, 3.33, 3.45, fill=TEAL_PALE, line=TEAL)
    # tightly packed chromatin
    for i in range(7):
        x = 4.89 + (i % 4) * 0.66
        y = 2.45 + (i // 4) * 0.72 + (i % 2) * 0.14
        add_circle(slide, x, y, 0.48, PURPLE, line=WHITE, line_width=1)
    add_line(slide, 4.82, 2.70, 7.52, 3.72, color=CORAL, width=3)
    add_line(slide, 4.82, 3.72, 7.52, 2.70, color=CORAL, width=3)
    # open chromatin
    add_line(slide, 9.18, 3.62, 11.80, 2.38, color=BLUE, width=4)
    for x, y in [(9.36, 3.52), (10.20, 3.12), (11.02, 2.73), (11.56, 2.47)]:
        add_circle(slide, x, y - 0.17, 0.34, GOLD, line=WHITE, line_width=1)
    add_pill(slide, "Tn5", 10.08, 3.90, 0.78, 0.38, CORAL, size=12)
    add_line(slide, 10.18, 3.91, 10.02, 3.50, color=CORAL, width=1.8)
    add_line(slide, 10.75, 3.91, 10.88, 3.13, color=CORAL, width=1.8)
    add_text(slide, "ATAC-seq uses Tn5 to cut and tag accessible DNA.", 5.04, 5.61, 6.82, 0.58, size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Source: Buenrostro et al., Nature Methods (2013). Diagram is conceptual.")
    add_notes(slide, "Students only need one idea here: accessible DNA is easier for the Tn5 enzyme to reach. ATAC-seq turns those accessible regions into measurable DNA fragments.")

    # 3 — ATAC-seq: aligned chromatin, fragments, and peak-count chart
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_title(
        slide,
        "ATAC-seq turns accessible DNA into peak counts",
        "ATAC-seq basics",
        slide_num=3,
    )
    add_text(
        slide,
        "The same genomic regions line up from chromatin → fragments → counts.",
        0.70,
        1.15,
        11.90,
        0.32,
        size=15.5,
        color=SLATE,
        align=PP_ALIGN.CENTER,
    )

    # Three full-width rows share one genomic x-coordinate system.
    row_x, row_w = 0.72, 11.90
    add_card(slide, row_x, 1.55, row_w, 1.39, fill=PALE, line=MID)
    add_card(slide, row_x, 3.05, row_w, 1.56, fill=WHITE, line=MID)
    add_card(slide, row_x, 4.72, row_w, 1.86, fill=BLUE_PALE, line=BLUE)

    # Subtle peak-aligned bands make the three rows refer to the same DNA.
    peak_columns = [
        (3.05, 2.30, "F2F6FE"),
        (6.05, 2.60, "EEF9F7"),
        (9.45, 2.35, "FFF8E8"),
    ]
    for x, w, fill in peak_columns:
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, 1.66, w, 1.17, fill, line=fill, line_width=0.2)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, 3.16, w, 1.34, fill, line=fill, line_width=0.2)

    # Row 1: chromatin. Purple cores are histone proteins; the dark DNA thread
    # visibly wraps around them, following the supplied ATAC-seq reference.
    add_circle(slide, 0.94, 1.76, 0.38, TEAL, line=TEAL, line_width=0.8)
    add_text(
        slide,
        "1",
        0.94,
        1.76,
        0.38,
        0.34,
        size=14,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )
    add_text(slide, "CHROMATIN", 1.43, 1.82, 1.40, 0.27, size=12, color=TEAL_DARK, bold=True)

    dna_y = 2.35
    add_line(slide, 2.72, dna_y, 12.18, dna_y, color="654640", width=1.7)
    for center_x in (4.16, 6.82, 7.82, 10.56, 12.02):
        add_nucleosome(slide, center_x, dna_y, scale=0.92)

    # Tn5 is shown only on exposed linker DNA, as in the reference figure.
    for marker_x in (3.48, 5.22, 6.20, 8.52, 9.72, 11.42):
        add_tn5_marker(slide, marker_x, dna_y)
    add_pill(slide, "Tn5", 5.02, 1.73, 0.56, 0.24, BLUE, size=8.2)
    add_text(
        slide,
        "DNA wrapped around histones",
        3.18,
        2.65,
        1.92,
        0.18,
        size=9.5,
        color=PURPLE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "exposed linker DNA",
        5.45,
        2.65,
        1.38,
        0.18,
        size=9.5,
        color=TEAL_DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Row 2: distinct fragments of different lengths on separate lanes.
    add_circle(slide, 0.94, 3.31, 0.38, CORAL, line=CORAL, line_width=0.8)
    add_text(
        slide,
        "2",
        0.94,
        3.31,
        0.38,
        0.34,
        size=14,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )
    add_text(slide, "FRAGMENTS", 1.43, 3.37, 1.40, 0.27, size=12, color=CORAL, bold=True)

    # Compact key: length can suggest whether a fragment spans nucleosomes.
    add_pill(slide, "SHORT • OPEN DNA", 2.88, 3.20, 1.34, 0.24, CORAL, size=8.0)
    add_pill(slide, "MEDIUM • SPANS 1 NUCLEOSOME", 5.52, 3.20, 2.12, 0.24, GOLD, color=NAVY, size=7.6)
    add_pill(slide, "LONG • SPANS 2 NUCLEOSOMES", 7.77, 3.20, 2.04, 0.24, PURPLE, size=7.6)

    fragments = [
        # Peak 1: two short open-DNA fragments and one that spans a nucleosome.
        (3.23, 3.70, 0.56, CORAL),
        (3.79, 4.03, 0.77, GOLD),
        (4.46, 4.35, 0.66, CORAL),
        # Peak 2: one mono-nucleosome-like and one longer protected fragment.
        (6.43, 3.73, 1.23, GOLD),
        (6.20, 4.22, 2.22, PURPLE),
        # Peak 3: one fragment with cuts flanking a nucleosome.
        (9.82, 3.99, 1.48, GOLD),
    ]
    for x, y, w, color in fragments:
        add_atac_fragment(slide, x, y, w, color)

    # Row 3 is the bar chart. Bar widths occupy the same genomic intervals as
    # the peak-aligned columns above, while bar heights encode endpoint counts.
    add_circle(slide, 0.94, 4.98, 0.38, BLUE, line=BLUE, line_width=0.8)
    add_text(
        slide,
        "3",
        0.94,
        4.98,
        0.38,
        0.34,
        size=14,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )
    add_text(slide, "PEAK COUNTS", 1.43, 5.04, 1.40, 0.27, size=12, color=BLUE, bold=True)
    add_pill(slide, "HEIGHT = COUNT  •  WIDTH = PEAK INTERVAL", 9.59, 4.84, 2.51, 0.24, NAVY, size=7.7)

    chart_x1, chart_x2, baseline, unit_h = 2.95, 12.18, 6.23, 0.17
    for value in (0, 2, 4, 6):
        y = baseline - value * unit_h
        add_line(slide, chart_x1, y, chart_x2, y, color=MID, width=0.7)
        add_text(slide, str(value), 2.62, y - 0.09, 0.24, 0.18, size=8.5, color=SLATE, align=PP_ALIGN.RIGHT)
    add_line(slide, chart_x1, baseline, chart_x1, baseline - 6 * unit_h, color=NAVY, width=1.0)
    chart_bars = [
        (3.05, 2.30, 6, BLUE, "Peak 1"),
        (6.05, 2.60, 4, TEAL, "Peak 2"),
        (9.45, 2.35, 2, GOLD, "Peak 3"),
    ]
    for x, width, value, color, label in chart_bars:
        height = value * unit_h
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, baseline - height, width, height, color, line=color)
        add_text(slide, str(value), x, baseline - height - 0.24, width, 0.20, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x, baseline + 0.05, width, 0.20, size=9.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_pill(
        slide,
        "COUNT THE ENDPOINT DOTS: 3 fragments → 6 counts   |   2 → 4   |   1 → 2",
        2.14,
        6.72,
        9.05,
        0.29,
        TEAL,
        size=9.7,
    )
    add_footer(
        slide,
        "Conceptual ATAC-seq schematic. Some pipelines count fragment overlaps instead; the convention must be stated.",
    )
    add_notes(
        slide,
        "Read the slide from top to bottom at the same horizontal genomic position. The purple shapes are histone protein cores; the dark DNA thread wraps around each core to form a nucleosome. Tn5 inserts adapters most readily in exposed linker DNA, and paired-end sequencing locates the two endpoints of each resulting fragment. Fragment lengths vary: short fragments are often nucleosome-free-like, approximately 200-base-pair fragments can span one nucleosome, and longer fragments can span two; length is informative but not a perfect label. The fragments are kept on separate lanes so they never look like one continuous segment. In the teaching convention used here, each endpoint inside a selected peak contributes one count. Peak 1 has three fragments and six endpoints, Peak 2 has two and four, and Peak 3 has one and two. The bar widths align with the selected genomic intervals and the gaps show DNA outside peaks. Some pipelines count fragments or overlaps instead, so a real analysis must state its counting convention.",
    )

    # 4 — Spatial mixture
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_title(slide, "One spatial spot can contain several cell types", "The problem", slide_num=4)
    add_text(slide, "A spatial assay preserves location—but each spot may collect fragments from multiple nearby nuclei.", 0.7, 1.18, 11.7, 0.48, size=19, color=SLATE)
    # tissue grid
    add_card(slide, 0.72, 1.86, 5.10, 4.54, fill=PALE, line=MID)
    add_text(slide, "Tissue section", 1.00, 2.08, 2.2, 0.35, size=18, color=NAVY, bold=True)
    grid_colors = [BLUE, BLUE, TEAL, TEAL, PURPLE, BLUE, CORAL, TEAL, PURPLE, CORAL, CORAL, GOLD, PURPLE, PURPLE, GOLD, GOLD]
    k = 0
    for row in range(4):
        for col in range(4):
            x = 1.08 + col * 0.93
            y = 2.72 + row * 0.78
            d = 0.58
            add_circle(slide, x, y, d, grid_colors[k], line=WHITE, line_width=2)
            if k == 6:
                ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.12), Inches(y - 0.12), Inches(d + 0.24), Inches(d + 0.24))
                ring.fill.background()
                ring.line.color.rgb = C(NAVY)
                ring.line.width = Pt(3)
            k += 1
    add_text(slide, "each circle = one measured spot", 1.02, 5.93, 3.82, 0.28, size=13, color=SLATE)
    # arrow and zoom
    arrow = add_shape(slide, MSO_SHAPE.RIGHT_ARROW, 5.95, 3.34, 1.10, 0.72, TEAL, line=TEAL)
    add_card(slide, 7.27, 1.86, 5.36, 4.54, fill=NAVY, line=NAVY)
    add_text(slide, "Zoom into one spot", 7.64, 2.08, 3.9, 0.35, size=18, color=WHITE, bold=True)
    for x, y, d, color in [
        (7.74, 2.84, 1.04, BLUE),
        (8.76, 3.26, 0.92, CORAL),
        (9.57, 2.70, 1.12, GOLD),
        (10.62, 3.40, 0.88, TEAL),
        (11.36, 2.69, 0.82, BLUE),
    ]:
        add_circle(slide, x, y, d, color, line=WHITE, line_width=2)
        add_circle(slide, x + d * 0.34, y + d * 0.33, d * 0.27, WHITE, line=WHITE, line_width=0.5)
    for i, (length, color) in enumerate([(0.44, BLUE), (0.72, CORAL), (0.93, GOLD), (0.56, TEAL), (0.80, BLUE)]):
        add_fragment(slide, 7.70 + i * 0.87, 4.70 + (i % 2) * 0.30, length, color, line_width=4, cut_color=WHITE)
    add_text(slide, "All fragments arrive mixed together", 7.68, 5.55, 4.48, 0.40, size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_pill(slide, "WHO IS IN THE SPOT—and in what proportions?", 7.68, 6.00, 4.46, 0.38, TEAL, size=11)
    add_footer(slide, "Source: Deng et al., Nature (2022). Spot drawing is conceptual.")
    add_notes(slide, "A spatial spot is like a collection bucket placed over a tiny region of tissue. It keeps the location, but it may mix fragments from several cell types. The task is to estimate the ingredients of that mixture.")

    # 5 — Deconvolution
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    add_title(slide, "Deconvolution works like identifying ingredients in a smoothie", "Core task", slide_num=5)
    add_text(slide, "Known reference fingerprints + an unknown mixture → estimated cell-type proportions", 0.71, 1.17, 11.9, 0.46, size=20, color=SLATE, align=PP_ALIGN.CENTER)
    # reference cards
    labels = [("B cell", BLUE, [4, 2, 1]), ("T cell", TEAL, [1, 4, 2]), ("Monocyte", CORAL, [1, 2, 4])]
    for i, (name, color, vals) in enumerate(labels):
        y = 1.92 + i * 1.32
        add_card(slide, 0.76, y, 3.08, 1.02, fill=WHITE, line=color)
        add_circle(slide, 0.99, y + 0.20, 0.58, color, line=WHITE, line_width=1.5)
        add_text(slide, name, 1.72, y + 0.17, 1.46, 0.30, size=17, color=NAVY, bold=True)
        for j, v in enumerate(vals):
            add_shape(slide, MSO_SHAPE.RECTANGLE, 1.75 + j * 0.43, y + 0.62 - v * 0.07, 0.24, v * 0.07, color, line=color)
        add_text(slide, "reference", 3.06, y + 0.68, 0.58, 0.18, size=9, color=SLATE, align=PP_ALIGN.RIGHT)
    add_text(slide, "reference library", 1.34, 6.16, 1.85, 0.27, size=13, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    # plus
    add_text(slide, "+", 4.16, 3.49, 0.52, 0.62, size=35, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    # mixed spot cup
    cup = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(4.92), Inches(2.25), Inches(2.22), Inches(3.26))
    cup.fill.solid(); cup.fill.fore_color.rgb = C(WHITE); cup.fill.transparency = 12
    cup.line.color.rgb = C(NAVY); cup.line.width = Pt(2)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 5.17, 3.49, 1.72, 1.75, PURPLE_PALE, line=PURPLE_PALE)
    for i, color in enumerate([BLUE, TEAL, CORAL, BLUE, CORAL, TEAL, BLUE, TEAL]):
        add_circle(slide, 5.28 + (i % 3) * 0.46, 3.78 + (i // 3) * 0.40, 0.28, color, line=WHITE, line_width=1)
    add_text(slide, "unknown mixed spot", 4.85, 5.74, 2.39, 0.31, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    # arrow model
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, 7.47, 3.22, 1.15, 0.70, TEAL, line=TEAL)
    add_pill(slide, "MODEL", 7.58, 4.12, 0.92, 0.34, NAVY, size=10.5)
    # output
    add_card(slide, 8.91, 2.10, 3.75, 3.62, fill=WHITE, line=NAVY)
    add_text(slide, "Estimated recipe", 9.22, 2.40, 3.12, 0.37, size=21, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_stacked_bar(slide, 9.36, 3.05, 2.80, 0.62, [50, 30, 20], [BLUE, TEAL, CORAL], ["50%", "30%", "20%"])
    for i, (name, color, pct) in enumerate([("B cell", BLUE, "50%"), ("T cell", TEAL, "30%"), ("Monocyte", CORAL, "20%")]):
        add_circle(slide, 9.35, 4.02 + i * 0.45, 0.20, color, line=color, line_width=0.5)
        add_text(slide, name, 9.66, 3.96 + i * 0.45, 1.54, 0.28, size=14, color=INK)
        add_text(slide, pct, 11.45, 3.96 + i * 0.45, 0.58, 0.28, size=14, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
    add_footer(slide, "Background: Ouologuem et al., Bioinformatics (2025). Values shown are a teaching example, not data.")
    add_notes(slide, "Use the smoothie analogy literally: references tell us what each pure ingredient tends to look like; the model asks which combination best matches the mixed measurement. The output is a proportion for every cell type.")

    # 6 — Existing algorithms: common setup
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_title(slide, "Most existing algorithms start from the same input table", "Existing algorithms", slide_num=6)
    add_text(
        slide,
        "Methods first collapse every ATAC peak to one count, then ask which reference mixture best reconstructs the spot.",
        0.72,
        1.14,
        11.84,
        0.48,
        size=19,
        color=SLATE,
        align=PP_ALIGN.CENTER,
    )
    # Reference matrix
    add_card(slide, 0.76, 1.89, 3.35, 4.34, fill=BLUE_PALE, line=BLUE)
    add_pill(slide, "REFERENCE", 1.00, 2.14, 1.18, 0.35, BLUE, size=9.5)
    add_text(slide, "cell type × peak", 2.31, 2.10, 1.45, 0.38, size=15, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, "A", 0.99, 2.70, 0.42, 0.32, size=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "B", 0.99, 3.29, 0.42, 0.32, size=16, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "C", 0.99, 3.88, 0.42, 0.32, size=16, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    matrix = [
        [90, 24, 42, 18],
        [20, 82, 31, 54],
        [36, 19, 88, 46],
    ]
    row_colors = [BLUE, TEAL, CORAL]
    for r_i, row in enumerate(matrix):
        for c_i, value in enumerate(row):
            x = 1.52 + c_i * 0.56
            y = 2.68 + r_i * 0.59
            add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.44, 0.40, row_colors[r_i], line=WHITE)
            add_text(slide, str(value), x, y + 0.02, 0.44, 0.28, size=8.7, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for c_i in range(4):
        add_text(slide, f"P{c_i + 1}", 1.50 + c_i * 0.56, 4.46, 0.48, 0.22, size=9.5, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "A fingerprint learned from labeled single cells", 1.06, 5.05, 2.75, 0.60, size=15.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
    # Mixed observation
    add_text(slide, "+", 4.25, 3.59, 0.47, 0.56, size=32, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, 4.88, 1.89, 3.13, 4.34, fill=GOLD_PALE, line=GOLD)
    add_pill(slide, "MIXED SPOT", 5.14, 2.14, 1.36, 0.35, GOLD, color=NAVY, size=9.5)
    add_text(slide, "spot × peak", 6.57, 2.10, 1.08, 0.38, size=15, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
    for i, height in enumerate([0.72, 1.44, 0.96, 1.18]):
        add_shape(slide, MSO_SHAPE.RECTANGLE, 5.34 + i * 0.55, 4.42 - height, 0.33, height, GOLD, line=GOLD)
        add_text(slide, f"P{i + 1}", 5.27 + i * 0.55, 4.56, 0.48, 0.22, size=9.5, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Only total cut sites in each peak are passed to the model", 5.21, 5.05, 2.47, 0.66, size=15.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
    # Output
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, 8.19, 3.47, 0.88, 0.64, TEAL, line=TEAL)
    add_card(slide, 9.24, 1.89, 3.31, 4.34, fill=TEAL_PALE, line=TEAL)
    add_pill(slide, "OUTPUT", 9.51, 2.14, 0.92, 0.35, TEAL, size=9.5)
    add_text(slide, "cell-type proportions", 10.50, 2.10, 1.68, 0.38, size=15, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
    add_stacked_bar(slide, 9.72, 3.00, 2.36, 0.56, [50, 30, 20], [BLUE, TEAL, CORAL], ["50", "30", "20"])
    for i, (name, color, pct) in enumerate([("Type A", BLUE, "50%"), ("Type B", TEAL, "30%"), ("Type C", CORAL, "20%")]):
        add_circle(slide, 9.74, 3.91 + i * 0.48, 0.22, color, line=color, line_width=0.5)
        add_text(slide, name, 10.08, 3.85 + i * 0.48, 1.10, 0.29, size=13.5, color=INK)
        add_text(slide, pct, 11.33, 3.85 + i * 0.48, 0.72, 0.29, size=13.5, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
    add_pill(slide, "SHARED LIMIT", 10.05, 5.50, 1.44, 0.34, CORAL, size=9)
    add_text(slide, "Fragment lengths are gone", 9.66, 5.82, 2.48, 0.34, size=12.5, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Background: Ouologuem et al., Bioinformatics (2025). Diagram is conceptual.")
    add_notes(slide, "This slide establishes the common setup before comparing individual methods. A labeled single-cell reference becomes a cell-type-by-peak fingerprint. The spatial measurement becomes one total count per peak. The algorithms differ in how they fit the mixture, but this transferred-ATAC setup usually discards the parent-fragment length before inference starts.")

    # 7 — Existing algorithms: families
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    add_title(slide, "Existing methods use different mathematical engines", "Existing algorithms", slide_num=7)
    add_text(slide, "All solve an inverse problem: work backward from a mixed measurement to the hidden recipe.", 0.72, 1.14, 11.85, 0.43, size=19, color=SLATE, align=PP_ALIGN.CENTER)
    families = [
        (
            "LINEAR FITTING",
            "NNLS  •  SpatialDWLS",
            "Choose nonnegative amounts so a weighted sum of references is close to the observed peak totals.",
            "y ≈ Aᵀz",
            "Fast and easy to inspect",
            "Noise is handled mainly through the loss and weights",
            BLUE,
            BLUE_PALE,
        ),
        (
            "PROBABILISTIC COUNTS",
            "RCTD  •  cell2location",
            "Describe how sequencing counts could be generated from a cell-type mixture, with explicit noise and technical effects.",
            "p(counts | mixture)",
            "Makes count noise explicit",
            "Needs distributions, priors, and more computation",
            PURPLE,
            PURPLE_PALE,
        ),
        (
            "MAPPING / LATENT MODELS",
            "Tangram  •  DestVI",
            "Optimize a cell-to-space map or learn hidden cell states that can reconstruct the spatial profile.",
            "learn a hidden map",
            "Flexible for complex biology",
            "More moving parts can be harder to interpret",
            TEAL,
            TEAL_PALE,
        ),
    ]
    for i, (family, methods, body, equation, strength, tradeoff, color, fill) in enumerate(families):
        x = 0.72 + i * 4.17
        add_card(slide, x, 1.78, 3.77, 4.91, fill=fill, line=color)
        add_pill(slide, family, x + 0.26, 2.03, 2.92, 0.35, color, size=9.2)
        add_text(slide, methods, x + 0.27, 2.63, 3.20, 0.38, size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_card(slide, x + 0.52, 3.24, 2.73, 0.72, fill=WHITE, line=color)
        add_text(slide, equation, x + 0.60, 3.39, 2.57, 0.38, size=20, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, body, x + 0.34, 4.22, 3.08, 0.88, size=14.3, color=INK, align=PP_ALIGN.CENTER)
        add_line(slide, x + 0.34, 5.28, x + 3.43, 5.28, color=MID, width=0.9)
        add_text(slide, "STRENGTH", x + 0.36, 5.48, 0.92, 0.24, size=9.5, color=color, bold=True)
        add_text(slide, strength, x + 1.29, 5.44, 2.06, 0.35, size=12.2, color=NAVY, bold=True)
        add_text(slide, "TRADE-OFF", x + 0.36, 6.02, 0.92, 0.24, size=9.5, color=CORAL, bold=True)
        add_text(slide, tradeoff, x + 1.29, 5.94, 2.06, 0.50, size=11.5, color=INK)
    add_footer(slide, "Sources: Tsoucas et al. (2019); Dong & Yuan (2021); Cable et al. (2022); Kleshchevnikov et al. (2022); Biancalani et al. (2021); Lopez et al. (2022).")
    add_notes(slide, "NNLS and SpatialDWLS are linear mixture methods; SpatialDWLS adds feature selection and weights. RCTD and cell2location write probabilistic count-generating models. Tangram optimizes a mapping from single cells to space, while DestVI uses a deep generative latent-variable model. These categories are simplified for teaching. In the published spatial-ATAC benchmark, peaks were substituted for genes, so these methods still used peak totals rather than parent-fragment shape.")

    # 8 — Shape clue
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_title(slide, "ShapeMix keeps an extra clue: fragment length", "The new idea", slide_num=8)
    add_rich_text(
        slide,
        [("Two samples can have the same peak total", {"bold": True, "color": NAVY}), (" but different fragment-length patterns.", {"color": SLATE})],
        0.71,
        1.16,
        11.9,
        0.46,
        size=20,
        align=PP_ALIGN.CENTER,
    )
    # two panels
    for x, title, fragment_counts in [(0.76, "Cell type A", [7, 3, 2]), (6.78, "Cell type B", [2, 4, 6])]:
        add_card(slide, x, 1.82, 5.78, 3.34, fill=PALE, line=MID)
        add_text(slide, title, x + 0.26, 2.07, 2.0, 0.34, size=19, color=NAVY, bold=True)
        cut_site_counts = [2 * count for count in fragment_counts]
        add_pill(slide, f"{sum(cut_site_counts)} cut sites in this peak", x + 3.03, 2.03, 2.20, 0.37, NAVY, size=10.5)
        colors = [TEAL] * fragment_counts[0] + [GOLD] * fragment_counts[1] + [CORAL] * fragment_counts[2]
        widths = [0.30] * fragment_counts[0] + [0.50] * fragment_counts[1] + [0.70] * fragment_counts[2]
        # Center each fragment in a fixed-width cell. Keeping the longest glyph
        # narrower than the 0.82-inch pitch prevents adjacent red fragments from
        # merging into one continuous bar while preserving the 6-by-2 count grid.
        for i, (color, width) in enumerate(zip(colors, widths)):
            row = i // 6
            col = i % 6
            cell_center = x + 0.84 + col * 0.82
            add_fragment(slide, cell_center - width / 2, 2.91 + row * 0.70, width, color, line_width=4, cut_color=INK)
        add_stacked_bar(slide, x + 0.42, 4.47, 4.90, 0.38, cut_site_counts, [TEAL, GOLD, CORAL])
    add_text(slide, "Count-only model sees: 24 = 24", 1.70, 5.40, 4.70, 0.38, size=17, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "ShapeMix also sees: the color pattern differs", 6.67, 5.40, 5.72, 0.38, size=17, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
    # bins
    bins = [
        ("SHORT", "< 100 bp", TEAL, "often nucleosome-free"),
        ("MIDDLE", "100–249 bp", GOLD, "mononucleosome-like"),
        ("LONG", "≥ 250 bp", CORAL, "longer protected fragments"),
    ]
    for i, (name, cutoff, color, note) in enumerate(bins):
        x = 1.03 + i * 4.05
        add_card(slide, x, 6.00, 3.72, 0.76, fill=WHITE, line=color)
        add_pill(slide, name, x + 0.12, 6.20, 0.92, 0.31, color, size=9.5)
        add_text(slide, cutoff, x + 1.18, 6.08, 2.34, 0.22, size=12.5, color=NAVY, bold=True)
        add_text(slide, note, x + 1.18, 6.40, 2.34, 0.19, size=9.5, color=SLATE)
    add_footer(slide, "Sources: Buenrostro et al. (2013); Martens et al., Nature Methods (2024). Diagram is conceptual.")
    add_notes(slide, "This is the project’s central idea. Each panel shows 12 fragments, and this conceptual example assigns both Tn5 cut sites from every fragment to the same peak, for 24 cut sites total. Cell type A allocates those cuts 14 short, 6 middle, and 4 long; cell type B allocates them 4 short, 8 middle, and 12 long. Ordinary deconvolution keeps only the total. ShapeMix also records the parent-fragment length pattern. The biological labels for bins are useful approximations, not perfect one-to-one assignments.")

    # 9 — Bayesian thinking
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    add_title(slide, "Bayesian modeling is a scientific story run backward", "ShapeMix: Bayesian model", slide_num=9)
    add_text(slide, "First describe how a hidden cell mixture could generate the data. Then use the data to infer the mixture.", 0.72, 1.14, 11.84, 0.48, size=19, color=SLATE, align=PP_ALIGN.CENTER)
    # Forward story
    add_card(slide, 0.76, 1.87, 5.29, 4.60, fill=WHITE, line=MID)
    add_pill(slide, "FORWARD MODEL", 1.02, 2.12, 1.45, 0.36, NAVY, size=9.5)
    forward = [
        ("1", "Hidden recipe", "possible cell amounts z", BLUE),
        ("2", "Expected signal", "predicted counts + length split", PURPLE),
        ("3", "Observed data", "noisy cut-site measurements", TEAL),
    ]
    for i, (num, title, body, color) in enumerate(forward):
        y = 2.75 + i * 1.04
        add_circle(slide, 1.12, y, 0.52, color, line=color, line_width=1)
        add_text(slide, num, 1.12, y + 0.01, 0.52, 0.40, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, title, 1.87, y - 0.06, 1.70, 0.40, size=17, color=NAVY, bold=True)
        add_text(slide, body, 3.52, y - 0.02, 2.15, 0.40, size=13.2, color=SLATE)
        if i < 2:
            add_shape(slide, MSO_SHAPE.DOWN_ARROW, 1.25, y + 0.56, 0.26, 0.39, color, line=color)
    add_text(slide, "If the recipe were known, what data would we expect?", 1.10, 5.92, 4.60, 0.38, size=14.5, color=SLATE, italic=True, align=PP_ALIGN.CENTER)
    # Reverse inference
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, 6.21, 3.66, 0.76, 0.62, TEAL, line=TEAL)
    add_text(slide, "reverse", 6.26, 4.36, 0.66, 0.23, size=9.5, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, 7.15, 1.87, 5.41, 4.60, fill=NAVY, line=NAVY)
    add_pill(slide, "BAYES' RULE", 7.49, 2.12, 1.24, 0.36, TEAL, size=9.5)
    add_text(slide, "posterior  ∝  likelihood  ×  prior", 7.52, 2.83, 4.68, 0.55, size=23, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    bayes_parts = [
        ("PRIOR", "plausible positive amounts before this spot", GOLD),
        ("LIKELIHOOD", "how well a recipe explains the observed data", PURPLE),
        ("POSTERIOR", "updated support after combining both", TEAL),
    ]
    for i, (label, body, color) in enumerate(bayes_parts):
        y = 3.62 + i * 0.73
        add_pill(slide, label, 7.54, y, 1.18, 0.33, color, color=NAVY if color == GOLD else WHITE, size=8.7)
        add_text(slide, body, 8.94, y - 0.02, 3.02, 0.39, size=13.3, color=LIGHT_TEXT)
    add_text(slide, "The prior guides weak evidence; it does not predetermine the answer.", 7.60, 5.96, 4.52, 0.29, size=13.5, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Bayesian vocabulary: prior + likelihood → posterior. ShapeMix uses a MAP point estimate from that posterior.")
    add_notes(slide, "Bayesian does not mean guessing. The model first states a forward scientific story: a hidden recipe produces expected ATAC signals, and measurements fluctuate around them. Inference runs that story backward. The prior rules out impossible negative amounts and discourages unsupported extremes; the likelihood lets the observed data update the answer. ShapeMix's MVP summarizes the posterior by its highest point, called MAP.")

    # 10 — Generative equations
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_title(slide, "The generative story links cell amounts to expected data", "ShapeMix: Bayesian model", slide_num=10)
    add_text(slide, "Two reference fingerprints are fixed from labeled training cells; only each spot's effective cell amounts are inferred.", 0.72, 1.14, 11.84, 0.46, size=18.5, color=SLATE, align=PP_ALIGN.CENTER)
    # Input cards
    add_card(slide, 0.76, 1.82, 3.00, 1.63, fill=BLUE_PALE, line=BLUE)
    add_pill(slide, "UNKNOWN", 1.02, 2.07, 0.95, 0.34, BLUE, size=9)
    add_text(slide, "z[s,c]", 2.36, 1.98, 0.90, 0.52, size=29, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "effective amount of cell type c in spot s", 1.05, 2.66, 2.42, 0.46, size=14, color=INK, align=PP_ALIGN.CENTER)
    add_card(slide, 0.76, 3.70, 3.00, 1.34, fill=PURPLE_PALE, line=PURPLE)
    add_pill(slide, "REFERENCE", 1.02, 3.96, 1.02, 0.34, PURPLE, size=9)
    add_text(slide, "A[c,p]", 2.34, 3.89, 0.96, 0.44, size=25, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "mean cut-site rate at peak p", 1.04, 4.45, 2.47, 0.42, size=13.2, color=INK, align=PP_ALIGN.CENTER)
    add_card(slide, 0.76, 5.27, 3.00, 1.34, fill=GOLD_PALE, line=GOLD)
    add_pill(slide, "REFERENCE", 1.02, 5.53, 1.02, 0.34, GOLD, color=NAVY, size=9)
    add_text(slide, "ω[c,p,b]", 2.24, 5.46, 1.17, 0.44, size=25, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "fraction expected in length bin b", 1.04, 6.02, 2.47, 0.42, size=13.2, color=INK, align=PP_ALIGN.CENTER)
    # Equations
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, 3.98, 3.65, 0.84, 0.64, TEAL, line=TEAL)
    add_card(slide, 5.02, 1.82, 4.34, 2.20, fill=PALE, line=NAVY)
    add_pill(slide, "EXPECTED PEAK TOTAL", 5.33, 2.08, 1.95, 0.34, NAVY, size=9)
    add_text(slide, "μ[s,p] = Σ over c  z[s,c] A[c,p]", 5.25, 2.67, 3.89, 0.52, size=22, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Add the cell types' expected contributions.", 5.55, 3.33, 3.29, 0.42, size=14.2, color=SLATE, align=PP_ALIGN.CENTER)
    add_card(slide, 5.02, 4.31, 4.34, 2.30, fill=TEAL_PALE, line=TEAL)
    add_pill(slide, "EXPECTED LENGTH MIX", 5.33, 4.58, 1.95, 0.34, TEAL, size=9)
    add_text(slide, "ρ[s,p,b] = Σ over c  z[s,c] A[c,p] ω[c,p,b] / μ[s,p]", 5.16, 5.12, 4.08, 0.72, size=16.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Weight each type's shape fingerprint, then divide by the total.", 5.40, 5.94, 3.58, 0.44, size=13.5, color=SLATE, align=PP_ALIGN.CENTER)
    # Observation
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, 9.57, 3.65, 0.84, 0.64, TEAL, line=TEAL)
    add_card(slide, 10.61, 1.82, 1.94, 4.79, fill=NAVY, line=NAVY)
    add_pill(slide, "OBSERVE", 10.98, 2.10, 1.18, 0.34, TEAL, size=9)
    for i, (name, color, count) in enumerate([("short", TEAL, "14"), ("middle", GOLD, "6"), ("long", CORAL, "4")]):
        y = 2.89 + i * 0.84
        add_circle(slide, 10.98, y, 0.35, color, line=WHITE, line_width=1)
        add_text(slide, name, 11.42, y - 0.02, 0.68, 0.27, size=11.5, color=LIGHT_TEXT)
        add_text(slide, count, 12.02, y - 0.03, 0.29, 0.29, size=13.5, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    add_line(slide, 10.96, 5.43, 12.18, 5.43, color=LIGHT_TEXT, width=0.8)
    add_text(slide, "N = 24", 10.96, 5.67, 1.24, 0.38, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Notation: s = spot, c = cell type, p = peak, b = fragment-length bin. Background is fixed to zero in model v1.")
    add_notes(slide, "z is the only unknown shown here. A is the learned mean peak count per training-reference cell type. Omega is the learned length-bin distribution for that type and peak. The model multiplies amount by accessibility to predict the total, and uses omega to predict how that total should be divided among short, middle, and long parent-fragment bins. The output z is depth-scaled effective abundance, not a calibrated nucleus count.")

    # 11 — Likelihoods
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    add_title(slide, "Two probability models score two kinds of evidence", "ShapeMix: likelihood", slide_num=11)
    add_text(slide, "The factorization protects the fair comparison: the count model stays identical, and ShapeMix adds only the conditional length term.", 0.72, 1.14, 11.84, 0.48, size=18.5, color=SLATE, align=PP_ALIGN.CENTER)
    # Negative binomial panel
    add_card(slide, 0.76, 1.84, 5.63, 4.56, fill=BLUE_PALE, line=BLUE)
    add_pill(slide, "1  PEAK TOTAL", 1.04, 2.11, 1.45, 0.36, BLUE, size=9.5)
    add_text(slide, "Negative binomial", 1.04, 2.70, 5.06, 0.43, size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "N[s,p] ~ NB(mean = μ[s,p])", 1.22, 3.35, 4.70, 0.50, size=23, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    for i, height in enumerate([0.55, 1.12, 0.82, 1.38, 0.69, 1.01]):
        add_shape(slide, MSO_SHAPE.RECTANGLE, 1.50 + i * 0.59, 5.08 - height, 0.34, height, BLUE, line=BLUE)
    add_text(slide, "Why not a simple fixed count?", 1.18, 5.38, 4.83, 0.31, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Sequencing counts vary more than a Poisson model usually allows. The negative binomial permits extra spread.", 1.24, 5.72, 4.72, 0.56, size=13.2, color=SLATE, align=PP_ALIGN.CENTER)
    # Multinomial panel
    add_card(slide, 6.69, 1.84, 5.86, 4.56, fill=TEAL_PALE, line=TEAL)
    add_pill(slide, "2  LENGTH SPLIT", 6.98, 2.11, 1.70, 0.36, TEAL, size=9.5)
    add_text(slide, "Conditional multinomial", 6.99, 2.70, 5.25, 0.43, size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Y[s,p,:] | N[s,p] ~ Multinomial(N[s,p], ρ[s,p,:])", 6.96, 3.31, 5.30, 0.62, size=17.5, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_stacked_bar(slide, 7.65, 4.26, 3.90, 0.62, [14, 6, 4], [TEAL, GOLD, CORAL], ["14", "6", "4"])
    add_text(slide, "Why condition on the total?", 7.12, 5.38, 5.04, 0.31, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "It asks only how the same 24 cuts split across bins, so the shape term cannot secretly replace the count model.", 7.16, 5.72, 4.96, 0.56, size=13.2, color=SLATE, align=PP_ALIGN.CENTER)
    add_pill(slide, "p(shape data) = p(total count) × p(length split | that total)", 3.36, 6.62, 6.62, 0.38, NAVY, size=10.8)
    add_footer(slide, "ShapeMix model v1: negative-binomial totals + conditional multinomial bins; not three independent negative-binomial models.")
    add_notes(slide, "The negative binomial is a probability model for noisy counts. It allows variance greater than the mean, which is common in sequencing data. The multinomial is like dividing a fixed number of colored marbles among three boxes. Conditioning on N means ShapeMix's extra term scores only the length composition. If there is just one bin, or every cell type has the same shape fingerprint, this extra term contributes no information.")

    # 12 — Inference and MAP
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_title(slide, "Inference searches for the most supported mixture", "ShapeMix: Bayesian inference", dark=True, slide_num=12)
    add_text(slide, "ShapeMix maximizes one score for every candidate abundance matrix z.", 0.72, 1.18, 11.84, 0.39, size=18.5, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)
    add_card(slide, 0.78, 1.80, 11.77, 1.23, fill=NAVY_2, line=TEAL)
    add_text(slide, "log posterior", 1.10, 2.06, 1.77, 0.38, size=19, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "=", 2.91, 2.05, 0.34, 0.38, size=21, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    score_parts = [
        ("count fit", "negative binomial", BLUE),
        ("+", "", WHITE),
        ("shape fit", "multinomial", TEAL),
        ("+", "", WHITE),
        ("Gamma prior", "positive, regularized z", PURPLE),
    ]
    x_positions = [3.28, 5.15, 5.63, 7.53, 8.04]
    widths = [1.66, 0.30, 1.67, 0.30, 2.24]
    for (label, sub, color), x, width in zip(score_parts, x_positions, widths):
        if label == "+":
            add_text(slide, label, x, 2.06, width, 0.37, size=21, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        else:
            add_pill(slide, label.upper(), x, 2.00, width, 0.34, color, size=8.5)
            add_text(slide, sub, x, 2.45, width, 0.24, size=9.3, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, "MAP targets the highest point of this posterior score", 10.37, 2.04, 1.72, 0.54, size=12.7, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    # Optimization path
    stages = [
        ("START", "NNLS estimate", "a sensible count-only guess", BLUE),
        ("CONSTRAIN", "log abundance", "z = exp(raw) stays positive", GOLD),
        ("UPDATE", "Adam optimizer", "moves uphill on the score", PURPLE),
        ("REPEAT", "3 restarts", "checks different starting paths", CORAL),
        ("REPORT", "best found MAP", "normalize z to percentages", TEAL),
    ]
    add_line(slide, 1.33, 4.77, 11.95, 4.77, color=LIGHT_TEXT, width=3)
    for i, (label, title, body, color) in enumerate(stages):
        x = 0.87 + i * 2.43
        add_circle(slide, x + 0.58, 4.37, 0.80, color, line=WHITE, line_width=1.5)
        add_text(slide, str(i + 1), x + 0.58, 4.39, 0.80, 0.61, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_pill(slide, label, x + 0.18, 3.55, 1.60, 0.32, color, size=8.5)
        add_text(slide, title, x - 0.02, 5.41, 2.02, 0.38, size=15.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x - 0.03, 5.87, 2.04, 0.48, size=11.7, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)
    add_card(slide, 2.52, 6.56, 8.31, 0.42, fill=CORAL, line=CORAL)
    add_text(slide, "Important: this MVP reports one MAP point estimate—not a credible interval or full posterior distribution.", 2.72, 6.59, 7.91, 0.32, size=11.7, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Frozen protocol v1: Gamma(2,1) prior; Adam; 3 deterministic restarts; identical settings in count-only and shape-aware arms.", dark=True)
    add_notes(slide, "ShapeMix performs MAP inference: it searches for the mixture with the largest posterior density. The count likelihood, shape likelihood, and Gamma prior are added on the log scale. The Gamma prior is defined only for positive amounts and discourages unsupported extremes. Optimization begins from a nonnegative least-squares solution, uses z = exp(raw) + 1e-8 to keep abundance positive, and runs Adam from three deterministic restarts. The best finite converged full objective is selected. This numerical search reduces but cannot eliminate the risk of a local maximum. It is not full posterior sampling, so the MVP does not report credible intervals.")

    # 13 — Algorithm
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_title(slide, "How ShapeMix works—four steps", "Algorithm overview", dark=True, slide_num=13)
    steps = [
        ("1", "Build references", "Use labeled single cells to learn each cell type’s peak counts and length pattern.", BLUE),
        ("2", "Read a mixed spot", "Count cut sites in each peak and sort them into three length bins.", GOLD),
        ("3", "Try mixtures", "Combine reference fingerprints in different proportions and score the fit.", PURPLE),
        ("4", "Report proportions", "Choose the most supported mixture and normalize it to percentages.", TEAL),
    ]
    for i, (num, title, body, color) in enumerate(steps):
        x = 0.67 + i * 3.16
        add_card(slide, x, 1.50, 2.82, 4.48, fill=NAVY_2, line=color)
        add_circle(slide, x + 0.20, 1.73, 0.58, color, line=color, line_width=1)
        add_text(slide, num, x + 0.20, 1.75, 0.58, 0.46, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        # icons
        if i == 0:
            for j, c in enumerate([BLUE, TEAL, CORAL]):
                add_circle(slide, x + 0.42 + j * 0.65, 2.64, 0.46, c, line=WHITE, line_width=1)
                for k in range(3):
                    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.40 + j * 0.65 + k * 0.10, 3.25 - k * 0.10, 0.07, 0.26 + k * 0.10, c, line=c)
        elif i == 1:
            for j, (ww, c) in enumerate([(0.36, TEAL), (0.72, GOLD), (1.05, CORAL)]):
                add_fragment(slide, x + 0.45 + j * 0.63, 3.02 + (j % 2) * 0.35, ww, c, line_width=4, cut_color=WHITE)
        elif i == 2:
            add_stacked_bar(slide, x + 0.45, 2.84, 1.88, 0.42, [4, 3, 2], [BLUE, TEAL, CORAL], outline=WHITE)
            add_text(slide, "↻", x + 1.04, 3.39, 0.70, 0.55, size=27, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
        else:
            add_stacked_bar(slide, x + 0.45, 2.84, 1.88, 0.54, [5, 3, 2], [BLUE, TEAL, CORAL], ["50", "30", "20"], outline=WHITE)
        add_text(slide, title, x + 0.25, 4.02, 2.32, 0.42, size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.28, 4.58, 2.26, 1.13, size=14, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)
        if i < 3:
            add_shape(slide, MSO_SHAPE.CHEVRON, x + 2.84, 3.45, 0.34, 0.58, color, line=color)
    add_pill(slide, "Under the hood: count fit + shape fit + a sensible abundance prior", 3.48, 6.36, 6.38, 0.42, TEAL, size=12)
    add_footer(slide, "Technical version: negative-binomial peak totals + conditional multinomial length bins; fitted by MAP.", dark=True)
    add_notes(slide, "Keep the math verbal unless someone asks. ShapeMix learns reference fingerprints, reads a mixed spot, tries possible combinations, and chooses the mixture that best explains both the total counts and the length composition. MAP means choosing the most supported values after combining the data with a reasonable prior.")

    # 14 — Fair test
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_title(slide, "A matched experiment isolates the value of fragment length", "Experimental design", slide_num=14)
    add_text(slide, "The two models see the same data and use the same settings. Only one ingredient changes.", 0.72, 1.16, 11.85, 0.42, size=19, color=SLATE, align=PP_ALIGN.CENTER)
    # split
    add_card(slide, 0.72, 1.77, 3.15, 4.56, fill=PALE, line=MID)
    add_text(slide, "1  Split labeled cells", 1.00, 2.03, 2.60, 0.34, size=19, color=NAVY, bold=True)
    add_stacked_bar(slide, 1.04, 2.69, 2.50, 0.52, [70, 30], [BLUE, GOLD], ["70%", "30%"])
    add_text(slide, "Reference pool", 1.03, 3.38, 1.30, 0.24, size=13, color=BLUE, bold=True)
    add_text(slide, "Held-out test pool", 2.17, 3.38, 1.42, 0.24, size=13, color=SLATE, bold=True, align=PP_ALIGN.RIGHT)
    add_bullets(slide, ["Learn fingerprints from 70%", "Make pseudo-spots from the other 30%", "Never reuse the same cell on both sides"], 1.00, 4.05, 2.55, 1.66, size=14.5, bullet_color=BLUE)
    # arms
    add_card(slide, 4.18, 1.77, 3.75, 4.56, fill=BLUE_PALE, line=BLUE)
    add_text(slide, "2A  Count-only", 4.49, 2.03, 3.13, 0.34, size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Uses peak totals", 4.48, 2.55, 3.15, 0.34, size=17, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    for i, height in enumerate([0.58, 1.10, 0.78, 1.38, 0.92]):
        add_shape(slide, MSO_SHAPE.RECTANGLE, 4.80 + i * 0.50, 4.55 - height, 0.30, height, BLUE, line=BLUE)
    add_text(slide, "Nₛₚ", 5.48, 4.88, 1.14, 0.42, size=25, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_pill(slide, "BASELINE", 5.31, 5.56, 1.46, 0.36, BLUE, size=10)
    add_card(slide, 8.24, 1.77, 4.38, 4.56, fill=TEAL_PALE, line=TEAL)
    add_text(slide, "2B  ShapeMix", 8.59, 2.03, 3.68, 0.34, size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Same totals + length composition", 8.56, 2.55, 3.74, 0.34, size=17, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_stacked_bar(slide, 8.96, 3.33, 2.90, 0.54, [5, 3, 2], [TEAL, GOLD, CORAL], outline=NAVY)
    add_text(slide, "Yₛₚᵦ", 9.87, 4.18, 1.12, 0.42, size=25, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_pill(slide, "ONLY EXTRA TERM: SHAPE", 9.39, 5.56, 2.18, 0.36, TEAL, size=9.5)
    add_pill(slide, "Same cells • peaks • seeds • prior • optimizer • compute budget", 3.11, 6.56, 7.14, 0.37, NAVY, size=11)
    add_footer(slide, "Design source: ShapeMix benchmark protocol and project proposal in this repository.")
    add_notes(slide, "This is the strongest part of the experimental design. If the models differ, it should be because ShapeMix saw fragment-length composition—not because it got different cells, peaks, or extra computing time.")

    # 15 — Main dataset
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    add_title(slide, "Primary dataset: human blood immune cells", "Dataset plan", slide_num=15)
    add_text(slide, "A public 10x Genomics PBMC Multiome dataset provides paired RNA labels and raw ATAC fragments.", 0.72, 1.14, 11.85, 0.44, size=19, color=SLATE, align=PP_ALIGN.CENTER)
    facts = [
        ("PBMCs", "blood immune cells", BLUE),
        ("1 donor", "conditional first test", CORAL),
        ("16 types", "including rare types", PURPLE),
        ("5,000 peaks", "selected from reference cells", TEAL),
        ("1,024 spots", "32 × 32 simulated grid", GOLD),
        ("~10 cells", "average per pseudo-spot", GREEN),
    ]
    for i, (big, small, color) in enumerate(facts):
        col = i % 3
        row = i // 3
        x = 0.78 + col * 4.16
        y = 1.83 + row * 1.53
        add_card(slide, x, y, 3.78, 1.22, fill=WHITE, line=color)
        add_circle(slide, x + 0.22, y + 0.30, 0.56, color, line=color, line_width=1)
        add_text(slide, big, x + 0.97, y + 0.19, 2.50, 0.40, size=23, color=NAVY, bold=True)
        add_text(slide, small, x + 0.97, y + 0.66, 2.50, 0.28, size=13, color=SLATE)
    add_card(slide, 0.79, 5.12, 11.75, 1.40, fill=NAVY, line=NAVY)
    add_text(slide, "Why simulate spatial spots?", 1.11, 5.43, 3.04, 0.36, size=21, color=WHITE, bold=True)
    add_text(slide, "Because the source cells are labeled, every pseudo-spot has an exact ingredient list. That gives a fair answer key.", 4.27, 5.38, 7.80, 0.68, size=18, color=LIGHT_TEXT)
    add_pill(slide, "KNOWN GROUND TRUTH", 9.87, 6.02, 1.84, 0.34, TEAL, size=9.5)
    add_footer(slide, "Source: 10x Genomics PBMC from a healthy donor, granulocytes removed, 10k, Cell Ranger ARC 2.0.0.")
    add_notes(slide, "PBMC means peripheral blood mononuclear cell. This dataset contains many related immune types, so it is a useful challenge. Simulated spots are essential because we know exactly which cells were mixed together. The one-donor limitation will be stated clearly.")

    # 16 — Additional datasets
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_title(slide, "Broader datasets test whether the idea generalizes", "Validation roadmap", slide_num=16)
    add_text(slide, "Different datasets answer different questions; not all provide exact ground truth.", 0.72, 1.13, 11.84, 0.42, size=19, color=SLATE, align=PP_ALIGN.CENTER)
    headers = [("DATASET", 0.82, 2.16), ("WHAT IT ADDS", 3.02, 5.10), ("HOW WE CHECK IT", 8.21, 4.23)]
    for text_value, x, w in headers:
        add_pill(slide, text_value, x, 1.73, w, 0.37, NAVY, size=10)
    rows = [
        ("Main PBMC", "Matched reference/test cells", "Exact pseudo-spot recipe", BLUE_PALE, BLUE),
        ("GSE129785", "Immune dilution series + sorted cells", "Nominal ratios or new pseudo-spots", GOLD_PALE, GOLD),
        ("GSE194122", "Multiple bone-marrow donors and sites", "Hold out one donor at a time", PURPLE_PALE, PURPLE),
        ("GSE205055 / 263333", "Real spatial tissues and multi-omics", "Anatomy, RNA, protein, or markers", TEAL_PALE, TEAL),
    ]
    for i, (name, adds, check, fill, color) in enumerate(rows):
        y = 2.29 + i * 0.91
        add_card(slide, 0.78, y, 11.78, 0.72, fill=fill, line=color)
        add_text(slide, name, 1.00, y + 0.16, 1.82, 0.27, size=15.5, color=NAVY, bold=True)
        add_text(slide, adds, 3.13, y + 0.13, 4.72, 0.34, size=14.5, color=INK)
        add_text(slide, check, 8.31, y + 0.13, 3.90, 0.34, size=14.5, color=INK)
    add_card(slide, 0.80, 6.05, 11.74, 0.64, fill=PALE, line=MID)
    add_rich_text(
        slide,
        [("Key rule: ", {"bold": True, "color": CORAL}), ("real spatial tissue usually has no exact answer key, so supporting evidence is validation—not perfect truth.", {"color": NAVY})],
        1.04,
        6.19,
        11.22,
        0.32,
        size=15,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "Sources: NCBI GEO accessions GSE129785, GSE194122, GSE205055, and GSE263333.")
    add_notes(slide, "Stress that each dataset has a different role. Simulations can provide exact composition truth. Physical dilutions provide planned mixing ratios. Real tissues need indirect checks using anatomy or other molecular measurements.")

    # 17 — Evaluation
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    add_title(slide, "Success means estimates are close to the known recipe", "Evaluation", slide_num=17)
    add_pill(slide, "TEACHING EXAMPLE — NOT RESULTS", 4.81, 1.20, 3.70, 0.39, CORAL, size=10.5)
    add_card(slide, 0.79, 1.88, 6.18, 4.72, fill=WHITE, line=MID)
    add_text(slide, "Compare truth with prediction", 1.11, 2.17, 5.50, 0.38, size=21, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    examples = [
        ("Truth", [50, 30, 20], ["50%", "30%", "20%"]),
        ("Model A", [46, 34, 20], ["46%", "34%", "20%"]),
        ("Model B", [28, 42, 30], ["28%", "42%", "30%"]),
    ]
    for i, (label, vals, labels2) in enumerate(examples):
        y = 2.91 + i * 0.94
        add_text(slide, label, 1.11, y + 0.08, 1.02, 0.28, size=15, color=NAVY, bold=True)
        add_stacked_bar(slide, 2.22, y, 3.92, 0.50, vals, [BLUE, TEAL, CORAL], labels2)
        if i == 1:
            add_pill(slide, "closer", 6.18, y + 0.05, 0.62, 0.34, GREEN, size=8.5)
        if i == 2:
            add_pill(slide, "farther", 6.18, y + 0.05, 0.62, 0.34, CORAL, size=8.5)
    for i, (name, color) in enumerate([("B cell", BLUE), ("T cell", TEAL), ("Monocyte", CORAL)]):
        add_circle(slide, 1.73 + i * 1.55, 5.85, 0.19, color, line=color, line_width=0.5)
        add_text(slide, name, 1.98 + i * 1.55, 5.79, 1.10, 0.25, size=11.5, color=SLATE)
    # metric cards
    metric_cards = [
        ("RMSE", "average size of the proportion errors", "LOWER is better", BLUE, BLUE_PALE),
        ("JSD", "how different the full mixtures are", "LOWER is better", PURPLE, PURPLE_PALE),
        ("Rare-cell F1", "balance of finding rare cells without false alarms", "HIGHER is better", TEAL, TEAL_PALE),
        ("Runtime", "extra computing cost for the added information", "LESS is better", GOLD, GOLD_PALE),
    ]
    for i, (name, desc, direction, color, fill) in enumerate(metric_cards):
        y = 1.88 + i * 1.18
        add_card(slide, 7.34, y, 5.20, 0.97, fill=fill, line=color)
        add_pill(slide, name, 7.61, y + 0.18, 1.18, 0.34, color, size=9.5)
        add_text(slide, desc, 9.02, y + 0.12, 3.08, 0.42, size=13.5, color=INK)
        add_text(slide, direction, 9.03, y + 0.59, 2.46, 0.22, size=10.5, color=color, bold=True)
    add_footer(slide, "Primary planned endpoints: RMSE and base-2 Jensen–Shannon divergence (JSD).")
    add_notes(slide, "The colored bars make the metrics intuitive: the model should reproduce the known recipe. RMSE looks at individual proportion errors; JSD compares the full mixture. Rare-cell F1 asks whether uncommon populations are detected without too many false alarms.")

    # 18 — Possible outcomes
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_title(slide, "Every possible outcome teaches us something", "Interpreting the study", dark=True, slide_num=18)
    outcomes = [
        ("SHAPE HELPS", "lower error", "Fragment length adds cell-type information beyond peak totals.", GREEN, GREEN_PALE, "↓"),
        ("ABOUT THE SAME", "similar error", "The extra clue may be too weak or redundant in this setting.", GOLD, GOLD_PALE, "≈"),
        ("SHAPE HURTS", "higher error", "Noise, sparsity, or reference mismatch may outweigh the added detail.", CORAL, CORAL_PALE, "↑"),
    ]
    for i, (title, sub, body, color, pale, symbol) in enumerate(outcomes):
        x = 0.73 + i * 4.17
        add_card(slide, x, 1.57, 3.78, 4.67, fill=NAVY_2, line=color)
        add_circle(slide, x + 1.40, 1.93, 0.98, color, line=color, line_width=1)
        add_text(slide, symbol, x + 1.40, 1.94, 0.98, 0.80, size=34, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, title, x + 0.33, 3.12, 3.10, 0.38, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_pill(slide, sub.upper(), x + 1.06, 3.70, 1.66, 0.34, color, size=9.5)
        add_text(slide, body, x + 0.42, 4.39, 2.94, 1.10, size=15.5, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, "A negative result is still a real result: it sets a boundary on when added complexity is useful.", 1.44, 6.58, 10.43, 0.34, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Outcome cards describe interpretations only; they do not report study results.", dark=True)
    add_notes(slide, "Science does not require the hypothesis to win. If ShapeMix ties or loses, that still tells us that this particular fragment representation did not justify its complexity under the tested conditions.")

    # 19 — Limitations and safeguards
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_title(slide, "Known risks—and how the design handles them", "Scientific caution", slide_num=19)
    pairs = [
        ("Sparse shape counts", "Use only three broad length bins and well-covered peaks.", GOLD),
        ("One main donor", "Limit claims; add multi-donor data before generalizing.", CORAL),
        ("Reference ≠ spatial assay", "Audit protocol differences and treat real tissue as qualitative validation.", PURPLE),
        ("More complexity can overfit", "Use held-out cells and an exactly matched count-only baseline.", TEAL),
    ]
    for i, (risk, safeguard, color) in enumerate(pairs):
        row = i // 2
        col = i % 2
        x = 0.78 + col * 6.08
        y = 1.53 + row * 2.46
        add_card(slide, x, y, 5.72, 2.06, fill=PALE, line=color)
        add_pill(slide, "RISK", x + 0.25, y + 0.25, 0.74, 0.33, color, size=9.5)
        add_text(slide, risk, x + 1.17, y + 0.18, 4.14, 0.42, size=19, color=NAVY, bold=True)
        add_line(slide, x + 0.27, y + 0.78, x + 5.43, y + 0.78, color=MID, width=0.8)
        add_pill(slide, "SAFEGUARD", x + 0.25, y + 1.03, 1.08, 0.32, NAVY, size=8.5)
        add_text(slide, safeguard, x + 1.50, y + 0.96, 3.82, 0.70, size=14.5, color=INK)
    add_pill(slide, "Claims will match the strength of the evidence", 4.30, 6.61, 4.73, 0.40, NAVY, size=11.5)
    add_footer(slide, "These limitations are part of the planned interpretation, not after-the-fact excuses.")
    add_notes(slide, "Explain that limitations do not invalidate a project; they define what can be concluded. The biggest limitation is that the first benchmark uses one donor, so it cannot establish broad biological generalization by itself.")

    # 20 — Results pending
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    add_title(slide, "This draft intentionally stops before the results", "Project status", slide_num=20)
    add_card(slide, 0.76, 1.47, 11.80, 1.08, fill=NAVY, line=NAVY)
    add_text(slide, "RESULTS PENDING", 1.05, 1.77, 2.70, 0.42, size=23, color=GOLD, bold=True)
    add_text(slide, "No performance claim is made in this slide deck.", 3.93, 1.78, 7.89, 0.38, size=19, color=WHITE, bold=True)
    timeline = [
        ("1", "Prepare data", "fragments, labels, splits", BLUE),
        ("2", "Build fingerprints", "counts + length profiles", PURPLE),
        ("3", "Run matched test", "count-only vs ShapeMix", TEAL),
        ("4", "Check controls", "fairness + diagnostics", GOLD),
        ("5", "Add result slides", "metrics + limitations", CORAL),
    ]
    add_line(slide, 1.42, 4.01, 11.85, 4.01, color=MID, width=5)
    for i, (num, title, sub, color) in enumerate(timeline):
        x = 1.08 + i * 2.42
        add_circle(slide, x, 3.63, 0.76, color, line=WHITE, line_width=2)
        add_text(slide, num, x, 3.65, 0.76, 0.60, size=21, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, title, x - 0.39, 4.65, 1.54, 0.34, size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, sub, x - 0.44, 5.10, 1.65, 0.52, size=12.5, color=SLATE, align=PP_ALIGN.CENTER)
    add_text(slide, "Planned update: replace this slide with observed results, uncertainty, controls, and an honest interpretation.", 1.28, 6.23, 10.78, 0.48, size=17, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Status language follows the requested results-free presentation scope.")
    add_notes(slide, "Be explicit with the audience: this is not a results talk. When results are ready, add the actual comparison, uncertainty, controls, and limitations rather than showing placeholder charts that could be mistaken for data.")

    # 21 — Take-home
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_pill(slide, "TAKE-HOME MESSAGE", 0.74, 0.56, 1.92, 0.38, TEAL, size=10.5)
    add_text(slide, "ShapeMix asks a simple, testable question", 0.74, 1.20, 8.62, 0.64, size=33, color=WHITE, bold=True, font=FONT_DISPLAY)
    add_text(slide, "Does fragment-length composition help identify the cell types inside a mixed spatial ATAC spot?", 0.76, 2.05, 8.32, 1.12, size=25, color=LIGHT_TEXT, font=FONT_DISPLAY)
    takeaways = [
        ("01", "ATAC fragments contain more than a peak count."),
        ("02", "ShapeMix adds length composition to a reference-guided mixture model."),
        ("03", "A matched, held-out test will decide whether that extra clue earns its complexity."),
    ]
    for i, (num, text_value) in enumerate(takeaways):
        y = 3.53 + i * 0.82
        add_pill(slide, num, 0.78, y, 0.62, 0.39, [BLUE, GOLD, TEAL][i], size=10)
        add_text(slide, text_value, 1.67, y - 0.03, 7.95, 0.50, size=17.5, color=WHITE, bold=True)
    # question mark graphic
    add_circle(slide, 10.10, 1.25, 2.40, NAVY_2, line=TEAL, line_width=3)
    add_text(slide, "?", 10.10, 1.27, 2.40, 1.92, size=70, color=GOLD, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Questions?", 9.86, 4.12, 2.90, 0.54, size=27, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "shape • mixture • evidence", 9.86, 4.73, 2.90, 0.28, size=12.5, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)
    add_footer(slide, "ShapeMix-ATAC project overview • Andy Zhuang • High School Science Research", dark=True)
    add_notes(slide, "Close by returning to the one-sentence question. The project is not trying to prove that a more complicated model must win; it is testing whether a specific ATAC-native clue adds measurable value.")

    # 22 — References and glossary
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_title(slide, "References and biological glossary", "Appendix", slide_num=22)
    add_text(slide, "Selected references", 0.77, 1.24, 5.85, 0.38, size=21, color=NAVY, bold=True)
    refs = [
        "Buenrostro JD et al. (2013). ATAC-seq and open chromatin. Nature Methods.",
        "Deng Y et al. (2022). Spatial chromatin-accessibility profiling. Nature.",
        "Ouologuem S et al. (2025). RNA deconvolution methods transferred to spatial ATAC. Bioinformatics.",
        "Tsoucas et al. (2019), DWLS; Dong & Yuan (2021), SpatialDWLS.",
        "Cable et al. (2022), RCTD; Kleshchevnikov et al. (2022), cell2location.",
        "Biancalani et al. (2021), Tangram; Lopez et al. (2022), DestVI.",
        "ShapeMix-ATAC model specification v1 and benchmark protocol in this repository.",
        "Data: 10x PBMC Multiome; GEO GSE129785, GSE194122, GSE205055, GSE263333.",
    ]
    add_bullets(slide, refs, 0.81, 1.78, 5.74, 4.78, size=10.8, bullet_color=TEAL)
    add_text(slide, "Biology / study glossary", 6.99, 1.24, 5.55, 0.38, size=21, color=NAVY, bold=True)
    glossary = [
        ("ATAC-seq", "measures accessible DNA"),
        ("Peak", "a genomic region with many ATAC cuts"),
        ("Fragment", "DNA piece between two Tn5 cut sites"),
        ("Deconvolution", "estimating ingredients of a mixture"),
        ("Reference", "fingerprint learned from labeled cells"),
        ("Pseudo-spot", "simulated spot with a known recipe"),
        ("MAP", "most supported parameter values under model + prior"),
    ]
    for i, (term, meaning) in enumerate(glossary):
        y = 1.78 + i * 0.67
        add_pill(slide, term, 7.03, y, 1.44, 0.34, [BLUE, TEAL, GOLD, CORAL, PURPLE, GREEN, NAVY][i], size=9)
        add_text(slide, meaning, 8.72, y - 0.01, 3.67, 0.38, size=13.2, color=INK)
    add_card(slide, 6.96, 6.57, 5.59, 0.40, fill=PALE, line=MID)
    add_text(slide, "Full project sources are documented in docs/ShapeMix and docs/research_class.", 7.12, 6.65, 5.25, 0.20, size=10.5, color=SLATE, align=PP_ALIGN.CENTER)
    add_footer(slide, "Links and full citations appear in the repository proposal and ShapeMix documentation.")
    add_notes(slide, "Use this slide if the audience asks about biological terms, datasets, or method sources. The algorithm papers are grouped compactly; full links and citations are in the repository proposal.")

    # 23 — Bayesian glossary
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    add_title(slide, "Bayesian glossary: translating the ShapeMix model", "Appendix", slide_num=23)
    add_text(slide, "posterior  ∝  likelihood  ×  prior", 3.35, 1.17, 6.63, 0.48, size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    glossary_bayes = [
        ("LATENT VARIABLE", "A quantity we cannot directly observe. Here, z is the hidden effective amount of each cell type.", BLUE, BLUE_PALE),
        ("PRIOR", "A probability rule used before this spot's data. It keeps z positive and regularizes weak evidence.", GOLD, GOLD_PALE),
        ("LIKELIHOOD", "A score for how probable the observed counts would be if a proposed mixture were true.", PURPLE, PURPLE_PALE),
        ("POSTERIOR", "The updated support for possible mixtures after likelihood and prior are combined.", TEAL, TEAL_PALE),
        ("NEGATIVE BINOMIAL", "A count distribution that allows more variability than a simple Poisson model.", BLUE, BLUE_PALE),
        ("MULTINOMIAL", "A distribution for splitting a fixed total among categories—here, three length bins.", CORAL, CORAL_PALE),
        ("MAP", "Maximum a posteriori: the single mixture at the highest point of the posterior density.", GREEN, GREEN_PALE),
        ("OPTIMIZER", "A numerical search procedure. Adam repeatedly changes z to raise the posterior score.", NAVY, PALE),
    ]
    for i, (term, meaning, color, fill) in enumerate(glossary_bayes):
        col = i % 2
        row = i // 2
        x = 0.78 + col * 6.06
        y = 1.86 + row * 1.18
        add_card(slide, x, y, 5.69, 0.94, fill=fill, line=color)
        add_pill(slide, term, x + 0.20, y + 0.24, 1.66, 0.34, color, color=NAVY if color == GOLD else WHITE, size=8.2)
        add_text(slide, meaning, x + 2.05, y + 0.11, 3.38, 0.68, size=11.8, color=INK)
    add_card(slide, 2.21, 6.55, 8.91, 0.48, fill=NAVY, line=NAVY)
    add_text(slide, "ShapeMix uses a Bayesian model with MAP inference; model v1 does not calculate credible intervals.", 2.43, 6.61, 8.47, 0.30, size=11.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Plain-language definitions are specific to how these terms are used in ShapeMix-ATAC.")
    add_notes(slide, "This appendix slide separates the model from the inference method. ShapeMix is Bayesian because it combines likelihoods with a prior. Its inference method is MAP optimization, which reports one best-supported point rather than sampling the full posterior. Negative binomial models total-count noise; multinomial models how a fixed total is divided among the three parent-fragment length bins.")

    # Core metadata
    prs.core_properties.title = "ShapeMix-ATAC: High School Science Research Deck"
    prs.core_properties.subject = "Results-free introduction, existing algorithms, Bayesian ShapeMix model, datasets, and evaluation plan"
    prs.core_properties.author = "Andy Zhuang"
    prs.core_properties.keywords = "ShapeMix, ATAC-seq, spatial deconvolution, Bayesian model, MAP inference, high school research"
    prs.core_properties.comments = "Draft deck generated from the deconvATAC repository; no study results included."

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build_deck()
