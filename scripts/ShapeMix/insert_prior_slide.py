#!/usr/bin/env python3
"""Insert a Gamma-prior slide immediately after slide 7 (the Bayes slide).

Content follows Section 1.4 of
docs/ShapeMix/tutorials/ShapeMix_ATAC_Bayesian_simple.pdf:
z ~ Gamma(shape 2, rate 1), the frozen version-1 prior (mean 2, variance 2,
mode 1), with a drawn density curve and how the prior is determined.

Existing slides are not modified; the new slide is added and moved into
position 8.
"""

from __future__ import annotations

import math
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from insert_count_only_bayesian_slides import (
    CREAM,
    FONT_DISPLAY,
    INK,
    KICKER,
    MID,
    NAVY,
    PURPLE,
    PURPLE_PALE,
    SLATE,
    TEAL,
    WHITE,
    C,
    add_bullets,
    add_card,
    add_notes,
    add_pill,
    add_text,
    set_bg,
)
from update_slide7_bayes_simplify import SLIDE_TITLE as BAYES_TITLE
from update_slide9_bayes_compact import add_line, slide_texts

REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

PRIOR_TITLE = "The prior: what is plausible before seeing any data"


def draw_density(slide, x0, y0, w, h):
    """Draw the Gamma(2,1) density x*exp(-x) on [0, 8] inside the given box."""
    x_max, y_max = 8.0, 0.42
    left_pad, bottom_pad = 0.35, 0.42  # room for axis labels
    plot_x, plot_w = x0 + left_pad, w - left_pad - 0.1
    plot_top, plot_h = y0 + 0.25, h - 0.25 - bottom_pad
    plot_bottom = plot_top + plot_h

    def px(x):
        return plot_x + (x / x_max) * plot_w

    def py(y):
        return plot_bottom - (y / y_max) * plot_h

    # Axes.
    add_line(slide, plot_x, plot_top - 0.05, plot_x, plot_bottom, SLATE, width=1.2)
    add_line(slide, plot_x, plot_bottom, plot_x + plot_w, plot_bottom, SLATE, width=1.2)

    # Density curve as an open freeform.
    points = []
    steps = 64
    for i in range(steps + 1):
        x = 0.02 + (x_max - 0.02) * i / steps
        y = x * math.exp(-x)
        points.append((Inches(px(x)), Inches(py(y))))
    builder = slide.shapes.build_freeform(
        Emu(int(points[0][0])), Emu(int(points[0][1])), scale=1.0
    )
    builder.add_line_segments(
        [(Emu(int(px_)), Emu(int(py_))) for px_, py_ in points[1:]], close=False
    )
    curve = builder.convert_to_shape()
    curve.fill.background()
    curve.line.color.rgb = C(PURPLE)
    curve.line.width = Pt(2.75)
    curve.shadow.inherit = False

    # Dashed marker at the mode (x = 1, density e^-1).
    mode_top = py(math.exp(-1))
    mode_line = add_line(slide, px(1.0), mode_top, px(1.0), plot_bottom, TEAL, width=1.5)
    mode_line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    add_text(slide, "mode = 1", px(1.0) + 0.08, mode_top - 0.28, 1.2, 0.24,
             size=11, color=TEAL, bold=True)
    add_text(slide, "positive values only", px(2.6), py(0.22), 2.2, 0.26,
             size=10.5, color=SLATE)

    # Tick labels and axis titles.
    for tick in (0, 2, 4, 6, 8):
        add_text(slide, str(tick), px(tick) - 0.12, plot_bottom + 0.03, 0.24, 0.2,
                 size=9, color=SLATE, align=PP_ALIGN.CENTER)
    add_text(slide, "1", px(1.0) - 0.12, plot_bottom + 0.03, 0.24, 0.2,
             size=9, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "possible amount z", plot_x + plot_w / 2 - 1.0,
             plot_bottom + 0.2, 2.0, 0.22, size=10.5, color=SLATE,
             align=PP_ALIGN.CENTER)
    add_text(slide, "prior density", x0, y0, 1.6, 0.22, size=10.5, color=SLATE)


def build(slide) -> None:
    set_bg(slide, CREAM)
    add_text(slide, KICKER, 0.68, 0.28, 8.8, 0.28, size=10.5, color=TEAL, bold=True)
    add_text(
        slide, PRIOR_TITLE, 0.65, 0.55, 11.8, 0.62,
        size=24, color=NAVY, bold=True, font=FONT_DISPLAY,
    )

    # Left card: the prior itself.
    add_card(slide, 0.7, 1.45, 4.75, 3.65)
    add_pill(slide, "PRIOR  p(z)", 0.95, 1.65, 1.55, 0.32, PURPLE)
    add_text(slide, "z ~ Gamma(shape = 2, rate = 1)",
             0.95, 2.15, 4.3, 0.42, size=19, color=NAVY, bold=True)
    add_text(slide, "the frozen version-1 prior, applied to every element of z",
             0.95, 2.62, 4.3, 0.3, size=11.5, color=SLATE)
    facts = [("mean", "2"), ("variance", "2"), ("mode", "1")]
    for i, (label, value) in enumerate(facts):
        fx = 0.95 + i * 1.42
        add_card(slide, fx, 3.05, 1.28, 0.78, fill=PURPLE_PALE, line=PURPLE)
        add_text(slide, value, fx, 3.12, 1.28, 0.34, size=17, color=NAVY,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, fx, 3.5, 1.28, 0.24, size=10, color=SLATE,
                 align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "Only positive values are allowed — a cell amount cannot be negative.",
        0.95, 4.05, 4.3, 0.55, size=11.5, color=INK,
    )

    # Right: the density plot.
    add_card(slide, 5.75, 1.45, 6.9, 3.65, fill=WHITE, line=MID)
    draw_density(slide, 6.0, 1.6, 6.4, 3.35)

    # Bottom: how the prior is determined.
    add_text(slide, "How is the prior determined?", 0.71, 5.3, 6.0, 0.3,
             size=14, color=NAVY, bold=True)
    add_bullets(
        slide,
        [
            "Chosen ahead of time and frozen: the Gamma(2, 1) prior is fixed before "
            "any spot's data is analyzed (version-1 protocol).",
            "Fair across cell types: every type gets the identical prior, so none "
            "starts with an advantage.",
            "An assumption, not a measurement: it gently prefers moderate positive "
            "amounts, and other reasonable priors can be tested in sensitivity checks.",
        ],
        0.71, 5.66, 11.9, 1.4, size=12, bullet_color=PURPLE,
    )

    hairline = add_card(slide, 0.68, 7.13, 11.97, 0.011, fill=MID, line=MID)
    hairline.line.width = Pt(0.25)

    add_notes(
        slide,
        "The prior is the third term of Bayes' rule from the previous slide. Every "
        "element of z gets the same Gamma(2,1) prior: mean 2, variance 2, mode 1, "
        "positive values only. It is determined by protocol, not by data: chosen and "
        "frozen before any spot is analyzed, identical for every cell type, and "
        "revisitable in sensitivity checks. Source: ShapeMix_ATAC_Bayesian_simple.pdf, "
        "Section 1.4.",
    )


def main() -> None:
    presentation = Presentation(DECK)
    for slide in presentation.slides:
        if any(PRIOR_TITLE in text for text in slide_texts(slide)):
            raise RuntimeError("The prior slide already exists; refusing to duplicate")

    bayes_positions = [
        index
        for index, slide in enumerate(presentation.slides)
        if any(BAYES_TITLE in text for text in slide_texts(slide))
    ]
    if bayes_positions != [6]:
        raise RuntimeError(f"Expected the Bayes slide only at position 7; found {bayes_positions}")

    layout = presentation.slides[6].slide_layout
    new_slide = presentation.slides.add_slide(layout)
    build(new_slide)

    slide_id_list = presentation.slides._sldIdLst
    new_id = slide_id_list[-1]
    slide_id_list.remove(new_id)
    slide_id_list.insert(7, new_id)

    presentation.save(DECK)

    check = Presentation(DECK)
    if len(check.slides) != 29:
        raise RuntimeError(f"Expected 29 slides; found {len(check.slides)}")
    texts = " ".join(t for t in slide_texts(check.slides[7]) if t.strip())
    for needle in (PRIOR_TITLE, "Gamma(shape = 2, rate = 1)", "mode = 1",
                   "How is the prior determined?"):
        if needle not in texts:
            raise RuntimeError(f"New slide 8 is missing {needle!r}")
    before = slide_texts(check.slides[6])
    after = slide_texts(check.slides[8])
    if not any(BAYES_TITLE in t for t in before):
        raise RuntimeError("Slide 7 is no longer the Bayes slide")
    if not any("Four symbols" in t for t in after):
        raise RuntimeError("Slide 9 is no longer the symbols slide")
    print(f"Inserted the Gamma-prior slide at position 8: {DECK}")


if __name__ == "__main__":
    main()
