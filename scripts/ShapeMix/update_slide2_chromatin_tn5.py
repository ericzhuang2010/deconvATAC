#!/usr/bin/env python3
"""Redraw slide 2's CLOSED/OPEN chromatin panels in slide 3's visual style.

Slide 3 draws chromatin as a dark DNA thread wrapped around purple histone
cores, with Tn5 shown as a small two-lobed light-blue enzyme that cuts only
exposed linker DNA. Slide 2 previously used a different, older vocabulary
(plain purple circles, a diagonal blue line with gold beads, a coral Tn5
pill). This patch replaces only the two panel diagrams on slide 2 so both
slides share one visual language. No other slide is touched.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

EMU_PER_INCH = 914400

# Shared palette (identical to build_shapemix_high_school_deck.py / slide 3).
NAVY = "13233B"
SLATE = "52627A"
TEAL_DARK = "0C7F75"
BLUE = "3B82F6"
PURPLE = "8B5CF6"
CORAL = "F26B5B"
WHITE = "FFFFFF"

DNA_COLOR = "654640"
LOBE_COLORS = ("6D43D6", PURPLE, "7650DE", PURPLE)
LOBE_LINE = "5634B5"
ENZYME_BLUE = "68B5E8"

FONT = "Aptos"

DNA_Y = 3.15
NUCLEOSOME_SCALE = 0.92


def C(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#").upper())


def inches(emu: int) -> float:
    return emu / EMU_PER_INCH


def add_shape(slide, kind, x, y, w, h, fill, line=None, line_width=1.0):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(fill)
    shape.line.color.rgb = C(line if line else fill)
    shape.line.width = Pt(line_width)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=13,
    color=NAVY,
    bold=False,
    align=PP_ALIGN.CENTER,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = C(color)
    return box


def add_line(slide, x1, y1, x2, y2, color=DNA_COLOR, width=1.5):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = C(color)
    line.line.width = Pt(width)
    return line


def add_circle(slide, x, y, d, fill, line=WHITE, line_width=1.5):
    return add_shape(slide, MSO_SHAPE.OVAL, x, y, d, d, fill, line=line, line_width=line_width)


def add_nucleosome(slide, center_x, center_y, scale=1.0):
    """Histone core with a visible DNA wrap — copied from the slide-3 builder."""
    lobe_w = 0.18 * scale
    lobe_h = 0.48 * scale
    for offset, color in zip((-0.22, -0.075, 0.075, 0.22), LOBE_COLORS):
        lobe = add_shape(
            slide,
            MSO_SHAPE.OVAL,
            center_x + offset * scale - lobe_w / 2,
            center_y - lobe_h / 2,
            lobe_w,
            lobe_h,
            color,
            line=LOBE_LINE,
            line_width=0.6,
        )
        lobe.rotation = 18

    wrap = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(center_x - 0.38 * scale),
        Inches(center_y - 0.31 * scale),
        Inches(0.76 * scale),
        Inches(0.62 * scale),
    )
    wrap.fill.background()
    wrap.line.color.rgb = C(DNA_COLOR)
    wrap.line.width = Pt(1.45)
    wrap.rotation = 9
    add_line(
        slide,
        center_x - 0.34 * scale,
        center_y + 0.12 * scale,
        center_x + 0.34 * scale,
        center_y - 0.10 * scale,
        color=DNA_COLOR,
        width=1.25,
    )
    return wrap


def add_tn5_marker(slide, x, y, with_tick=True):
    """Two-lobed Tn5 enzyme — copied from the slide-3 builder."""
    add_circle(slide, x - 0.11, y - 0.14, 0.16, ENZYME_BLUE, line=WHITE, line_width=0.6)
    add_circle(slide, x - 0.01, y - 0.10, 0.16, ENZYME_BLUE, line=WHITE, line_width=0.6)
    if with_tick:
        add_line(slide, x + 0.02, y + 0.01, x + 0.02, y + 0.15, color=CORAL, width=1.2)


def add_blocked_mark(slide, cx, cy, arm=0.07):
    """Small coral X meaning Tn5 cannot bind here."""
    add_line(slide, cx - arm, cy - arm, cx + arm, cy + arm, color=CORAL, width=1.6)
    add_line(slide, cx - arm, cy + arm, cx + arm, cy - arm, color=CORAL, width=1.6)


def find_slide2(presentation):
    slide = presentation.slides[1]
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    joined = " ".join(texts)
    if "Chromatin controls which DNA instructions" not in joined:
        raise RuntimeError("Slide 2 is not the expected chromatin background slide")
    if "CLOSED" not in joined or "OPEN" not in joined:
        raise RuntimeError("Slide 2 no longer has the CLOSED/OPEN panels")
    return slide


def panel_card(slide, expected_x, expected_y):
    for shape in slide.shapes:
        if (
            abs(inches(shape.left) - expected_x) < 0.05
            and abs(inches(shape.top) - expected_y) < 0.05
            and inches(shape.width) > 3.0
            and inches(shape.height) > 3.0
        ):
            return shape
    raise RuntimeError(f"Panel card at ({expected_x}, {expected_y}) not found on slide 2")


def remove_old_diagrams(slide, keep_ids):
    """Delete every shape drawn inside the two panels except the cards."""
    doomed = []
    for shape in slide.shapes:
        if shape.shape_id in keep_ids:
            continue
        x, y = inches(shape.left), inches(shape.top)
        bottom = y + inches(shape.height)
        if x >= 4.4 and y >= 1.8 and bottom <= 5.4:
            doomed.append(shape)
    if len(doomed) != 17:
        names = [(s.shape_id, s.name) for s in doomed]
        raise RuntimeError(
            f"Expected 17 old diagram shapes inside the panels; found {len(doomed)}: {names}"
        )
    for shape in doomed:
        shape._element.getparent().remove(shape._element)


def draw_closed_panel(slide):
    # Tightly packed nucleosomes: the DNA thread has no exposed linker.
    add_line(slide, 4.88, DNA_Y, 7.72, DNA_Y, color=DNA_COLOR, width=1.7)
    for center_x in (5.25, 5.95, 6.65, 7.35):
        add_nucleosome(slide, center_x, DNA_Y, scale=NUCLEOSOME_SCALE)

    # Tn5 floats nearby but has nowhere to land.
    for marker_x, marker_y in ((5.62, 2.42), (6.68, 2.35)):
        add_tn5_marker(slide, marker_x, marker_y, with_tick=False)
        add_blocked_mark(slide, marker_x + 0.02, marker_y + 0.24)

    add_text(
        slide,
        "DNA wrapped around histones",
        4.70,
        3.62,
        3.00,
        0.20,
        size=9.5,
        color=PURPLE,
        bold=True,
    )
    add_text(
        slide,
        "No exposed linker → Tn5 is blocked",
        4.64,
        4.38,
        3.13,
        0.30,
        size=13,
        color=CORAL,
        bold=True,
    )


def draw_open_panel(slide):
    # Spread-out nucleosomes leave a long stretch of exposed linker DNA.
    add_line(slide, 9.08, DNA_Y, 11.94, DNA_Y, color=DNA_COLOR, width=1.7)
    for center_x in (9.42, 11.60):
        add_nucleosome(slide, center_x, DNA_Y, scale=NUCLEOSOME_SCALE)

    # Tn5 lands on the linker and cuts, exactly as drawn on slide 3.
    for marker_x in (10.12, 10.60, 11.08):
        add_tn5_marker(slide, marker_x, DNA_Y, with_tick=True)

    pill = add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        9.72,
        2.38,
        0.56,
        0.24,
        BLUE,
        line=BLUE,
        line_width=1,
    )
    tf = pill.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Tn5"
    r.font.name = FONT
    r.font.size = Pt(8.2)
    r.font.bold = True
    r.font.color.rgb = C(WHITE)

    add_text(
        slide,
        "exposed linker DNA",
        9.73,
        3.62,
        1.56,
        0.20,
        size=9.5,
        color=TEAL_DARK,
        bold=True,
    )
    add_text(
        slide,
        "Tn5 cuts and tags the exposed DNA",
        8.93,
        4.38,
        3.13,
        0.30,
        size=13,
        color=TEAL_DARK,
        bold=True,
    )


def update_notes(slide):
    try:
        notes = slide.notes_slide.notes_text_frame
        note_text = notes.text
        addition = (
            "The closed and open panels now use the same drawing vocabulary as slide 3: "
            "a dark DNA thread wraps around purple histone cores to form nucleosomes, and "
            "Tn5 is the small two-lobed blue enzyme. In closed chromatin the nucleosomes sit "
            "side by side, so Tn5 has no exposed linker DNA to land on (coral X). In open "
            "chromatin the nucleosomes are spread apart and Tn5 cuts the exposed linker."
        )
        if "same drawing vocabulary as slide 3" not in note_text:
            notes.text = f"{note_text.rstrip()}\n\n{addition}".strip()
    except Exception:
        pass


def main() -> None:
    presentation = Presentation(DECK)
    if len(presentation.slides) < 3:
        raise RuntimeError("The presentation has fewer than three slides")

    slide = find_slide2(presentation)
    closed_card = panel_card(slide, 4.54, 1.83)
    open_card = panel_card(slide, 8.83, 1.83)

    remove_old_diagrams(slide, keep_ids={closed_card.shape_id, open_card.shape_id})
    draw_closed_panel(slide)
    draw_open_panel(slide)
    update_notes(slide)

    temporary = DECK.with_name(f".{DECK.stem}.slide2-chromatin.tmp.pptx")
    presentation.save(temporary)
    check = Presentation(temporary)
    if len(check.slides) != len(presentation.slides):
        raise RuntimeError("Slide count changed while saving the patched deck")
    temporary.replace(DECK)
    print(f"Updated only slide 2: {DECK}")


if __name__ == "__main__":
    main()
