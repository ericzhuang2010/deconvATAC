#!/usr/bin/env python3
"""Insert three count-only Bayesian model slides immediately after slide 6.

The slides cover the Bayesian model *without* fragment length bins from
docs/ShapeMix/tutorials/ShapeMix_ATAC_Bayesian_simple.pdf (Section 1),
introducing the Section 1.3 symbols R, z, n, N.

The final package is assembled from the original ZIP so every pre-existing
slide and notes XML payload remains byte-for-byte unchanged.
"""

from __future__ import annotations

import re
from hashlib import sha256
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

KICKER = "BAYESIAN MODEL • WITHOUT LENGTH BINS"
INSERT_AFTER_INDEX = 5  # zero-based index of slide 6
NEW_TITLES = [
    "Four symbols describe one spot's counts",
    "Predict each peak's count, then allow noise",
    "Bayes' rule picks the best-supported amounts",
]

NAVY = "13233B"
INK = "213047"
SLATE = "52627A"
TEAL = "18A999"
TEAL_PALE = "DDF4F1"
BLUE = "3B82F6"
BLUE_PALE = "E7F0FE"
PURPLE = "8B5CF6"
PURPLE_PALE = "EEE8FF"
CORAL = "F26B5B"
CORAL_PALE = "FDE9E6"
GOLD = "F4B942"
GOLD_PALE = "FFF3D6"
CREAM = "F7F4EC"
WHITE = "FFFFFF"
MID = "D7E0E8"
FONT = "Aptos"
FONT_DISPLAY = "Aptos Display"


def C(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C(color)


def add_shape(slide, kind, x, y, w, h, fill, line=None, line_width=1.0):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(fill)
    shape.line.color.rgb = C(line if line else fill)
    shape.line.width = Pt(line_width)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=20,
    color=INK,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for side in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, side, Inches(0.04))
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = C(color)
    return box


def add_bullets(slide, items, x, y, w, h, size=12, color=INK, bullet_color=TEAL):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        p.line_spacing = 1.05
        p.alignment = PP_ALIGN.LEFT
        rb = p.add_run()
        rb.text = "•  "
        rb.font.name = FONT
        rb.font.size = Pt(size)
        rb.font.bold = True
        rb.font.color.rgb = C(bullet_color)
        rt = p.add_run()
        rt.text = item
        rt.font.name = FONT
        rt.font.size = Pt(size)
        rt.font.color.rgb = C(color)
    return box


def add_pill(slide, text, x, y, w, h, fill, color=WHITE, size=11):
    shp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill)
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = C(color)
    return shp


def add_card(slide, x, y, w, h, fill=WHITE, line=MID):
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line=line)


def add_title(slide, title):
    add_text(slide, KICKER, 0.68, 0.28, 8.8, 0.28, size=10.5, color=TEAL, bold=True)
    add_text(
        slide, title, 0.65, 0.55, 11.8, 0.62,
        size=27, color=NAVY, bold=True, font=FONT_DISPLAY,
    )


def add_footer(slide):
    line = add_shape(slide, MSO_SHAPE.RECTANGLE, 0.68, 7.13, 11.97, 0.011, MID)
    line.line.width = Pt(0.25)
    add_text(
        slide,
        "Model reference: docs/ShapeMix/tutorials/ShapeMix_ATAC_Bayesian_simple.pdf • Section 1",
        0.7, 7.17, 11.8, 0.18, size=7.5, color=SLATE,
    )


def add_takeaway(slide, text, y, x=0.7, w=11.93):
    bar = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.5, NAVY)
    tf = bar.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = C(WHITE)


def add_notes(slide, text):
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        pass


def build_symbols_slide(slide):
    set_bg(slide, CREAM)
    add_title(slide, NEW_TITLES[0])
    add_text(
        slide,
        "A fixed reference, a hidden amount, a calculated prediction, and one measurement.",
        0.71, 1.17, 11.9, 0.4, size=14, color=SLATE,
    )

    cards = [
        ("REFERENCE", BLUE, BLUE_PALE, "R[c,p]", "reference peak rate",
         "expected cut sites per unit of cell type c at peak p, learned from labeled reference cells",
         "array: C × P  •  fixed"),
        ("UNKNOWN", PURPLE, PURPLE_PALE, "z[s,c]", "effective abundance",
         "hidden positive amount of cell type c in spot s — the quantity the model infers",
         "array: S × C  •  inferred"),
        ("CALCULATED", TEAL, TEAL_PALE, "n[s,p]", "predicted (mean) counts",
         "the count a candidate recipe z predicts for spot s at peak p",
         "array: S × P  •  computed"),
        ("DATA", GOLD, GOLD_PALE, "N[s,p]", "observed total counts",
         "the cut-site count actually measured in spot s at peak p",
         "array: S × P  •  measured"),
    ]
    x, w, gap, y, h = 0.7, 2.88, 0.14, 1.72, 3.55
    for i, (role, color, pale, symbol, name, meaning, shape_line) in enumerate(cards):
        cx = x + i * (w + gap)
        add_card(slide, cx, y, w, h, fill=pale, line=color)
        add_pill(slide, role, cx + (w - 1.5) / 2, y + 0.2, 1.5, 0.3, color,
                 color=NAVY if color == GOLD else WHITE)
        add_text(slide, symbol, cx + 0.1, y + 0.66, w - 0.2, 0.5,
                 size=25, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, name, cx + 0.1, y + 1.24, w - 0.2, 0.3,
                 size=12.5, color=color if color != GOLD else "B8860B",
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, meaning, cx + 0.14, y + 1.62, w - 0.28, 1.35,
                 size=11, color=INK, align=PP_ALIGN.CENTER)
        add_text(slide, shape_line, cx + 0.1, y + 3.08, w - 0.2, 0.3,
                 size=10.5, color=SLATE, bold=True, align=PP_ALIGN.CENTER)

    add_takeaway(
        slide,
        "PREDICTED = HIDDEN AMOUNT × KNOWN RATE   •   OBSERVED = PREDICTED + NOISE",
        5.55,
    )
    add_text(
        slide,
        "Indices: s = spot, p = peak, c = cell type. ShapeMix later adds b = fragment-length bin.",
        0.71, 6.35, 11.9, 0.3, size=10.5, color=SLATE, align=PP_ALIGN.CENTER,
    )
    add_footer(slide)
    add_notes(
        slide,
        "These four arrays are the entire count-only model. R is fixed from labeled "
        "reference cells, z is the hidden per-spot mixture we infer, n is the model's "
        "predicted mean count, and N is the measured count. Color coding matches the "
        "tutorial: blue fixed, purple unknown, teal calculated, gold observed data.",
    )


def build_prediction_slide(slide):
    set_bg(slide, CREAM)
    add_title(slide, NEW_TITLES[1])
    add_text(
        slide,
        "Two steps: compute the expected count from the recipe, then treat the measurement "
        "as a noisy draw around it.",
        0.71, 1.17, 11.9, 0.4, size=14, color=SLATE,
    )

    # Panel 1: prediction.
    add_card(slide, 0.7, 1.72, 5.85, 4.1)
    add_pill(slide, "1 • PREDICT", 0.95, 1.92, 1.55, 0.32, TEAL)
    add_text(slide, "n[s,p]  =  Σ over c   z[s,c] · R[c,p]",
             0.95, 2.4, 5.35, 0.4, size=17, color=NAVY, bold=True)
    add_text(slide, "a dot product: each cell type's amount times its rate, added up",
             0.95, 2.84, 5.35, 0.3, size=11.5, color=SLATE)
    add_card(slide, 0.95, 3.3, 5.35, 1.5, fill=TEAL_PALE, line=TEAL)
    add_text(slide, "One spot, one peak, two cell types:",
             1.12, 3.42, 5.0, 0.28, size=11, color=SLATE)
    add_text(slide, "z_T = 1.5,  z_B = 0.5,  R_T = R_B = 12",
             1.12, 3.72, 5.0, 0.32, size=13.5, color=NAVY, bold=True)
    add_text(slide, "n = 1.5 × 12 + 0.5 × 12 = 24",
             1.12, 4.1, 5.0, 0.32, size=15, color=NAVY, bold=True)
    add_text(slide, "24 expected cut sites at this peak",
             1.12, 4.47, 5.0, 0.28, size=10.5, color=SLATE)
    add_text(slide, "Add each cell type's expected contribution.",
             0.95, 5.05, 5.35, 0.3, size=11, color=SLATE)

    # Panel 2: noise model.
    add_card(slide, 6.78, 1.72, 5.85, 4.1)
    add_pill(slide, "2 • ALLOW NOISE", 7.03, 1.92, 1.95, 0.32, CORAL)
    add_text(slide, "N[s,p]  ~  NegBinomial( mean = n[s,p],",
             7.03, 2.4, 5.45, 0.35, size=15.5, color=NAVY, bold=True)
    add_text(slide, "inverse-dispersion = φref · Σ over c  z[s,c] )",
             7.35, 2.78, 5.15, 0.35, size=15.5, color=NAVY, bold=True)
    add_bullets(
        slide,
        [
            "Real sequencing counts fluctuate: a prediction of 24 can appear as 21 or 27.",
            "The negative binomial allows more spread than a Poisson count with the same mean.",
            "φref is a fixed inverse-dispersion estimated once from training reference cells; "
            "bigger spots (larger Σ z) get steadier counts.",
        ],
        7.03, 3.35, 5.4, 2.2, size=11.5, bullet_color=CORAL,
    )

    add_takeaway(
        slide,
        "THIS COUNT-ONLY MODEL IS THE BASELINE ARM — SHAPEMIX ADDS LENGTH BINS ON TOP OF IT",
        6.2,
    )
    add_footer(slide)
    add_notes(
        slide,
        "Step one is deterministic: the predicted mean n is the dot product of amounts and "
        "rates. With the running example z_T = 1.5, z_B = 0.5 and both rates 12, the "
        "prediction is 24. Step two is probabilistic: the observed N is modeled as a "
        "negative-binomial draw centered on n, with spread controlled by the fixed "
        "reference inverse-dispersion times the spot's total abundance.",
    )


def build_bayes_slide(slide):
    set_bg(slide, CREAM)
    add_title(slide, NEW_TITLES[2])
    add_text(
        slide,
        "Every candidate recipe z receives one score; the reported answer is the "
        "highest-scoring recipe.",
        0.71, 1.17, 11.9, 0.4, size=14, color=SLATE,
    )

    cards = [
        (0.7, "POSTERIOR", TEAL, TEAL_PALE, "p(z | N)",
         "updated support for each candidate recipe after seeing the data",
         "report the MAP: the single best-supported z, normalized to percentages"),
        (4.66, "LIKELIHOOD", GOLD, GOLD_PALE, "p(N | z)",
         "negative-binomial score: how well the prediction n = z · R explains the observed N",
         "computed with the model from the previous slide"),
        (8.62, "PRIOR", PURPLE, PURPLE_PALE, "p(z)",
         "z ~ Gamma(2, 1): positive amounts only — mean 2, mode 1",
         "the same prior for every cell type, so none is favored"),
    ]
    y, w, h = 1.75, 3.5, 3.3
    for x, role, color, pale, formula, line1, line2 in cards:
        add_card(slide, x, y, w, h, fill=pale, line=color)
        add_pill(slide, role, x + (w - 1.6) / 2, y + 0.2, 1.6, 0.3, color,
                 color=NAVY if color == GOLD else WHITE)
        add_text(slide, formula, x + 0.1, y + 0.66, w - 0.2, 0.45,
                 size=22, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, line1, x + 0.16, y + 1.28, w - 0.32, 1.0,
                 size=11.5, color=INK, align=PP_ALIGN.CENTER)
        add_text(slide, line2, x + 0.16, y + 2.35, w - 0.32, 0.85,
                 size=10.5, color=SLATE, align=PP_ALIGN.CENTER)
    add_text(slide, "∝", 4.22, 3.0, 0.42, 0.6, size=28, color=NAVY,
             bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "×", 8.18, 3.0, 0.42, 0.6, size=28, color=NAVY,
             bold=True, align=PP_ALIGN.CENTER)

    add_takeaway(
        slide,
        "MAP = THE HIGHEST POINT OF THE POSTERIOR — ONE POINT ESTIMATE, NOT A CONFIDENCE INTERVAL",
        5.45,
    )
    add_text(
        slide,
        "ShapeMix keeps this exact machinery and adds one extra likelihood term for "
        "fragment-length bins (later sections).",
        0.71, 6.25, 11.9, 0.3, size=10.5, color=SLATE, align=PP_ALIGN.CENTER,
    )
    add_footer(slide)
    add_notes(
        slide,
        "Bayes' rule combines the two ingredients: the likelihood measures data fit and "
        "the Gamma(2,1) prior measures plausibility of the positive amounts. The model "
        "reports the MAP — the location of the posterior's highest point — normalized to "
        "cell-type percentages. It is a point estimate, not an interval.",
    )


BUILDERS = [build_symbols_slide, build_prediction_slide, build_bayes_slide]


def slide_texts(slide) -> list[str]:
    return [shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)]


def build_in_memory(presentation: Presentation) -> None:
    for slide in presentation.slides:
        if any(KICKER in text for text in slide_texts(slide)):
            raise RuntimeError("Count-only Bayesian slides already exist; refusing to duplicate")

    layout = presentation.slides[INSERT_AFTER_INDEX].slide_layout
    new_ids = []
    for builder in BUILDERS:
        slide = presentation.slides.add_slide(layout)
        builder(slide)
        new_ids.append(presentation.slides._sldIdLst[-1])

    slide_id_list = presentation.slides._sldIdLst
    for offset, sld_id in enumerate(new_ids):
        slide_id_list.remove(sld_id)
        slide_id_list.insert(INSERT_AFTER_INDEX + 1 + offset, sld_id)


def payload_hashes(path: Path, prefix: str) -> dict[str, str]:
    with ZipFile(path) as archive:
        return {
            name: sha256(archive.read(name)).hexdigest()
            for name in archive.namelist()
            if name.startswith(prefix) and name.endswith(".xml")
        }


NEW_PART_PATTERN = re.compile(
    r"^ppt/(slides|notesSlides)/(_rels/)?(slide|notesSlide)\d+\.xml(\.rels)?$"
)
REPLACED_PARTS = {
    "[Content_Types].xml",
    "ppt/presentation.xml",
    "ppt/_rels/presentation.xml.rels",
    "docProps/app.xml",
}


def assemble_without_rewriting_existing_slides(
    original_path: Path, generated_path: Path, output_path: Path
) -> None:
    with ZipFile(original_path, "r") as original, ZipFile(
        generated_path, "r"
    ) as generated, ZipFile(output_path, "w", compression=ZIP_DEFLATED) as output:
        original_names = set(original.namelist())
        generated_names = set(generated.namelist())

        new_parts = generated_names - original_names
        bad = [name for name in new_parts if not NEW_PART_PATTERN.match(name)]
        if bad:
            raise RuntimeError(f"Unexpected new package parts: {sorted(bad)}")
        missing = REPLACED_PARTS - generated_names
        if missing:
            raise RuntimeError(f"Generated package is missing required parts: {sorted(missing)}")

        for info in original.infolist():
            data = (
                generated.read(info.filename)
                if info.filename in REPLACED_PARTS
                else original.read(info.filename)
            )
            output.writestr(info, data)
        for name in sorted(new_parts):
            output.writestr(generated.getinfo(name), generated.read(name))


def main() -> None:
    before_slides = payload_hashes(DECK, "ppt/slides/slide")
    before_notes = payload_hashes(DECK, "ppt/notesSlides/notesSlide")
    original_count = len(Presentation(DECK).slides)

    with TemporaryDirectory(prefix="shapemix_count_only_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"

        presentation = Presentation(DECK)
        build_in_memory(presentation)
        presentation.save(generated)
        assemble_without_rewriting_existing_slides(DECK, generated, assembled)

        check = Presentation(assembled)
        expected_count = original_count + len(BUILDERS)
        if len(check.slides) != expected_count:
            raise RuntimeError(
                f"Expected {expected_count} slides; found {len(check.slides)}"
            )
        for offset, title in enumerate(NEW_TITLES):
            position = INSERT_AFTER_INDEX + 1 + offset
            texts = slide_texts(check.slides[position])
            if not any(title in text for text in texts):
                raise RuntimeError(
                    f"Slide {position + 1} does not contain the expected title {title!r}"
                )

        after_slides = payload_hashes(assembled, "ppt/slides/slide")
        after_notes = payload_hashes(assembled, "ppt/notesSlides/notesSlide")
        for name, digest in before_slides.items():
            if after_slides.get(name) != digest:
                raise RuntimeError(f"Existing slide payload changed: {name}")
        for name, digest in before_notes.items():
            if after_notes.get(name) != digest:
                raise RuntimeError(f"Existing notes payload changed: {name}")

        assembled.replace(DECK)

    print(
        f"Inserted {len(BUILDERS)} count-only Bayesian model slides after slide "
        f"{INSERT_AFTER_INDEX + 1}: {DECK}"
    )


if __name__ == "__main__":
    main()
