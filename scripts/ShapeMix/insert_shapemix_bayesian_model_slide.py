#!/usr/bin/env python3
"""Insert the ShapeMix Bayesian-model introduction after slide 13.

The slide introduces the Section 2 symbols R_LB, n_lb, and N_LB and shows the
compact Bayes formula from the beginning of Section 2.2. It deliberately does
not introduce the conditional multinomial.

The final package is assembled from the original ZIP so every pre-existing
slide and notes payload remains byte-for-byte unchanged.
"""

from __future__ import annotations

import re
from hashlib import sha256
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from insert_count_only_bayesian_slides import (
    BLUE,
    BLUE_PALE,
    C,
    CREAM,
    FONT_DISPLAY,
    GOLD,
    GOLD_PALE,
    INK,
    MID,
    NAVY,
    PURPLE,
    SLATE,
    TEAL,
    TEAL_PALE,
    WHITE,
    add_card,
    add_pill,
    add_shape,
    add_text,
    set_bg,
)
from update_slide9_bayes_compact import add_line

REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

INSERT_AFTER_INDEX = 12  # zero-based index of slide 13
PREVIOUS_TITLE = "Deconvolution works like identifying ingredients in a smoothie"
NEXT_TITLE = "Four symbols describe one spot's counts"
TITLE = "Bayesian model for ShapeMix"
DARK_GOLD = "B8860B"


def add_centered_text(slide, text, x, y, w, h, *, size, color, bold=False):
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_symbol_card(slide, x, role, symbol, name, meaning, fill, line):
    add_card(slide, x, 1.78, 3.75, 2.22, fill=fill, line=line)
    add_pill(
        slide,
        role,
        x + 1.10,
        2.00,
        1.55,
        0.31,
        line,
        color=NAVY if line == GOLD else WHITE,
    )
    add_centered_text(
        slide,
        symbol,
        x + 0.18,
        2.48,
        3.39,
        0.48,
        size=23,
        color=NAVY,
        bold=True,
    )
    add_centered_text(
        slide,
        name,
        x + 0.18,
        3.03,
        3.39,
        0.28,
        size=12,
        color=line if line != GOLD else DARK_GOLD,
        bold=True,
    )
    add_centered_text(
        slide,
        meaning,
        x + 0.24,
        3.40,
        3.27,
        0.42,
        size=10.8,
        color=INK,
    )


def add_definition(slide, symbol, text, y):
    add_text(
        slide,
        symbol,
        1.55,
        y,
        1.45,
        0.34,
        size=15,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.RIGHT,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        text,
        3.20,
        y,
        8.55,
        0.34,
        size=14,
        color=SLATE,
        valign=MSO_ANCHOR.MIDDLE,
    )


def build(slide) -> None:
    set_bg(slide, CREAM)
    add_text(
        slide,
        "SHAPEMIX • BAYESIAN MODEL",
        0.68,
        0.28,
        8.8,
        0.28,
        size=10.5,
        color=TEAL,
        bold=True,
    )
    add_text(
        slide,
        "Bayesian model for ShapeMix",
        0.65,
        0.55,
        11.8,
        0.62,
        size=27,
        color=NAVY,
        bold=True,
        font=FONT_DISPLAY,
    )

    add_definition(
        slide,
        "R_LB",
        "reference rates for the short, middle, and long fragment-length bins",
        1.20,
    )
    add_definition(
        slide,
        "n_lb",
        "predicted mean counts in each fragment-length bin",
        1.55,
    )
    add_definition(
        slide,
        "N_LB",
        "the cut-site counts actually observed in each fragment-length bin",
        1.90,
    )

    terms = [
        ("p(z | N_LB, N)", TEAL, 1.35, 3.15, "posterior", 2.93),
        ("∝", NAVY, 4.60, 0.50, None, None),
        ("p(N_LB, N | z)", DARK_GOLD, 5.12, 3.30, "likelihood", 6.77),
        ("×", NAVY, 8.45, 0.50, None, None),
        ("p(z)", PURPLE, 8.97, 1.70, "prior", 9.82),
    ]
    for text, color, x, width, label, center in terms:
        add_centered_text(
            slide,
            text,
            x,
            2.55,
            width,
            0.68,
            size=28,
            color=color,
            bold=True,
        )
        if label:
            add_centered_text(
                slide,
                label,
                center - 0.90,
                3.27,
                1.80,
                0.28,
                size=11,
                color=color,
                bold=True,
            )

    callouts = [
        (
            0.70,
            TEAL,
            TEAL_PALE,
            "POSTERIOR",
            "best supported z after seeing total counts N and length-bin counts N_LB",
            2.93,
            2.575,
        ),
        (
            4.79,
            DARK_GOLD,
            GOLD_PALE,
            "LIKELIHOOD",
            "given z, what is the probability of seeing data N_LB and N",
            6.77,
            6.665,
        ),
        (
            8.88,
            PURPLE,
            "EEE8FF",
            "PRIOR",
            "how plausible z is before seeing any data",
            9.82,
            10.755,
        ),
    ]
    for x, color, fill, heading, body, term_center, card_center in callouts:
        add_line(slide, term_center, 3.60, card_center, 4.10, color)
        add_card(slide, x, 4.10, 3.75, 1.30, fill=fill, line=color)
        add_text(
            slide,
            heading,
            x + 0.16,
            4.20,
            3.43,
            0.26,
            size=11,
            color=color,
            bold=True,
        )
        add_text(
            slide,
            body,
            x + 0.16,
            4.52,
            3.43,
            0.65,
            size=14,
            color=INK,
        )

    hairline = add_shape(slide, MSO_SHAPE.RECTANGLE, 0.68, 7.13, 11.97, 0.011, MID)
    hairline.line.width = Pt(0.25)


def slide_texts(slide) -> list[str]:
    return [
        shape.text_frame.text.strip()
        for shape in slide.shapes
        if shape.has_text_frame and shape.text_frame.text.strip()
    ]


def part_hashes(path: Path, prefixes: tuple[str, ...]) -> dict[str, str]:
    with ZipFile(path) as archive:
        return {
            name: sha256(archive.read(name)).hexdigest()
            for name in archive.namelist()
            if name.startswith(prefixes)
        }


NEW_PART_PATTERN = re.compile(
    r"^ppt/slides/(?:_rels/)?slide\d+\.xml(?:\.rels)?$"
)
REPLACED_PARTS = {
    "[Content_Types].xml",
    "ppt/presentation.xml",
    "ppt/_rels/presentation.xml.rels",
    "docProps/app.xml",
}


def assemble(original_path: Path, generated_path: Path, output_path: Path) -> None:
    with ZipFile(original_path, "r") as original, ZipFile(
        generated_path, "r"
    ) as generated, ZipFile(output_path, "w", compression=ZIP_DEFLATED) as output:
        original_names = set(original.namelist())
        generated_names = set(generated.namelist())
        new_parts = generated_names - original_names
        bad = [name for name in new_parts if not NEW_PART_PATTERN.match(name)]
        if bad:
            raise RuntimeError(f"Unexpected new package parts: {sorted(bad)}")
        missing = REPLACED_PARTS - generated_names
        if missing:
            raise RuntimeError(f"Generated package is missing parts: {sorted(missing)}")

        for info in original.infolist():
            data = (
                generated.read(info.filename)
                if info.filename in REPLACED_PARTS
                else original.read(info.filename)
            )
            output.writestr(info, data)
        for name in sorted(new_parts):
            output.writestr(generated.getinfo(name), generated.read(name))


def main() -> None:
    original = Presentation(DECK)
    original_count = len(original.slides)
    if original_count != 32:
        raise RuntimeError(f"Expected 32 slides; found {original_count}")
    for slide in original.slides:
        if any(TITLE in text for text in slide_texts(slide)):
            raise RuntimeError("The ShapeMix Bayesian-model slide already exists")
    if not any(PREVIOUS_TITLE in text for text in slide_texts(original.slides[12])):
        raise RuntimeError("Slide 13 is not the expected deconvolution slide")
    if not any(NEXT_TITLE in text for text in slide_texts(original.slides[13])):
        raise RuntimeError("Slide 14 is not the expected symbols slide")

    preserved_before = part_hashes(
        DECK,
        ("ppt/slides/", "ppt/notesSlides/"),
    )

    with TemporaryDirectory(prefix="shapemix_bayesian_intro_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"

        presentation = Presentation(DECK)
        new_slide = presentation.slides.add_slide(presentation.slides[12].slide_layout)
        build(new_slide)
        slide_ids = presentation.slides._sldIdLst
        new_id = slide_ids[-1]
        slide_ids.remove(new_id)
        slide_ids.insert(INSERT_AFTER_INDEX + 1, new_id)
        presentation.save(generated)

        assemble(DECK, generated, assembled)

        check = Presentation(assembled)
        if len(check.slides) != original_count + 1:
            raise RuntimeError(
                f"Expected {original_count + 1} slides; found {len(check.slides)}"
            )
        new_text = " ".join(slide_texts(check.slides[13]))
        for needle in (
            TITLE,
            "R_LB[c,p,b]",
            "n_lb[s,p,b]",
            "N_LB[s,p,b]",
            "p(z | N_LB, N)",
        ):
            if needle not in new_text:
                raise RuntimeError(f"New slide 14 is missing {needle!r}")
        if "Multinomial" in new_text:
            raise RuntimeError("New slide 14 must not introduce the multinomial")
        if not any(PREVIOUS_TITLE in text for text in slide_texts(check.slides[12])):
            raise RuntimeError("The original slide 13 was changed")
        if not any(NEXT_TITLE in text for text in slide_texts(check.slides[14])):
            raise RuntimeError("The original slide 14 did not move intact to position 15")

        preserved_after = part_hashes(
            assembled,
            ("ppt/slides/", "ppt/notesSlides/"),
        )
        for name, digest in preserved_before.items():
            if preserved_after.get(name) != digest:
                raise RuntimeError(f"Existing slide or notes payload changed: {name}")

        assembled.replace(DECK)

    print(f"Inserted ShapeMix Bayesian-model slide after slide 13: {DECK}")
    print(f"Slide count: {original_count} -> {original_count + 1}")
    print("Verified: every pre-existing slide and notes payload is byte-identical")


if __name__ == "__main__":
    main()
