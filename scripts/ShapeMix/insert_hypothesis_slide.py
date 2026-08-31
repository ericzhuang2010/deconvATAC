#!/usr/bin/env python3
"""Insert a hypothesis slide after slide 8 without rewriting existing slides."""

from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from insert_count_only_bayesian_slides import (
    BLUE,
    BLUE_PALE,
    CORAL,
    CORAL_PALE,
    CREAM,
    FONT_DISPLAY,
    GOLD,
    GOLD_PALE,
    INK,
    MID,
    NAVY,
    SLATE,
    TEAL,
    TEAL_PALE,
    WHITE,
    add_card,
    add_pill,
    add_shape,
    add_text,
    set_bg,
)
from insert_shapemix_multinomial_slide import (
    assemble_with_insert,
    next_slide_part,
    package_hashes,
    slide_part_at,
    slide_texts,
)

REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

INSERT_AFTER_INDEX = 7  # zero-based index of slide 8
PREVIOUS_TITLE = "ShapeMix uses an ATAC-specific signal that RNA methods miss"
NEXT_TITLE = "Bayesian model to find the number of cells for each cell type"
TITLE = "Our hypothesis"
HYPOTHESIS = (
    "A Bayesian model using ATAC fragment-length signals will recover cell-type "
    "mixtures in spatial ATAC-seq spots more accurately than models using peak "
    "counts alone."
)


def add_centered_text(slide, text, x, y, w, h, *, size, color, bold=False):
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_signal_card(slide, x, pill, heading, body, fill, line):
    add_card(slide, x, 4.47, 3.26, 1.32, fill=fill, line=line)
    add_pill(slide, pill, x + 0.17, 4.61, 1.20, 0.28, line, size=9)
    add_text(
        slide,
        heading,
        x + 0.17,
        5.00,
        2.92,
        0.29,
        size=14,
        color=NAVY,
        bold=True,
    )
    add_text(
        slide,
        body,
        x + 0.17,
        5.30,
        2.92,
        0.42,
        size=10.5,
        color=INK,
        valign=MSO_ANCHOR.MIDDLE,
    )


def build(slide) -> None:
    set_bg(slide, CREAM)
    add_text(
        slide,
        "RESEARCH QUESTION • HYPOTHESIS",
        0.68,
        0.30,
        8.0,
        0.28,
        size=10.5,
        color=TEAL,
        bold=True,
    )
    add_text(
        slide,
        TITLE,
        0.65,
        0.62,
        8.0,
        0.63,
        size=29,
        color=NAVY,
        bold=True,
        font=FONT_DISPLAY,
    )

    add_card(slide, 0.68, 1.54, 11.98, 2.35, fill=TEAL_PALE, line=TEAL)
    add_pill(slide, "PREDICTION", 1.00, 1.84, 1.36, 0.33, TEAL, size=9.5)
    add_centered_text(
        slide,
        HYPOTHESIS,
        1.05,
        2.24,
        11.23,
        1.14,
        size=23,
        color=NAVY,
        bold=True,
    )

    add_signal_card(
        slide,
        0.68,
        "BASELINE",
        "Peak counts",
        "how many cut sites fall in each peak",
        BLUE_PALE,
        BLUE,
    )
    add_centered_text(
        slide,
        "+",
        4.05,
        4.77,
        0.55,
        0.58,
        size=28,
        color=SLATE,
        bold=True,
    )
    add_signal_card(
        slide,
        4.70,
        "ATAC SIGNAL",
        "Fragment lengths",
        "how those same fragments split by length",
        CORAL_PALE,
        CORAL,
    )

    arrow = add_shape(
        slide,
        MSO_SHAPE.CHEVRON,
        8.13,
        4.82,
        0.66,
        0.57,
        GOLD,
        GOLD,
    )
    arrow.line.width = Pt(0)

    add_signal_card(
        slide,
        8.92,
        "EXPECTED RESULT",
        "Better deconvolution",
        "more accurate estimates of the cell-type mixture",
        GOLD_PALE,
        GOLD,
    )

    add_centered_text(
        slide,
        "The comparison makes the hypothesis testable: fragment-length model vs. count-only model.",
        0.80,
        6.25,
        11.72,
        0.38,
        size=12,
        color=SLATE,
        bold=True,
    )

    hairline = add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        0.68,
        7.13,
        11.97,
        0.011,
        MID,
    )
    hairline.line.width = Pt(0.25)


def refresh_inserted_slide() -> None:
    """Rebuild only the inserted slide while preserving every other package part."""
    presentation = Presentation(DECK)
    slide_count = len(presentation.slides)
    target_index = INSERT_AFTER_INDEX + 1
    if not any(TITLE == text for text in slide_texts(presentation.slides[target_index])):
        raise RuntimeError("Slide 9 is not the inserted hypothesis slide")

    target_slide_part = slide_part_at(DECK, target_index)
    target_rels_part = (
        f"ppt/slides/_rels/{Path(target_slide_part).name}.rels"
    )
    before = package_hashes(DECK)

    with TemporaryDirectory(prefix="shapemix_hypothesis_refresh_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"

        source = Presentation(DECK)
        rebuilt = source.slides.add_slide(source.slides[7].slide_layout)
        build(rebuilt)
        source.save(generated)
        generated_slide_part = slide_part_at(generated, slide_count)
        generated_rels_part = (
            f"ppt/slides/_rels/{Path(generated_slide_part).name}.rels"
        )

        with ZipFile(DECK, "r") as original, ZipFile(
            generated, "r"
        ) as generated_archive, ZipFile(assembled, "w") as output:
            if generated_rels_part not in generated_archive.namelist():
                raise RuntimeError("Generated hypothesis slide has no relationships")
            for info in original.infolist():
                if info.filename == target_slide_part:
                    payload = generated_archive.read(generated_slide_part)
                elif info.filename == target_rels_part:
                    payload = generated_archive.read(generated_rels_part)
                else:
                    payload = original.read(info.filename)
                output.writestr(info, payload)

        check = Presentation(assembled)
        if len(check.slides) != slide_count:
            raise RuntimeError("Refreshing the hypothesis slide changed slide count")
        refreshed_text = " ".join(slide_texts(check.slides[target_index]))
        if HYPOTHESIS not in refreshed_text:
            raise RuntimeError("Refreshed slide is missing the hypothesis")

        after = package_hashes(assembled)
        allowed_changes = {target_slide_part, target_rels_part}
        for name, digest in before.items():
            if name in allowed_changes:
                continue
            if after.get(name) != digest:
                raise RuntimeError(
                    f"Unexpected change while refreshing hypothesis slide: {name}"
                )
        assembled.replace(DECK)


def main() -> None:
    original = Presentation(DECK)
    original_count = len(original.slides)
    if original_count < 9:
        raise RuntimeError(f"Expected at least 9 slides; found {original_count}")
    for slide in original.slides:
        if any(TITLE == text for text in slide_texts(slide)):
            raise RuntimeError("The hypothesis slide already exists")
    if not any(
        PREVIOUS_TITLE in text for text in slide_texts(original.slides[7])
    ):
        raise RuntimeError("Slide 8 is not the expected ATAC-specific signal slide")
    if not any(NEXT_TITLE in text for text in slide_texts(original.slides[8])):
        raise RuntimeError("Slide 9 is not the expected Bayesian model slide")

    before = package_hashes(DECK)
    with ZipFile(DECK) as archive:
        new_slide_part = next_slide_part(set(archive.namelist()))

    with TemporaryDirectory(prefix="shapemix_hypothesis_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"

        presentation = Presentation(DECK)
        new_slide = presentation.slides.add_slide(
            presentation.slides[7].slide_layout
        )
        build(new_slide)
        presentation.save(generated)
        generated_slide_part = slide_part_at(generated, original_count)

        assemble_with_insert(
            DECK,
            generated,
            assembled,
            generated_slide_part,
            new_slide_part,
            INSERT_AFTER_INDEX + 1,
            original_count + 1,
        )

        check = Presentation(assembled)
        if len(check.slides) != original_count + 1:
            raise RuntimeError(
                f"Expected {original_count + 1} slides; found {len(check.slides)}"
            )
        new_text = " ".join(slide_texts(check.slides[8]))
        for needle in (
            TITLE,
            "fragment-length signals",
            "peak counts alone",
            "Better deconvolution",
            "testable",
        ):
            if needle not in new_text:
                raise RuntimeError(f"New slide 9 is missing {needle!r}")
        if not any(
            PREVIOUS_TITLE in text for text in slide_texts(check.slides[7])
        ):
            raise RuntimeError("The original slide 8 was changed")
        if not any(NEXT_TITLE in text for text in slide_texts(check.slides[9])):
            raise RuntimeError(
                "The original slide 9 did not move intact to position 10"
            )

        after = package_hashes(assembled)
        allowed_changes = {
            "[Content_Types].xml",
            "ppt/presentation.xml",
            "ppt/_rels/presentation.xml.rels",
            "docProps/app.xml",
        }
        for name, digest in before.items():
            if name in allowed_changes:
                continue
            if after.get(name) != digest:
                raise RuntimeError(
                    f"Unexpected change to existing package part: {name}"
                )

        assembled.replace(DECK)

    print(f"Inserted hypothesis slide after slide 8: {DECK}")
    print(f"Slide count: {original_count} -> {original_count + 1}")
    print("Verified: every pre-existing slide and notes payload is byte-identical")


if __name__ == "__main__":
    main()
