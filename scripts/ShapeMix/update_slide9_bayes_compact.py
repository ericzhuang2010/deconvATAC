#!/usr/bin/env python3
"""Rebuild slide 9 (Bayes' rule) as a compact formula with per-term callouts.

Only slide 9's XML payload (and its speaker notes) are replaced; every other
package part is copied byte-for-byte from the current deck.
"""

from __future__ import annotations

from hashlib import sha256
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches, Pt

from insert_count_only_bayesian_slides import (
    CREAM,
    INK,
    KICKER,
    NAVY,
    PURPLE,
    PURPLE_PALE,
    SLATE,
    TEAL,
    TEAL_PALE,
    GOLD_PALE,
    C,
    add_card,
    add_footer,
    add_notes,
    add_takeaway,
    add_text,
    add_title,
    set_bg,
)
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

SLIDE_INDEX = 8  # zero-based; slide 9
TITLE = "Bayes' rule picks the best-supported amounts"
DARK_GOLD = "B8860B"


def add_line(slide, x1, y1, x2, y2, color, width=1.2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = C(color)
    line.line.width = Pt(width)
    return line


def rebuild(slide) -> None:
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)

    set_bg(slide, CREAM)
    add_title(slide, TITLE)
    add_text(
        slide,
        "Every candidate recipe z receives one score; the reported answer is the "
        "highest-scoring recipe.",
        0.71, 1.17, 11.9, 0.4, size=14, color=SLATE,
    )

    # Compact formula: separate boxes per term so the labels and callout leader
    # lines can align exactly under each term.
    terms = [
        ("p(z | N)", TEAL, 2.80, 2.30, "posterior", 3.95),
        ("∝", NAVY, 5.20, 0.50, None, None),
        ("p(N | z)", DARK_GOLD, 5.80, 2.30, "likelihood", 6.95),
        ("×", NAVY, 8.20, 0.50, None, None),
        ("p(z)", PURPLE, 8.80, 1.70, "prior", 9.65),
    ]
    for text, color, x, w, label, center in terms:
        add_text(
            slide, text, x, 2.12, w, 0.68,
            size=30, color=color, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
        )
        if label:
            add_text(
                slide, label, center - 0.9, 2.86, 1.8, 0.28,
                size=11, color=color, bold=True, align=PP_ALIGN.CENTER,
            )

    callouts = [
        (0.70, TEAL, TEAL_PALE, "POSTERIOR",
         "Updated support for each candidate recipe z after seeing the data. "
         "ShapeMix reports the MAP — the single best-supported z — normalized to "
         "cell-type percentages.",
         3.95, 2.575),
        (4.79, DARK_GOLD, GOLD_PALE, "LIKELIHOOD",
         "The data-fit score: the negative-binomial probability of the observed "
         "count N when the prediction is n = z · R (previous slide).",
         6.95, 6.665),
        (8.88, PURPLE, PURPLE_PALE, "PRIOR",
         "Plausibility before seeing this spot: z ~ Gamma(2, 1) — positive "
         "amounts, mean 2, mode 1. The same prior for every cell type, so none "
         "is favored.",
         9.65, 10.755),
    ]
    y, w, h = 3.75, 3.75, 1.95
    for x, color, pale, heading, body, term_center, card_center in callouts:
        add_line(slide, term_center, 3.18, card_center, y, color)
        add_card(slide, x, y, w, h, fill=pale, line=color)
        add_text(slide, heading, x + 0.16, y + 0.12, w - 0.32, 0.26,
                 size=11, color=color, bold=True)
        add_text(slide, body, x + 0.16, y + 0.44, w - 0.32, h - 0.56,
                 size=10.5, color=INK)

    add_takeaway(
        slide,
        "MAP = THE HIGHEST POINT OF THE POSTERIOR — ONE POINT ESTIMATE, NOT A CONFIDENCE INTERVAL",
        6.05,
    )
    add_text(
        slide,
        "ShapeMix keeps this exact machinery and adds one extra likelihood term for "
        "fragment-length bins (later sections).",
        0.71, 6.68, 11.9, 0.3, size=10.5, color=SLATE, align=PP_ALIGN.CENTER,
    )
    add_footer(slide)
    add_notes(
        slide,
        "Bayes' rule shown as one compact line: posterior is proportional to "
        "likelihood times prior. The teal posterior is the updated support whose "
        "highest point (MAP) is reported; the gold likelihood is the "
        "negative-binomial data-fit score from the previous slide; the purple "
        "prior is the fixed Gamma(2,1) preference for moderate positive amounts.",
    )


def slide_texts(slide) -> list[str]:
    return [shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)]


def payload_hashes(path: Path) -> dict[str, str]:
    with ZipFile(path) as archive:
        return {name: sha256(archive.read(name)).hexdigest() for name in archive.namelist()}


def assemble(original_path: Path, generated_path: Path, output_path: Path, replaced: set[str]) -> None:
    with ZipFile(original_path, "r") as original, ZipFile(
        generated_path, "r"
    ) as generated, ZipFile(output_path, "w", compression=ZIP_DEFLATED) as output:
        missing = replaced - set(generated.namelist())
        if missing:
            raise RuntimeError(f"Generated package is missing parts: {sorted(missing)}")
        for info in original.infolist():
            data = (
                generated.read(info.filename)
                if info.filename in replaced
                else original.read(info.filename)
            )
            output.writestr(info, data)


def main() -> None:
    before = payload_hashes(DECK)

    with TemporaryDirectory(prefix="shapemix_slide9_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"

        presentation = Presentation(DECK)
        slide = presentation.slides[SLIDE_INDEX]
        texts = slide_texts(slide)
        if not any(TITLE in text for text in texts) or not any(KICKER in text for text in texts):
            raise RuntimeError("Slide 9 does not look like the Bayes slide; refusing to rewrite")

        slide_part = str(slide.part.partname).lstrip("/")
        notes_part = str(slide.notes_slide.part.partname).lstrip("/")
        replaced = {slide_part, notes_part}

        rebuild(slide)
        presentation.save(generated)
        assemble(DECK, generated, assembled, replaced)

        after = payload_hashes(assembled)
        for name, digest in before.items():
            if name in replaced:
                continue
            if after.get(name) != digest:
                raise RuntimeError(f"Unexpected change to package part: {name}")

        check = Presentation(assembled)
        if len(check.slides) != 28:
            raise RuntimeError(f"Expected 28 slides; found {len(check.slides)}")
        new_texts = " ".join(slide_texts(check.slides[SLIDE_INDEX]))
        for needle in ("∝", "p(z | N)", "Gamma(2, 1)", TITLE):
            if needle not in new_texts:
                raise RuntimeError(f"Rebuilt slide 9 is missing {needle!r}")

        assembled.replace(DECK)

    print(f"Rebuilt slide 9 with a compact Bayes formula and callouts: {DECK}")


if __name__ == "__main__":
    main()
