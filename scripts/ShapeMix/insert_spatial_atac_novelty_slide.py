#!/usr/bin/env python3
"""Insert a slide on ShapeMix-ATAC's novelty immediately after slide 7.

The final package is assembled from the original ZIP so every pre-existing
slide and notes payload remains byte-for-byte unchanged.
"""

from __future__ import annotations

import re
from hashlib import sha256
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from insert_count_only_bayesian_slides import (
    BLUE,
    BLUE_PALE,
    C,
    CORAL,
    CORAL_PALE,
    CREAM,
    FONT_DISPLAY,
    GOLD,
    GOLD_PALE,
    INK,
    MID,
    NAVY,
    PURPLE,
    PURPLE_PALE,
    SLATE,
    TEAL,
    TEAL_PALE,
    WHITE,
    add_card,
    add_pill,
    add_shape,
    add_text,
    set_bg,
)

REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

INSERT_AFTER_INDEX = 6  # zero-based index of slide 7
PREVIOUS_TITLE = "A better deconvolution method can help many studies"
NEXT_TITLE = "Bayesian model to find the number of cells for each cell type"
TITLE = "ShapeMix uses an ATAC-specific signal that RNA methods miss"


def add_centered_text(slide, text, x, y, w, h, *, size, color, bold=False):
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_peak_totals(slide):
    baseline_y = 4.57
    add_shape(slide, MSO_SHAPE.RECTANGLE, 1.35, baseline_y, 4.45, 0.015, MID)
    heights = [0.50, 1.16, 0.72, 1.42, 0.91]
    for i, height in enumerate(heights):
        x = 1.62 + i * 0.83
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            baseline_y - height,
            0.36,
            height,
            BLUE,
            line=BLUE,
        )
        add_centered_text(
            slide,
            f"P{i + 1}",
            x - 0.08,
            baseline_y + 0.06,
            0.52,
            0.22,
            size=9,
            color=SLATE,
        )


def add_fragment(slide, x, y, width, color):
    # Parent fragment with two dark Tn5 cut-site endpoints.
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y + 0.09, width, 0.045, color)
    for endpoint_x in (x - 0.04, x + width - 0.04):
        add_shape(
            slide,
            MSO_SHAPE.OVAL,
            endpoint_x,
            y + 0.025,
            0.13,
            0.13,
            NAVY,
            line=NAVY,
        )


def add_nucleosome(slide, x, y, label):
    circle = add_shape(
        slide,
        MSO_SHAPE.OVAL,
        x,
        y,
        0.42,
        0.42,
        PURPLE_PALE,
        line=PURPLE,
    )
    circle.line.width = Pt(1.2)
    add_centered_text(
        slide,
        label,
        x,
        y + 0.01,
        0.42,
        0.40,
        size=8.5,
        color=PURPLE,
        bold=True,
    )


def add_length_row(slide, y, label, description, color, kind):
    add_pill(
        slide,
        label,
        7.16,
        y,
        0.93,
        0.28,
        color,
        color=NAVY if color == GOLD else WHITE,
        size=9.5,
    )
    if kind == "short":
        add_fragment(slide, 8.38, y + 0.01, 0.82, color)
    elif kind == "one":
        add_fragment(slide, 8.24, y + 0.01, 1.28, color)
        add_nucleosome(slide, 8.67, y - 0.06, "1")
    else:
        add_fragment(slide, 8.08, y + 0.01, 1.62, color)
        add_nucleosome(slide, 8.39, y - 0.06, "1")
        add_nucleosome(slide, 8.94, y - 0.06, "2")
    add_text(
        slide,
        description,
        9.92,
        y - 0.01,
        2.23,
        0.35,
        size=10.5,
        color=INK,
        valign=MSO_ANCHOR.MIDDLE,
    )


def build(slide) -> None:
    set_bg(slide, CREAM)
    add_text(
        slide,
        "WHY THIS WORK IS NEW",
        0.68,
        0.28,
        8.8,
        0.28,
        size=10.5,
        color=TEAL,
        bold=True,
    )
    add_text(
        slide,
        TITLE,
        0.65,
        0.55,
        11.8,
        0.62,
        size=25,
        color=NAVY,
        bold=True,
        font=FONT_DISPLAY,
    )
    add_centered_text(
        slide,
        "Spatial ATAC deconvolution is still an open problem.",
        0.75,
        1.25,
        11.83,
        0.38,
        size=16,
        color=SLATE,
        bold=True,
    )

    # Current state of the art.
    add_card(slide, 0.70, 1.82, 5.78, 3.55, fill=BLUE_PALE, line=BLUE)
    add_pill(slide, "2025 STATE OF THE ART", 0.98, 2.05, 2.18, 0.31, BLUE)
    add_text(
        slide,
        "Borrow spatial-RNA methods",
        0.98,
        2.55,
        5.18,
        0.42,
        size=19,
        color=NAVY,
        bold=True,
    )
    add_text(
        slide,
        "The benchmark applies methods designed for gene-expression data to an ATAC peak-count matrix.",
        0.98,
        3.03,
        5.18,
        0.66,
        size=12,
        color=INK,
    )
    add_peak_totals(slide)
    add_centered_text(
        slide,
        "Peak totals only  →  fragment lengths are collapsed away",
        0.98,
        4.82,
        5.18,
        0.34,
        size=11.5,
        color=BLUE,
        bold=True,
    )

    # ShapeMix-ATAC novelty.
    add_card(slide, 6.82, 1.82, 5.80, 3.55, fill=TEAL_PALE, line=TEAL)
    add_pill(slide, "SHAPEMIX-ATAC", 7.10, 2.05, 1.70, 0.31, TEAL)
    add_text(
        slide,
        "Keep an ATAC-specific signal",
        7.10,
        2.55,
        5.20,
        0.42,
        size=19,
        color=NAVY,
        bold=True,
    )
    add_text(
        slide,
        "DNA wrapping around nucleosomes creates informative fragment-size patterns:",
        7.10,
        3.00,
        5.20,
        0.48,
        size=12,
        color=INK,
    )
    add_length_row(
        slide,
        3.64,
        "SHORT",
        "open, nucleosome-free DNA",
        CORAL,
        "short",
    )
    add_length_row(
        slide,
        4.12,
        "MIDDLE",
        "spans about one nucleosome",
        GOLD,
        "one",
    )
    add_length_row(
        slide,
        4.60,
        "LONG",
        "can span two nucleosomes",
        PURPLE,
        "two",
    )

    novelty = add_card(slide, 0.70, 5.72, 11.92, 0.72, fill=CORAL_PALE, line=CORAL)
    novelty.line.width = Pt(1.3)
    add_centered_text(
        slide,
        "NOVELTY: model the fragment-length composition inside each peak—not peak counts alone",
        0.95,
        5.83,
        11.42,
        0.48,
        size=15,
        color=NAVY,
        bold=True,
    )
    add_centered_text(
        slide,
        "RNA measurements have no equivalent fragment-length feature.",
        0.70,
        6.56,
        11.92,
        0.30,
        size=11.5,
        color=SLATE,
    )

    hairline = add_shape(slide, MSO_SHAPE.RECTANGLE, 0.68, 7.13, 11.97, 0.011, MID)
    hairline.line.width = Pt(0.25)
    add_text(
        slide,
        "Source: Ouologuem et al., Bioinformatics (2025) • ShapeMix research proposal",
        0.70,
        7.17,
        11.80,
        0.18,
        size=7.5,
        color=SLATE,
    )


def slide_texts(slide) -> list[str]:
    return [
        shape.text_frame.text.strip()
        for shape in slide.shapes
        if shape.has_text_frame and shape.text_frame.text.strip()
    ]


def part_hashes(path: Path, prefixes: tuple[str, ...]) -> dict[str, str]:
    with ZipFile(path) as archive:
        return {
            name: sha256(archive.read(name)).hexdigest()
            for name in archive.namelist()
            if name.startswith(prefixes)
        }


NEW_PART_PATTERN = re.compile(
    r"^ppt/slides/(?:_rels/)?slide\d+\.xml(?:\.rels)?$"
)
REPLACED_PARTS = {
    "[Content_Types].xml",
    "ppt/presentation.xml",
    "ppt/_rels/presentation.xml.rels",
    "docProps/app.xml",
}


def assemble(original_path: Path, generated_path: Path, output_path: Path) -> None:
    with ZipFile(original_path, "r") as original, ZipFile(
        generated_path, "r"
    ) as generated, ZipFile(output_path, "w", compression=ZIP_DEFLATED) as output:
        original_names = set(original.namelist())
        generated_names = set(generated.namelist())
        new_parts = generated_names - original_names
        bad = [name for name in new_parts if not NEW_PART_PATTERN.match(name)]
        if bad:
            raise RuntimeError(f"Unexpected new package parts: {sorted(bad)}")
        missing = REPLACED_PARTS - generated_names
        if missing:
            raise RuntimeError(f"Generated package is missing parts: {sorted(missing)}")

        for info in original.infolist():
            data = (
                generated.read(info.filename)
                if info.filename in REPLACED_PARTS
                else original.read(info.filename)
            )
            output.writestr(info, data)
        for name in sorted(new_parts):
            output.writestr(generated.getinfo(name), generated.read(name))


def main() -> None:
    original = Presentation(DECK)
    original_count = len(original.slides)
    if original_count != 31:
        raise RuntimeError(f"Expected 31 slides; found {original_count}")
    for slide in original.slides:
        if any(TITLE in text for text in slide_texts(slide)):
            raise RuntimeError("The novelty slide already exists")
    if not any(PREVIOUS_TITLE in text for text in slide_texts(original.slides[6])):
        raise RuntimeError("Slide 7 is not the expected impact slide")
    if not any(NEXT_TITLE in text for text in slide_texts(original.slides[7])):
        raise RuntimeError("Slide 8 is not the expected Bayesian-model slide")

    preserved_before = part_hashes(
        DECK,
        ("ppt/slides/", "ppt/notesSlides/"),
    )

    with TemporaryDirectory(prefix="shapemix_novelty_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"

        presentation = Presentation(DECK)
        new_slide = presentation.slides.add_slide(presentation.slides[6].slide_layout)
        build(new_slide)
        slide_ids = presentation.slides._sldIdLst
        new_id = slide_ids[-1]
        slide_ids.remove(new_id)
        slide_ids.insert(INSERT_AFTER_INDEX + 1, new_id)
        presentation.save(generated)

        assemble(DECK, generated, assembled)

        check = Presentation(assembled)
        if len(check.slides) != original_count + 1:
            raise RuntimeError(
                f"Expected {original_count + 1} slides; found {len(check.slides)}"
            )
        new_text = " ".join(slide_texts(check.slides[7]))
        for needle in (
            TITLE,
            "2025 STATE OF THE ART",
            "SHAPEMIX-ATAC",
            "fragment-length composition inside each peak",
        ):
            if needle not in new_text:
                raise RuntimeError(f"New slide 8 is missing {needle!r}")
        if not any(PREVIOUS_TITLE in text for text in slide_texts(check.slides[6])):
            raise RuntimeError("The original slide 7 was changed")
        if not any(NEXT_TITLE in text for text in slide_texts(check.slides[8])):
            raise RuntimeError("The original slide 8 did not move intact to position 9")

        preserved_after = part_hashes(
            assembled,
            ("ppt/slides/", "ppt/notesSlides/"),
        )
        for name, digest in preserved_before.items():
            if preserved_after.get(name) != digest:
                raise RuntimeError(f"Existing slide or notes payload changed: {name}")

        assembled.replace(DECK)

    print(f"Inserted spatial-ATAC novelty slide after slide 7: {DECK}")
    print(f"Slide count: {original_count} -> {original_count + 1}")
    print("Verified: every pre-existing slide and notes payload is byte-identical")


if __name__ == "__main__":
    main()
