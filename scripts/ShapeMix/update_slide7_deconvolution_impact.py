#!/usr/bin/env python3
"""Simplify slide 7 to emphasize the broad impact of better deconvolution.

Only slide 7's XML payload is replaced. Every other package part remains
byte-for-byte unchanged.
"""

from __future__ import annotations

from hashlib import sha256
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from insert_count_only_bayesian_slides import (
    BLUE,
    BLUE_PALE,
    C,
    CREAM,
    FONT_DISPLAY,
    INK,
    MID,
    NAVY,
    SLATE,
    TEAL,
    TEAL_PALE,
    add_card,
    add_pill,
    add_shape,
    add_text,
    set_bg,
)

REPO_ROOT = Path(__file__).resolve().parents[2]
DECK = (
    REPO_ROOT
    / "docs"
    / "ShapeMix"
    / "presentations"
    / "ShapeMix_High_School_Research_Deck.pptx"
)

SLIDE_INDEX = 6
OLD_TITLE = "Why deconvolution matters for spatial ATAC-seq"
TITLE = "A better deconvolution method can help many studies"
GREEN = "43AA6A"
GREEN_PALE = "E5F5E9"


def add_centered_text(slide, text, x, y, w, h, *, size, color, bold=False):
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_impact_card(slide, x, pill, heading, body, fill, line):
    add_card(slide, x, 2.23, 3.56, 2.70, fill=fill, line=line)
    add_pill(slide, pill, x + 0.78, 2.47, 2.00, 0.32, line)
    add_centered_text(
        slide,
        heading,
        x + 0.20,
        3.08,
        3.16,
        0.48,
        size=20,
        color=NAVY,
        bold=True,
    )
    add_centered_text(
        slide,
        body,
        x + 0.28,
        3.73,
        3.00,
        0.74,
        size=12.5,
        color=INK,
    )


def build(slide) -> None:
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)

    set_bg(slide, CREAM)
    add_text(
        slide,
        "WHY THIS WORK MATTERS",
        0.68,
        0.28,
        8.8,
        0.28,
        size=10.5,
        color=TEAL,
        bold=True,
    )
    add_text(
        slide,
        TITLE,
        0.65,
        0.55,
        11.8,
        0.62,
        size=27,
        color=NAVY,
        bold=True,
        font=FONT_DISPLAY,
    )
    add_centered_text(
        slide,
        "Researchers across genomics use deconvolution to interpret mixed measurements.",
        0.75,
        1.35,
        11.83,
        0.40,
        size=16,
        color=SLATE,
        bold=True,
    )

    add_impact_card(
        slide,
        0.70,
        "COMMON PROBLEM",
        "Mixed data",
        "A sample or spatial spot can contain several cell types.",
        BLUE_PALE,
        BLUE,
    )
    add_centered_text(slide, "→", 4.31, 3.17, 0.48, 0.55, size=30, color=TEAL, bold=True)
    add_impact_card(
        slide,
        4.88,
        "WIDELY USED TOOL",
        "Deconvolution",
        "Estimate how much of each cell type is present.",
        TEAL_PALE,
        TEAL,
    )
    add_centered_text(slide, "→", 8.49, 3.17, 0.48, 0.55, size=30, color=TEAL, bold=True)
    add_impact_card(
        slide,
        9.06,
        "BROAD IMPACT",
        "Better conclusions",
        "More reliable tissue maps, comparisons, and biological interpretation.",
        GREEN_PALE,
        GREEN,
    )

    callout = add_card(slide, 0.70, 5.40, 11.92, 0.78, fill=TEAL_PALE, line=TEAL)
    callout.line.width = Pt(1.4)
    add_centered_text(
        slide,
        "Improve a method used by many researchers  →  improve many downstream results",
        0.95,
        5.53,
        11.42,
        0.50,
        size=16,
        color=NAVY,
        bold=True,
    )
    add_centered_text(
        slide,
        "ShapeMix asks whether fragment-length information can make this widely used step more accurate.",
        0.70,
        6.48,
        11.92,
        0.35,
        size=11.5,
        color=SLATE,
    )

    hairline = add_shape(slide, MSO_SHAPE.RECTANGLE, 0.68, 7.13, 11.97, 0.011, MID)
    hairline.line.width = Pt(0.25)


def slide_texts(slide) -> list[str]:
    return [
        shape.text_frame.text.strip()
        for shape in slide.shapes
        if shape.has_text_frame and shape.text_frame.text.strip()
    ]


def package_hashes(path: Path) -> dict[str, str]:
    with ZipFile(path) as archive:
        return {
            name: sha256(archive.read(name)).hexdigest()
            for name in archive.namelist()
        }


def assemble(
    original_path: Path,
    generated_path: Path,
    output_path: Path,
    slide_part: str,
) -> None:
    with ZipFile(original_path, "r") as original, ZipFile(
        generated_path, "r"
    ) as generated, ZipFile(output_path, "w", compression=ZIP_DEFLATED) as output:
        if slide_part not in generated.namelist():
            raise RuntimeError(f"Generated package is missing {slide_part}")
        for info in original.infolist():
            data = (
                generated.read(info.filename)
                if info.filename == slide_part
                else original.read(info.filename)
            )
            output.writestr(info, data)


def main() -> None:
    before = package_hashes(DECK)
    current = Presentation(DECK)
    if len(current.slides) != 31:
        raise RuntimeError(f"Expected 31 slides; found {len(current.slides)}")
    if not any(OLD_TITLE in text for text in slide_texts(current.slides[SLIDE_INDEX])):
        raise RuntimeError("Slide 7 is not the expected deconvolution-importance slide")
    slide_part = str(current.slides[SLIDE_INDEX].part.partname).lstrip("/")

    with TemporaryDirectory(prefix="shapemix_impact_slide_") as temp_dir:
        temp_root = Path(temp_dir)
        generated = temp_root / "generated.pptx"
        assembled = temp_root / "assembled.pptx"

        presentation = Presentation(DECK)
        build(presentation.slides[SLIDE_INDEX])
        presentation.save(generated)
        assemble(DECK, generated, assembled, slide_part)

        check = Presentation(assembled)
        if len(check.slides) != 31:
            raise RuntimeError("Slide count changed while updating slide 7")
        joined = " ".join(slide_texts(check.slides[SLIDE_INDEX]))
        for needle in (
            TITLE,
            "WIDELY USED TOOL",
            "Improve a method used by many researchers",
            "ShapeMix asks whether fragment-length information",
        ):
            if needle not in joined:
                raise RuntimeError(f"Revised slide 7 is missing {needle!r}")

        after = package_hashes(assembled)
        for name, digest in before.items():
            if name == slide_part:
                continue
            if after.get(name) != digest:
                raise RuntimeError(f"Unexpected change outside slide 7: {name}")

        assembled.replace(DECK)

    print(f"Revised slide 7 to emphasize broad deconvolution impact: {DECK}")
    print("Verified: every package part except slide 7 is byte-identical")


if __name__ == "__main__":
    main()
